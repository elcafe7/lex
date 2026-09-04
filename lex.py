#!/usr/bin/env python3
"""
Lex: The Elegant Bible Terminal
A source-aware, local-first CLI tool for Bible study, language inspection,
and traversal of the Christian tradition.

Supports multiple Bible versions, interlinear study, global FTS5 search,
historical creeds/confessions, and manifest-driven auto-updates.
"""
import sqlite3
import os
import sys
import re
import json
import argparse
from collections import Counter
import shlex
import time
import html
import subprocess
import shutil
import plistlib
from rich.align import Align
from rich.console import Console, Group
from rich.panel import Panel
from rich.text import Text
from rich.table import Table
from rich import box
from rich.markdown import Markdown
from rich.theme import Theme
from rich.prompt import Prompt, IntPrompt
from rich.rule import Rule

import urllib.request
import hashlib

# ... (rest of imports)

# ---------------------------------------------------------------------------
# Update Manager
# ---------------------------------------------------------------------------
class LexUpdateManager:
    MANIFEST_URL = "https://raw.githubusercontent.com/elcafe7/lex/main/manifest.json"
    RAW_BASE_URL = "https://raw.githubusercontent.com/elcafe7/lex/main/"

    def __init__(self, console, data_dir=None):
        self.console = console
        self.data_dir = data_dir or RUNTIME_DATA_DIR

    def get_local_hash(self, filepath):
        if not os.path.exists(filepath):
            return None
        sha256_hash = hashlib.sha256()
        with open(filepath, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()

    def fetch_remote_manifest(self):
        try:
            with urllib.request.urlopen(self.MANIFEST_URL) as response:
                return json.loads(response.read().decode())
        except Exception as e:
            self.console.print(f"[warning]Failed to fetch update manifest: {e}[/]")
            return None

    def resolve_asset_path(self, rel_path):
        if not rel_path.startswith("runtime-data/"):
            return None
        actual_rel = rel_path.replace("runtime-data/", "", 1)
        if not actual_rel or actual_rel.startswith(("/", "\\")):
            return None

        data_root = os.path.realpath(self.data_dir)
        target_path = os.path.realpath(os.path.join(data_root, actual_rel))
        try:
            if os.path.commonpath([data_root, target_path]) != data_root:
                return None
        except ValueError:
            return None
        return target_path

    def check_for_updates(self):
        remote = self.fetch_remote_manifest()
        if not remote: return None, None

        updates_needed = []
        for rel_path, info in remote.get("assets", {}).items():
            local_path = self.resolve_asset_path(rel_path)
            if not local_path:
                continue

            local_hash = self.get_local_hash(local_path)
            if local_hash != info["hash"]:
                updates_needed.append(rel_path)

        return updates_needed, remote["version"]

    def ensure_data(self):
        """Ensures that essential data files exist. If not, trigger a full update."""
        critical_file = os.path.join(self.data_dir, "lexicon.db")
        if not os.path.exists(critical_file):
            self.console.print("[info]First run detected: Downloading Bible databases (approx 280MB)...[/]")
            if not self.perform_update():
                self.console.print("[error]Lex data is not installed and the automatic download did not complete.[/]")
                self.console.print(f"[info]Expected data directory: {self.data_dir}[/]")
                self.console.print("[info]For the most reliable install, clone the full repository and run ./setup.sh.[/]")
                return False
        return os.path.exists(critical_file)

    def perform_update(self):
        updates, remote_version = self.check_for_updates()
        if updates is None:
            return False

        if not updates:
            self.console.print("[success]Lex is already up to date.[/]")
            return True

        self.console.print(f"[info]Updating Lex data to {remote_version}... ({len(updates)} files)[/]")

        for rel_path in updates:
            self.console.print(f"  → Downloading {rel_path}...")
            url = self.RAW_BASE_URL + rel_path

            target_path = self.resolve_asset_path(rel_path)
            if not target_path:
                self.console.print(f"[error]Skipping unsafe manifest path: {rel_path}[/]")
                return False

            os.makedirs(os.path.dirname(target_path), exist_ok=True)

            try:
                urllib.request.urlretrieve(url, target_path + ".tmp")
                os.replace(target_path + ".tmp", target_path)
            except Exception as e:
                self.console.print(f"[error]Failed to download {rel_path}: {e}[/]")
                return False

        self.console.print(f"[success]Successfully updated Lex to {remote_version}![/]")
        return True

# ---------------------------------------------------------------------------
# Runtime paths and bundled-data adapters
# ---------------------------------------------------------------------------
# Lex is currently a single-file CLI that reads several local SQLite/JSON data
# stores. Keep these paths centralized so future packaging can replace them
# with config/env-driven paths without touching feature code.
VERSION = "2.6.0"
HISTORY_FILE = os.path.expanduser("~/.lex_history")
QUERY_HISTORY_FILE = os.path.expanduser("~/.lex_query_history")
QUERY_HISTORY_LIMIT = 200
CONFIG_FILE = os.path.expanduser("~/.lex_config.json")
MANUSCRIPT_CACHE_DIR = os.path.expanduser("~/.cache/lex/manuscripts")
# The browser UI lives on apocalypse.press, while its generated data assets are
# served from Poeta's static /lex-web/ tree.
MANUSCRIPT_WEB_BASE = "https://poeta.icu/lex-web/"

# Local-first path resolution. Clones ship the compact runtime data bundle
# (SQLite DBs and JSON) under runtime-data/, while local developer worktrees
# may also have full upstream data checkouts beside lex.py.
BASE_DIR = os.path.dirname(os.path.realpath(__file__))
RUNTIME_DATA_DIR = os.path.join(BASE_DIR, "runtime-data")
HOME_FALLBACK = os.path.expanduser("~/bible-lexicon-data")

# Determine which data directory to use
if os.path.exists(RUNTIME_DATA_DIR) and os.path.exists(os.path.join(RUNTIME_DATA_DIR, "lexicon.db")):
    DATA_DIR = RUNTIME_DATA_DIR
else:
    DATA_DIR = HOME_FALLBACK

def get_lex_path(relative_path):
    return os.path.join(DATA_DIR, relative_path)

LEXICON_DB_PATH = get_lex_path("lexicon.db")

# Bible Versions Configuration
BIBLE_VERSIONS = {
    "esv": {"name": "English Standard Version", "file": "bible_versions/esv.db"},
    "kjv": {"name": "King James Version (Oxford 1769)", "file": "bible_versions/kjv.db"},
    "kj16": {"name": "King James Version (1611)", "file": "bible_versions/kj16.db"},
    "nasb": {"name": "New American Standard Bible (1995)", "file": "bible_versions/nasb.db"},
    "gen": {"name": "Geneva Bible (1587)", "file": "bible_versions/gen.db"},
    "lxx": {"name": "Septuagint (Rahlfs 1935)", "file": "bible_versions/lxx.db"},
    "vulg": {"name": "Clementine Vulgate", "file": "bible_versions/vulg.db", "latin_db": "latin.db", "type": "latin"},
    "vulgate": {"name": "Clementine Vulgate", "file": "bible_versions/vulg.db", "latin_db": "latin.db", "type": "latin", "alias": "vulg"},
}

def get_bible_path(bible_id):
    if bible_id in BIBLE_VERSIONS:
        return get_lex_path(BIBLE_VERSIONS[bible_id]["file"])
    return get_lex_path("bible_versions/esv.db")

def bible_version_available(bible_id):
    return os.path.exists(get_bible_path(bible_id))

def print_missing_bible_version(bible_id):
    info = BIBLE_VERSIONS.get(bible_id, {"name": bible_id})
    console.print(
        f"[error]Bible version [bold cyan]{bible_id}[/] ({info['name']}) is not installed.[/]"
    )
    console.print("[info]Run `lex update`, choose another version with `lex -v`, or reinstall the runtime data.[/]")

ENCYCLOPEDIA_DB_PATH = get_lex_path("encyclopedia.db")
CROSS_REFS_DB_PATH = get_lex_path("cross_refs.db")
STRONGS_DB_PATH = get_lex_path("strongs.db")
DICTIONARY_DB_PATH = get_lex_path("dictionary.db")
CREEDS_DB_PATH = get_lex_path("creeds.db")
PLACES_DB_PATH = get_lex_path("places.db")
NAVES_DB_PATH = get_lex_path("naves.db")
LXX_DB_PATH = get_lex_path("lxx.db")
HENRY_DB_PATH = get_lex_path("commentaries/matthew_henry.db")
CALVIN_DB_PATH = get_lex_path("commentaries/john_calvin.db")
INTERLINEAR_PATH = get_lex_path("esv-data/data/esv/esv-interlinear.json")
INTERLINEAR_STRONGS_PATH = get_lex_path("esv-data/data/interlinear/strongs.json")
STEP_GREEK_PATH = get_lex_path("theolog-ai/data/biblical-languages/stepbible-lexicons/tbesg-greek.json")
STEP_HEBREW_PATH = get_lex_path("theolog-ai/data/biblical-languages/stepbible-lexicons/tbesh-hebrew.json")
HISTORICAL_DOCS_DIR = get_lex_path("theolog-ai/data/historical-documents")

# The creeds table in lexicon.db has placeholder rows for some documents. This
# map lets the UI fall back to the complete local JSON document when needed.
HISTORICAL_DOC_FILES = {
    "The Apostles' Creed": "apostles-creed.json",
    "The Nicene Creed": "nicene-creed.json",
    "Athanasian Creed": "athanasian-creed.json",
    "Chalcedonian Definition": "chalcedonian-definition.json",
    "Augsburg Confession": "augsburg-confession.json",
    "Baltimore Catechism": "baltimore-catechism.json",
    "Belgic Confession": "belgic-confession.json",
    "Canons of Dort": "canons-of-dort.json",
    "Confession of Dositheus": "confession-of-dositheus.json",
    "Council of Trent": "council-of-trent.json",
    "Heidelberg Catechism": "heidelberg-catechism.json",
    "London Baptist Confession of Faith": "london-baptist-1689.json",
    "The Longer Catechism of the Orthodox Church": "philaret-catechism.json",
    "Thirty-Nine Articles": "39-articles.json",
    "Westminster Confession of Faith": "westminster-confession.json",
    "Westminster Larger Catechism": "westminster-larger-catechism.json",
    "Westminster Shorter Catechism": "westminster-shorter-catechism.json",
}

# TSK cross-reference data uses abbreviated references like "John.3.16", while
# the Bible DB uses "esv:John:3:16". These maps are the bridge between them.
TSK_BOOK_ABBR = {
    "Genesis": "Gen.", "Exodus": "Ex.", "Leviticus": "Lev.", "Numbers": "Num.",
    "Deuteronomy": "Deut.", "Joshua": "Josh.", "Judges": "Judg.", "Ruth": "Ruth",
    "1 Samuel": "1Sam.", "2 Samuel": "2Sam.", "1 Kings": "1Kgs.", "2 Kings": "2Kgs.",
    "1 Chronicles": "1Chr.", "2 Chronicles": "2Chr.", "Ezra": "Ezra", "Nehemiah": "Neh.",
    "Esther": "Est.", "Job": "Job", "Psalms": "Ps.", "Proverbs": "Prov.",
    "Ecclesiastes": "Eccl.", "Song of Solomon": "Song", "Isaiah": "Isa.", "Jeremiah": "Jer.",
    "Lamentations": "Lam.", "Ezekiel": "Ezek.", "Daniel": "Dan.", "Hosea": "Hos.",
    "Joel": "Joel", "Amos": "Amos", "Obadiah": "Obad.", "Jonah": "Jonah",
    "Micah": "Mic.", "Nahum": "Nah.", "Habakkuk": "Hab.", "Zephaniah": "Zeph.",
    "Haggai": "Hag.", "Zechariah": "Zech.", "Malachi": "Mal.", "Matthew": "Matt.",
    "Mark": "Mark", "Luke": "Luke", "John": "John", "Acts": "Acts", "Romans": "Rom.",
    "1 Corinthians": "1Cor.", "2 Corinthians": "2Cor.", "Galatians": "Gal.",
    "Ephesians": "Eph.", "Philippians": "Phil.", "Colossians": "Col.",
    "1 Thessalonians": "1Thess.", "2 Thessalonians": "2Thess.", "1 Timothy": "1Tim.",
    "2 Timothy": "2Tim.", "Titus": "Titus", "Philemon": "Phlm.", "Hebrews": "Heb.",
    "James": "Jas.", "1 Peter": "1Pet.", "2 Peter": "2Pet.", "1 John": "1John.",
    "2 John": "2John.", "3 John": "3John.", "Jude": "Jude", "Revelation": "Rev.",
}
NAVES_BOOK_ABBR = {
    "Genesis": "GEN", "Exodus": "EXO", "Leviticus": "LEV", "Numbers": "NUM",
    "Deuteronomy": "DEU", "Joshua": "JOS", "Judges": "JDG", "Ruth": "RUT",
    "1 Samuel": "1SA", "2 Samuel": "2SA", "1 Kings": "1KI", "2 Kings": "2KI",
    "1 Chronicles": "1CH", "2 Chronicles": "2CH", "Ezra": "EZR", "Nehemiah": "NEH",
    "Esther": "EST", "Job": "JOB", "Psalms": "PSA", "Proverbs": "PRO",
    "Ecclesiastes": "ECC", "Song of Solomon": "So", "Isaiah": "ISA", "Jeremiah": "JER",
    "Lamentations": "LAM", "Ezekiel": "EZK", "Daniel": "DAN", "Hosea": "HOS",
    "Joel": "JOL", "Amos": "AMO", "Obadiah": "OBA", "Jonah": "JON",
    "Micah": "MIC", "Nahum": "NAM", "Habakkuk": "HAB", "Zephaniah": "ZEP",
    "Haggai": "HAG", "Zechariah": "ZEC", "Malachi": "MAL", "Matthew": "MAT",
    "Mark": "MRK", "Luke": "LUK", "John": "JHN", "Acts": "ACT", "Romans": "ROM",
    "1 Corinthians": "1CO", "2 Corinthians": "2CO", "Galatians": "GAL",
    "Ephesians": "EPH", "Philippians": "PHP", "Colossians": "COL",
    "1 Thessalonians": "1TH", "2 Thessalonians": "2TH", "1 Timothy": "1TI",
    "2 Timothy": "2TI", "Titus": "TIT", "Philemon": "PHM", "Hebrews": "HEB",
    "James": "JAS", "1 Peter": "1PE", "2 Peter": "2PE", "1 John": "1JN",
    "2 John": "2JN", "3 John": "3JN", "Jude": "Jude", "Revelation": "REV",
}
LXX_BOOK_CODES = {
    "Genesis": "GEN",
    "Exodus": "EXO",
    "Leviticus": "LEV",
    "Numbers": "NUM",
    "Deuteronomy": "DEU",
    "Joshua": "JOS",
    "Judges": "JDG",
    "Ruth": "RUT",
    "1 Samuel": "1SA",
    "2 Samuel": "2SA",
    "1 Kings": "1KI",
    "2 Kings": "2KI",
    "1 Chronicles": "1CH",
    "2 Chronicles": "2CH",
    "Ezra": "EZR",
    "Nehemiah": "NEH",
    "Esther": "EST",
    "Job": "JOB",
    "Psalms": "PSA",
    "Proverbs": "PRO",
    "Ecclesiastes": "ECC",
    "Song of Solomon": "SNG",
    "Isaiah": "ISA",
    "Jeremiah": "JER",
    "Lamentations": "LAM",
    "Ezekiel": "EZK",
    "Daniel": "DAN",
    "Hosea": "HOS",
    "Joel": "JOL",
    "Amos": "AMO",
    "Obadiah": "OBA",
    "Jonah": "JON",
    "Micah": "MIC",
    "Nahum": "NAM",
    "Habakkuk": "HAB",
    "Zephaniah": "ZEP",
    "Haggai": "HAG",
    "Zechariah": "ZEC",
    "Malachi": "MAL",
    "Joshua A": "JOSA",
    "Joshua B": "JOS",
    "Judges A": "JDGA",
    "Judges B": "JDG",
    "1 Samuel (1 Kingdoms)": "1SA",
    "2 Samuel (2 Kingdoms)": "2SA",
    "1 Kings (3 Kingdoms)": "1KI",
    "2 Kings (4 Kingdoms)": "2KI",
    "Canticle (Song of Solomon)": "SNG",
    "Ecclesiastes (Preacher)": "ECC",
    "Lamentations (Threni)": "LAM",
    "Daniel LXX": "DAN",
    "Daniel TH": "DANTH",
    "Esther (with additions)": "EST",
    "Ezra (Esdras B/II: 1-10)": "2ES",
    "Nehemiah (Esdras B/II: 11-23)": "2ES",
    "Esdras B/II": "2ES",
    "Esdras A/I": "1ES",
    "Baruch": "BAR",
    "Epistle of Jeremiah": "EPJ",
    "Judith": "JDT",
    "Tobit BA": "TOB",
    "Tobit S": "TOBS",
    "Wisdom of Solomon": "WIS",
    "Wisdom of Sirach": "SIR",
    "I Maccabees": "1MA",
    "II Maccabees": "2MA",
    "III Maccabees": "3MA",
    "IV Maccabees": "4MA",
    "Bel LXX": "BEL",
    "Bel TH": "BELTH",
    "Susanna LXX": "SUS",
    "Susanna TH": "SUSTH",
    "Odes": "ODE",
    "Psalms of Solomon": "PSS",
}
LXX_REFERENCE_BOOK_ALIASES_RAW = {
    "1 Samuel (1 Kingdoms)": ["1 Samuel", "1 Kingdoms", "1 Kingdom", "I Kingdoms", "I Kingdom"],
    "2 Samuel (2 Kingdoms)": ["2 Samuel", "2 Kingdoms", "2 Kingdom", "II Kingdoms", "II Kingdom"],
    "1 Kings (3 Kingdoms)": ["1 Kings", "3 Kingdoms", "3 Kingdom", "III Kingdoms", "III Kingdom"],
    "2 Kings (4 Kingdoms)": ["2 Kings", "4 Kingdoms", "4 Kingdom", "IV Kingdoms", "IV Kingdom"],
    "Joshua B": ["Joshua", "Joshua B"],
    "Joshua A": ["Joshua A"],
    "Judges B": ["Judges", "Judges B"],
    "Judges A": ["Judges A"],
    "Canticle (Song of Solomon)": ["Song of Solomon", "Canticle", "Canticles", "Song", "Songs"],
    "Ecclesiastes (Preacher)": ["Ecclesiastes", "Preacher", "Qoheleth"],
    "Lamentations (Threni)": ["Lamentations", "Threni"],
    "Daniel LXX": ["Daniel", "Daniel LXX", "Daniel OG"],
    "Daniel TH": ["Daniel TH", "Daniel Theodotion", "Theodotion Daniel"],
    "Esther (with additions)": ["Esther", "Greek Esther", "Additions to Esther"],
    "Ezra (Esdras B/II: 1-10)": ["Ezra", "Esdras B", "Esdras II", "2 Esdras", "II Esdras"],
    "Nehemiah (Esdras B/II: 11-23)": ["Nehemiah", "Esdras B Nehemiah"],
    "Esdras A/I": ["Esdras A", "Esdras I", "1 Esdras", "I Esdras", "Greek Esdras"],
    "Baruch": ["Baruch"],
    "Epistle of Jeremiah": ["Epistle of Jeremiah", "Letter of Jeremiah", "Jeremiah Letter"],
    "Judith": ["Judith"],
    "Tobit BA": ["Tobit", "Tobit BA"],
    "Tobit S": ["Tobit S", "Tobit Sinaiticus"],
    "Wisdom of Solomon": ["Wisdom of Solomon", "Wisdom", "Wisdom Solomon"],
    "Wisdom of Sirach": ["Wisdom of Sirach", "Sirach", "Ecclesiasticus", "Ben Sira"],
    "I Maccabees": ["I Maccabees", "1 Maccabees", "1 Macc", "1 Mac", "First Maccabees"],
    "II Maccabees": ["II Maccabees", "2 Maccabees", "2 Macc", "2 Mac", "Second Maccabees"],
    "III Maccabees": ["III Maccabees", "3 Maccabees", "3 Macc", "3 Mac", "Third Maccabees"],
    "IV Maccabees": ["IV Maccabees", "4 Maccabees", "4 Macc", "4 Mac", "Fourth Maccabees"],
    "Bel LXX": ["Bel", "Bel LXX", "Bel OG", "Bel and the Dragon", "Bel Dragon"],
    "Bel TH": ["Bel TH", "Bel Theodotion", "Theodotion Bel"],
    "Susanna LXX": ["Susanna", "Susanna LXX", "Susanna OG"],
    "Susanna TH": ["Susanna TH", "Susanna Theodotion", "Theodotion Susanna"],
    "Odes": ["Odes"],
    "Psalms of Solomon": ["Psalms of Solomon", "Psalm of Solomon", "Pss Sol", "Psalms Solomon"],
}
TSK_TO_BOOK = {abbr.rstrip("."): book for book, abbr in TSK_BOOK_ABBR.items()}
BIBLE_BOOKS = list(TSK_BOOK_ABBR.keys())
BIBLE_BOOK_INDEX = {book: idx for idx, book in enumerate(BIBLE_BOOKS)}
MANUSCRIPT_NT_BOOKS = set(BIBLE_BOOKS[BIBLE_BOOK_INDEX["Matthew"]:])
PROTESTANT_OT_BOOKS = set(BIBLE_BOOKS[:39])
PROTESTANT_NT_BOOKS = set(BIBLE_BOOKS[39:])
SOURCE_VARIANT_RANGES = [
    ("Mark", 16, 9, "Mark", 16, 20, "Longer Ending of Mark"),
    ("John", 7, 53, "John", 8, 11, "Pericope Adulterae"),
    ("Acts", 8, 37, "Acts", 8, 37, "Ethiopian Eunuch Confession"),
    ("1 John", 5, 7, "1 John", 5, 8, "Comma Johanneum"),
]

# Pre-compile normalization pattern for performance
NORM_RE = re.compile(r"[^a-z0-9]+")

LXX_REFERENCE_BOOK_ALIASES = {}
for target, aliases in LXX_REFERENCE_BOOK_ALIASES_RAW.items():
    for alias in [target, *aliases]:
        LXX_REFERENCE_BOOK_ALIASES[NORM_RE.sub("-", alias.lower()).strip("-")] = target
        LXX_REFERENCE_BOOK_ALIASES[NORM_RE.sub("", alias.lower())] = target

BOOK_SCOPE_ALIASES = {}
for book in BIBLE_BOOKS:
    low = book.lower()
    book_key = NORM_RE.sub("-", low).strip("-")
    compact_key = NORM_RE.sub("", low)
    BOOK_SCOPE_ALIASES[book_key] = book
    BOOK_SCOPE_ALIASES[compact_key] = book
    abbr = TSK_BOOK_ABBR.get(book, "").rstrip(".").lower()
    if abbr:
        BOOK_SCOPE_ALIASES[NORM_RE.sub("-", abbr).strip("-")] = book
        BOOK_SCOPE_ALIASES[NORM_RE.sub("", abbr)] = book

BOOK_SCOPE_ALIASES.update({
    "ge": "Genesis",
    "gn": "Genesis",
    "gen": "Genesis",
    "ex": "Exodus",
    "exo": "Exodus",
    "exod": "Exodus",
    "le": "Leviticus",
    "lev": "Leviticus",
    "nu": "Numbers",
    "num": "Numbers",
    "de": "Deuteronomy",
    "dt": "Deuteronomy",
    "deut": "Deuteronomy",
    "jos": "Joshua",
    "josh": "Joshua",
    "jdg": "Judges",
    "judg": "Judges",
    "ru": "Ruth",
    "1sa": "1 Samuel",
    "1sam": "1 Samuel",
    "2sa": "2 Samuel",
    "2sam": "2 Samuel",
    "1ki": "1 Kings",
    "1kgs": "1 Kings",
    "2ki": "2 Kings",
    "2kgs": "2 Kings",
    "1ch": "1 Chronicles",
    "1chr": "1 Chronicles",
    "2ch": "2 Chronicles",
    "2chr": "2 Chronicles",
    "ezr": "Ezra",
    "neh": "Nehemiah",
    "est": "Esther",
    "psalm": "Psalms",
    "ps": "Psalms",
    "psa": "Psalms",
    "psm": "Psalms",
    "pss": "Psalms",
    "pr": "Proverbs",
    "pro": "Proverbs",
    "prov": "Proverbs",
    "ec": "Ecclesiastes",
    "ecc": "Ecclesiastes",
    "eccl": "Ecclesiastes",
    "song": "Song of Solomon",
    "sos": "Song of Solomon",
    "canticles": "Song of Solomon",
    "is": "Isaiah",
    "isa": "Isaiah",
    "jr": "Jeremiah",
    "jer": "Jeremiah",
    "lam": "Lamentations",
    "eze": "Ezekiel",
    "ezek": "Ezekiel",
    "ezk": "Ezekiel",
    "da": "Daniel",
    "dn": "Daniel",
    "dan": "Daniel",
    "hos": "Hosea",
    "jl": "Joel",
    "am": "Amos",
    "ob": "Obadiah",
    "obad": "Obadiah",
    "jon": "Jonah",
    "mi": "Micah",
    "mic": "Micah",
    "na": "Nahum",
    "nah": "Nahum",
    "hab": "Habakkuk",
    "zep": "Zephaniah",
    "zeph": "Zephaniah",
    "hag": "Haggai",
    "zec": "Zechariah",
    "zech": "Zechariah",
    "mal": "Malachi",
    "mt": "Matthew",
    "mat": "Matthew",
    "matt": "Matthew",
    "mk": "Mark",
    "mrk": "Mark",
    "lk": "Luke",
    "lu": "Luke",
    "jn": "John",
    "jhn": "John",
    "joh": "John",
    "ac": "Acts",
    "ro": "Romans",
    "rom": "Romans",
    "1co": "1 Corinthians",
    "1cor": "1 Corinthians",
    "2co": "2 Corinthians",
    "2cor": "2 Corinthians",
    "gal": "Galatians",
    "eph": "Ephesians",
    "php": "Philippians",
    "phil": "Philippians",
    "col": "Colossians",
    "1th": "1 Thessalonians",
    "1thess": "1 Thessalonians",
    "2th": "2 Thessalonians",
    "2thess": "2 Thessalonians",
    "1ti": "1 Timothy",
    "1tim": "1 Timothy",
    "2ti": "2 Timothy",
    "2tim": "2 Timothy",
    "tit": "Titus",
    "phm": "Philemon",
    "phlm": "Philemon",
    "heb": "Hebrews",
    "jas": "James",
    "jam": "James",
    "1pe": "1 Peter",
    "1pet": "1 Peter",
    "2pe": "2 Peter",
    "2pet": "2 Peter",
    "1jn": "1 John",
    "1jhn": "1 John",
    "2jn": "2 John",
    "2jhn": "2 John",
    "3jn": "3 John",
    "3jhn": "3 John",
    "rev": "Revelation",
    "rv": "Revelation",
    "re": "Revelation",
    "revelations": "Revelation",
    "tobit": "Tob",
    "tob": "Tob",
    "judith": "Jdt",
    "jdt": "Jdt",
    "wisdom": "Wis",
    "wisdomofsolomon": "Wis",
    "wis": "Wis",
    "sirach": "Sir",
    "ecclesiasticus": "Sir",
    "sir": "Sir",
    "baruch": "Bar",
    "bar": "Bar",
    "1maccabees": "1Macc",
    "1macc": "1Macc",
    "1mac": "1Macc",
    "firstmaccabees": "1Macc",
    "2maccabees": "2Macc",
    "2macc": "2Macc",
    "2mac": "2Macc",
    "secondmaccabees": "2Macc",
})

BOOK_SCOPE_GROUPS = {
    "ot": BIBLE_BOOKS[:39],
    "old-testament": BIBLE_BOOKS[:39],
    "nt": BIBLE_BOOKS[39:],
    "new-testament": BIBLE_BOOKS[39:],
    "law": BIBLE_BOOKS[:5],
    "pentateuch": BIBLE_BOOKS[:5],
    "torah": BIBLE_BOOKS[:5],
    "history": BIBLE_BOOKS[5:17],
    "wisdom": BIBLE_BOOKS[17:22],
    "poetry": BIBLE_BOOKS[17:22],
    "major": BIBLE_BOOKS[22:27],
    "major-prophets": BIBLE_BOOKS[22:27],
    "minor": BIBLE_BOOKS[27:39],
    "minor-prophets": BIBLE_BOOKS[27:39],
    "prophets": BIBLE_BOOKS[22:39],
    "gospels": BIBLE_BOOKS[39:43],
    "gospel": BIBLE_BOOKS[39:43],
    "epistles": BIBLE_BOOKS[44:66],
    "letters": BIBLE_BOOKS[44:66],
    "pauline": BIBLE_BOOKS[44:57],
    "general-epistles": BIBLE_BOOKS[57:66],
}

# Original-language creed text is only stored for short documents where
# side-by-side display is useful. Longer confessions stay English-only for now.
CREED_ORIGINALS = {
    "The Apostles' Creed": {
        "language": "Latin",
        "sections": {
            "God the Father": "Credo in Deum Patrem omnipotentem, Creatorem caeli et terrae.",
            "Jesus Christ": "Et in Iesum Christum, Filium eius unicum, Dominum nostrum, qui conceptus est de Spiritu Sancto, natus ex Maria Virgine, passus sub Pontio Pilato, crucifixus, mortuus, et sepultus; descendit ad inferos; tertia die resurrexit a mortuis; ascendit ad caelos; sedet ad dexteram Dei Patris omnipotentis; inde venturus est iudicare vivos et mortuos.",
            "The Holy Spirit and the Church": "Credo in Spiritum Sanctum, sanctam Ecclesiam catholicam, sanctorum communionem, remissionem peccatorum, carnis resurrectionem, vitam aeternam. Amen.",
        },
    },
    "The Nicene Creed": {
        "language": "Greek",
        "sections": {
            "God the Father": "Πιστεύομεν εἰς ἕνα Θεόν, Πατέρα, παντοκράτορα, ποιητὴν οὐρανοῦ καὶ γῆς, ὁρατῶν τε πάντων καὶ ἀοράτων.",
            "Jesus Christ the Son": "Καὶ εἰς ἕνα Κύριον Ἰησοῦν Χριστόν, τὸν Υἱὸν τοῦ Θεοῦ τὸν μονογενῆ, τὸν ἐκ τοῦ Πατρὸς γεννηθέντα πρὸ πάντων τῶν αἰώνων· φῶς ἐκ φωτός, Θεὸν ἀληθινὸν ἐκ Θεοῦ ἀληθινοῦ, γεννηθέντα, οὐ ποιηθέντα, ὁμοούσιον τῷ Πατρί, δι᾿ οὗ τὰ πάντα ἐγένετο. Τὸν δι᾿ ἡμᾶς τοὺς ἀνθρώπους καὶ διὰ τὴν ἡμετέραν σωτηρίαν κατελθόντα ἐκ τῶν οὐρανῶν καὶ σαρκωθέντα ἐκ Πνεύματος ἁγίου καὶ Μαρίας τῆς Παρθένου καὶ ἐνανθρωπήσαντα. Σταυρωθέντα τε ὑπὲρ ἡμῶν ἐπὶ Ποντίου Πιλάτου καὶ παθόντα καὶ ταφέντα. Καὶ ἀναστάντα τῇ τρίτῃ ἡμέρᾳ, κατὰ τὰς Γραφάς. Καὶ ἀνελθόντα εἰς τοὺς οὐρανοὺς καὶ καθεζόμενον ἐκ δεξιῶν τοῦ Πατρός. Καὶ πάλιν ἐρχόμενον μετὰ δόξης κρῖναι ζῶντας καὶ νεκρούς, οὗ τῆς βασιλείας οὐκ ἔσται τέλος.",
            "The Holy Spirit": "Καὶ εἰς τὸ Πνεῦμα τὸ Ἅγιον, τὸ κύριον, τὸ ζωοποιόν, τὸ ἐκ τοῦ Πατρὸς ἐκπορευόμενον, τὸ σὺν Πατρὶ καὶ Υἱῷ συμπροσκυνούμενον καὶ συνδοξαζόμενον, τὸ λαλῆσαν διὰ τῶν προφητῶν.",
            "The Church and Final Hope": "Εἰς μίαν, ἁγίαν, καθολικὴν καὶ ἀποστολικὴν Ἐκκλησίαν. Ὁμολογοῦμεν ἓν βάπτισμα εἰς ἄφεσιν ἁμαρτιῶν. Προσδοκοῦμεν ἀνάστασιν νεκρῶν καὶ ζωὴν τοῦ μέλλοντος αἰῶνος. Ἀμήν.",
        },
    },
    "Athanasian Creed": {
        "language": "Latin",
        "sections": {
            "Opening": "Quicunque vult salvus esse, ante omnia opus est, ut teneat catholicam fidem: Quam nisi quisque integram inviolatamque servaverit, absque dubio in aeternum peribit.",
            "The Doctrine of the Trinity": "Fides autem catholica haec est: ut unum Deum in Trinitate, et Trinitatem in unitate veneremur. Neque confundentes personas, neque substantiam separantes. Alia est enim persona Patris alia Filii, alia Spiritus Sancti: Sed Patris, et Filii, et Spiritus Sancti una est divinitas, aequalis gloria, coaeterna maiestas. Qualis Pater, talis Filius, talis Spiritus Sanctus. Increatus Pater, increatus Filius, increatus Spiritus Sanctus. Immensus Pater, immensus Filius, immensus Spiritus Sanctus. Aeternus Pater, aeternus Filius, aeternus Spiritus Sanctus. Et tamen non tres aeterni, sed unus aeternus. Sicut non tres increati, nec tres immensi, sed unus increatus, et unus immensus. Similiter omnipotens Pater, omnipotens Filius, omnipotens Spiritus Sanctus. Et tamen non tres omnipotentes, sed unus omnipotens. Ita Deus Pater, Deus Filius, Deus Spiritus Sanctus. Et tamen non tres Dii, sed unus est Deus. Ita Dominus Pater, Dominus Filius, Dominus Spiritus Sanctus. Et tamen non tres Domini, sed unus est Dominus. Quia, sicut singillatim unamquamque personam Deum ac Dominum confiteri christiana veritate compellimur: ita tres Deos aut Dominos dicere catholica religione prohibemur. Pater a nullo est factus: nec creatus, nec genitus. Filius a Patre solo est: non factus, nec creatus, sed genitus. Spiritus Sanctus a Patre et Filio: non factus, nec creatus, nec genitus, sed procedens. Unus ergo Pater, non tres Patres: unus Filius, non tres Filii: unus Spiritus Sanctus, non tres Spiritus Sancti. Et in hac Trinitate nihil prius aut posterius, nihil maius aut minus: sed totae tres personae coaeternae sibi sunt et coaequales. Ita ut per omnia, sicut iam supra dictum est, et unitas in Trinitate, et Trinitas in unitate veneranda sit. Qui vult ergo salvus esse, ita de Trinitate sentiat.",
            "The Incarnation of Christ": "Sed necessarium est ad aeternam salutem, ut incarnationem quoque Domini nostri Iesu Christi fideliter credat. Est ergo fides recta ut credamus et confiteamur, quia Dominus noster Iesus Christus, Dei Filius, Deus et homo est. Deus est ex substantia Patris ante saecula genitus: et homo est ex substantia matris in saeculo natus. Perfectus Deus, perfectus homo: ex anima rationabili et humana carne subsistens. Aequalis Patri secundum divinitatem: minor Patre secundum humanitatem. Qui licet Deus sit et homo, non duo tamen, sed unus est Christus. Unus autem non conversione divinitatis in carnem, sed assumptione humanitatis in Deum. Unus omnino, non confusione substantiae, sed unitate personae. Nam sicut anima rationabilis et caro unus est homo: ita Deus et homo unus est Christus. Qui passus est pro salute nostra: descendit ad inferos: tertia die resurrexit a mortuis. Ascendit ad caelos, sedet ad dexteram Dei Patris omnipotentis: inde venturus est iudicare vivos et mortuos. Ad cuius adventum omnes homines resurgere habent cum corporibus suis: et reddituri sunt de factis propriis rationem. Et qui bona egerunt, ibunt in vitam aeternam: qui vero mala, in ignem aeternum.",
            "Conclusion": "Haec est fides catholica, quam nisi quisque fideliter firmiterque crediderit, salvus esse non poterit.",
        },
    },
    "Chalcedonian Definition": {
        "language": "Greek",
        "sections": {
            "Introduction": "Ἡ ἁγία καὶ μεγάλη καὶ οἰκουμενικὴ σύνοδος, ἡ κατὰ θεοῦ χάριν καὶ θέσπισμα τῶν εὐσεβεστάτων καὶ φιλοχρίστων ἡμῶν βασιλέων Μαρκιανοῦ καὶ Οὐαλεντινιανοῦ Αὐγούστων ἐν Χαλκηδόνι τῇ μητροπόλει τῆς Βιθυνῶν ἐπαρχίας συναχθεῖσα ἐν τῷ μαρτυρίῳ τῆς ἁγίας καὶ καλλινίκου μάρτυρος Εὐφημίας, ὥρισε τὰ ὑποτεταγμένα. Ὁ κύριος ἡμῶν καὶ σωτὴρ Ἰησοῦς Χριστός, τὴν τῆς πίστεως γνῶσιν τοῖς μαθηταῖς βεβαιῶν, ἔφη πρὸς αὐτούς· Εἰρήνην τὴν ἐμὴν δίδωμι ὑμῖν, εἰρήνην τὴν ἐμὴν ἀφίημι ὑμῖν, πρὸς τὸ μηδένα πρὸς τὸν πλησίον διχονοεῖν ἐν τοῖς τῆς εὐσεβείας δόγμασιν, ἀλλ' ἐπίσης τοῖς πᾶσι τὸ τῆς ἀληθείας ἐπιδείκνυσθαι κήρυγμα.",
            "The Definition": "Ἑπόμενοι τοίνυν τοῖς ἁγίοις πατράσιν, ἕνα καὶ τὸν αὐτὸν ὁμολογεῖν υἱὸν τὸν κύριον ἡμῶν Ἰησοῦν Χριστὸν συμφώνως ἅπαντες ἐκδιδάσκομεν, τέλειον τὸν αὐτὸν ἐν θεότητι καὶ τέλειον τὸν αὐτὸν ἐν ἀνθρωπότητι, θεὸν ἀληθῶς καὶ ἄνθρωπον ἀληθῶς τὸν αὐτόν, ἐκ ψυχῆς λογικῆς καὶ σώματος, ὁμοούσιον τῷ πατρὶ κατὰ τὴν θεότητα καὶ ὁμοούσιον τὸν αὐτὸν ἡμῖν κατὰ τὴν ἀνθρωπότητα, κατὰ πάντα ὅμοιον ἡμῖν χωρὶς ἁμαρτίας· πρὸ αἰώνων μὲν ἐκ τοῦ πατρὸς γεννηθέντα κατὰ τὴν θεότητα, ἐπ' ἐσχάτων δὲ τῶν ἡμερῶν τὸν αὐτὸν δι' ἡμᾶς καὶ διὰ τὴν ἡμετέραν σωτηρίαν ἐκ Μαρίας τῆς παρθένου τῆς θεοτόκου κατὰ τὴν ἀνθρωπότητα, ἕνα καὶ τὸν αὐτὸν Χριστόν, υἱόν, κύριον, μονογενῆ, ἐν δύο φύσεσιν ἀσυγχύτως, ἀτρέπτως, ἀδιαιρέτως, ἀχωρίστως γνωριζόμενον· οὐδαμοῦ τῆς τῶν φύσεων διαφορᾶς ἀνῃρημένης διὰ τὴν ἕνωσιν, σωζομένης δὲ μᾶλλον τῆς ἰδιότητος ἑκατέρας φύσεως καὶ εἰς ἓν πρόσωπον καὶ μίαν ὑπόστασιν συντρεχούσης, οὐκ εἰς δύο πρόσωπα μεριζόμενον ἢ διαιρούμενον, ἀλλ' ἕνα καὶ τὸν αὐτὸν υἱὸν καὶ μονογενῆ θεὸν λόγον, κύριον Ἰησοῦν Χριστόν· καθάπερ ἄνωθεν οἱ προφῆται περὶ αὐτοῦ καὶ αὐτὸς ἡμᾶς ὁ κύριος Ἰησοῦς Χριστὸς ἐξεπαίδευσε καὶ τὸ τῶν πατέρων ἡμῖν παραδέδωκε σύμβολον.",
            "Prohibitions": "Τούτων τοίνυν μετὰ πάσης πανταχόθεν ἀκριβείας τε καὶ ἐμμελείας διατυπωθέντων ἡμῖν, ὥρισεν ἡ ἁγία καὶ οἰκουμενικὴ σύνοδος, ἑτέραν πίστιν μηδενὶ ἐξεῖναι προφέρειν, ἤγουν συγγράφειν ἢ συντιθέναι ἢ φρονεῖν ἢ διδάσκειν ἑτέρους· τοὺς δὲ τολμῶντας ἢ συντιθέναι πίστιν ἑτέραν, ἤγουν προκομίζειν ἢ διδάσκειν ἢ παραδιδόναι ἕτερον σύμβολον τοῖς ἐθέλουσιν ἐπιστρέφειν εἰς ἐπίγνωσιν τῆς ἀληθείας ἐξ Ἑλληνισμοῦ ἢ ἐξ Ἰουδαϊσμοῦ ἢ γοῦν ἐξ αἱρέσεως οἱασδηποτοῦν, τούτους, εἰ μὲν εἶεν ἐπίσκοποι ἢ κληρικοί, ἀλλοτρίους εἶναι τοὺς ἐπισκόπους τῆς ἐπισκοπῆς καὶ τοὺς κληρικοὺς τοῦ κλήρου, εἰ δὲ μονάζοντες ἢ λαϊκοὶ εἶεν, ἀναθεματίζεσθαι.",
        },
    },
}

# Topic-level notes for creed texts. These are intentionally separated from the
# creed body so we can explain textual/traditional variants without altering the
# source document text.
CREED_NOTES = {
    "The Nicene Creed": (
        "**Filioque note:** This local English text includes the Filioque clause "
        "('and the Son') in the line on the Holy Spirit: 'He proceeds from the "
        "Father and the Son.' The Greek text shown here preserves the older "
        "conciliar wording, 'from the Father,' without the later Latin addition.\n\n"
        "**Generally accept/use the Filioque:** Roman Catholic/Latin Western "
        "tradition and many Western Protestant traditions, including much "
        "Anglican, Lutheran, Reformed, Methodist, and Baptist usage.\n\n"
        "**Generally deny or omit the Filioque:** Eastern Orthodox churches, "
        "Oriental Orthodox churches, and the Church of the East. Some Eastern "
        "Catholic churches may omit it liturgically while remaining in communion "
        "with Rome."
    )
}

# Rich styles used across panels/tables. Keep style names stable; rendering
# methods reference these string keys directly.
def load_config():
    try:
        with open(CONFIG_FILE, "r", encoding="utf-8") as fh:
            data = json.load(fh)
        return data if isinstance(data, dict) else {}
    except (OSError, json.JSONDecodeError):
        return {}

def save_config(config):
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as fh:
            json.dump(config, fh, indent=2, sort_keys=True)
            fh.write("\n")
    except OSError:
        pass

def save_theme_preference(theme_mode):
    config = load_config()
    config["theme"] = theme_mode
    save_config(config)

def clear_theme_preference():
    config = load_config()
    config.pop("theme", None)
    save_config(config)

def save_bible_preference(bible_id):
    config = load_config()
    config["bible"] = bible_id
    save_config(config)

def load_bible_preference():
    return load_config().get("bible", "esv")

def normalize_theme_value(value):
    value = (value or "").strip().lower()
    if not value:
        return None
    if value in {"light", "bright", "day"}:
        return "light"
    if value in {"dark", "black", "night"}:
        return "dark"
    if re.search(r"(^|[^a-z])(light|bright|day)([^a-z]|$)", value):
        return "light"
    if re.search(r"(^|[^a-z])(dark|black|night)([^a-z]|$)", value):
        return "dark"
    return None

def rgb_luminance(red, green, blue):
    def linearize(channel):
        return channel / 12.92 if channel <= 0.03928 else ((channel + 0.055) / 1.055) ** 2.4
    r, g, b = linearize(red), linearize(green), linearize(blue)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b

def theme_from_rgb(red, green, blue):
    return "light" if rgb_luminance(red, green, blue) >= 0.45 else "dark"

def rgb_from_archived_color(data):
    if not isinstance(data, (bytes, bytearray)):
        return None
    matches = re.findall(rb"([01](?:\.\d+)?) ([01](?:\.\d+)?) ([01](?:\.\d+)?) ([01](?:\.\d+)?)", data)
    if not matches:
        return None
    red, green, blue, _alpha = (float(part) for part in matches[-1])
    return red, green, blue

def theme_from_colorfgbg():
    colorfgbg = os.environ.get("COLORFGBG", "")
    if colorfgbg:
        try:
            background = int(colorfgbg.split(";")[-1])
            return "light" if background in range(7, 16) else "dark"
        except ValueError:
            pass
    return None

def theme_from_env_hints():
    for key in (
        "LEX_THEME",
        "TERMINAL_THEME",
        "COLOR_SCHEME",
        "THEME",
        "VSCODE_COLOR_THEME",
        "ITERM_PROFILE",
        "WT_PROFILE",
        "TERM_PROGRAM",
    ):
        theme = normalize_theme_value(os.environ.get(key))
        if theme:
            return theme
    return None

def theme_from_apple_terminal_profile():
    if sys.platform != "darwin" or os.environ.get("TERM_PROGRAM") != "Apple_Terminal":
        return None
    try:
        exported = subprocess.run(
            ["defaults", "export", "com.apple.Terminal", "-"],
            capture_output=True,
            timeout=0.75,
            check=False,
        )
        if exported.returncode != 0 or not exported.stdout:
            return None
        prefs = plistlib.loads(exported.stdout)
    except (OSError, plistlib.InvalidFileException, subprocess.SubprocessError):
        return None

    profile_name = (
        os.environ.get("TERM_PROFILE")
        or prefs.get("Startup Window Settings")
        or prefs.get("Default Window Settings")
    )
    profile = (prefs.get("Window Settings") or {}).get(profile_name or "")
    if not isinstance(profile, dict):
        return normalize_theme_value(profile_name)

    rgb = rgb_from_archived_color(profile.get("BackgroundColor"))
    if rgb:
        return theme_from_rgb(*rgb)

    return normalize_theme_value(profile_name) or "light"

def theme_from_iterm_profile():
    if sys.platform != "darwin" or os.environ.get("TERM_PROGRAM") != "iTerm.app":
        return None
    try:
        exported = subprocess.run(
            ["defaults", "export", "com.googlecode.iterm2", "-"],
            capture_output=True,
            timeout=0.75,
            check=False,
        )
        if exported.returncode != 0 or not exported.stdout:
            return None
        prefs = plistlib.loads(exported.stdout)
    except (OSError, plistlib.InvalidFileException, subprocess.SubprocessError):
        return None

    profiles = prefs.get("New Bookmarks") or []
    profile_hint = os.environ.get("ITERM_PROFILE")
    default_guid = prefs.get("Default Bookmark Guid")
    for profile in profiles:
        if not isinstance(profile, dict):
            continue
        name = profile.get("Name")
        guid = profile.get("Guid")
        if profile_hint and name != profile_hint:
            continue
        if not profile_hint and default_guid and guid != default_guid:
            continue
        red = profile.get("Background Color", {}).get("Red Component")
        green = profile.get("Background Color", {}).get("Green Component")
        blue = profile.get("Background Color", {}).get("Blue Component")
        if None not in (red, green, blue):
            return theme_from_rgb(float(red), float(green), float(blue))
        return normalize_theme_value(name)
    return normalize_theme_value(profile_hint)

def theme_from_linux_desktop():
    if not sys.platform.startswith("linux"):
        return None

    commands = (
        ["gsettings", "get", "org.gnome.desktop.interface", "color-scheme"],
        ["gsettings", "get", "org.gnome.desktop.interface", "gtk-theme"],
        ["kreadconfig6", "--group", "General", "--key", "ColorScheme"],
        ["kreadconfig5", "--group", "General", "--key", "ColorScheme"],
    )
    for command in commands:
        executable = shutil.which(command[0])
        if not executable:
            continue
        try:
            proc = subprocess.run(
                [executable, *command[1:]],
                capture_output=True,
                text=True,
                timeout=0.5,
                check=False,
            )
        except (OSError, subprocess.SubprocessError):
            continue
        if proc.returncode == 0:
            theme = normalize_theme_value(proc.stdout.strip().strip("'\""))
            if theme:
                return theme

    return None

def theme_from_macos_appearance():
    if sys.platform != "darwin":
        return None
    try:
        proc = subprocess.run(
            ["defaults", "read", "-g", "AppleInterfaceStyle"],
            capture_output=True,
            text=True,
            timeout=0.5,
            check=False,
        )
        if proc.returncode == 0 and proc.stdout.strip().lower() == "dark":
            return "dark"
        return "light"
    except (OSError, subprocess.SubprocessError):
        return None

def detect_terminal_theme():
    for detector in (
        theme_from_env_hints,
        theme_from_colorfgbg,
        theme_from_apple_terminal_profile,
        theme_from_iterm_profile,
        theme_from_linux_desktop,
        theme_from_macos_appearance,
    ):
        theme = detector()
        if theme in {"light", "dark"}:
            return theme

    return "dark"

def resolve_theme_mode(raw_argv):
    if "-light" in raw_argv:
        return "light"
    if "-dark" in raw_argv:
        return "dark"
    if "-auto" in raw_argv:
        return detect_terminal_theme()

    env_theme = os.environ.get("LEX_THEME", "").strip().lower()
    if env_theme in {"light", "dark"}:
        return env_theme

    saved_theme = load_config().get("theme")
    if saved_theme in {"light", "dark"}:
        return saved_theme

    return detect_terminal_theme()

def has_theme_override(raw_argv):
    if "-auto" in raw_argv:
        return False
    if "-light" in raw_argv or "-dark" in raw_argv:
        return True

    env_theme = os.environ.get("LEX_THEME", "").strip().lower()
    if env_theme in {"light", "dark"}:
        return True

    return load_config().get("theme") in {"light", "dark"}

def resolve_no_color(raw_argv):
    if os.environ.get("LEX_NO_COLOR"):
        return True
    return False

def build_theme(theme_mode):
    if theme_mode == "light":
        text_style = "rgb(31,31,31)"
        strong_text_style = "bold rgb(31,31,31)"
        muted_text_style = "rgb(100,100,100)"
        accent_style = "rgb(180,30,40)" # Crimson / Red
        accent_strong_style = "bold rgb(180,30,40)"
        success_style = "rgb(30,100,50)" # Deep Green
        warning_style = "rgb(160,100,0)" # Ochre
        border_style = "rgb(200,190,170)" # Warm divider
        verse_ref_style = accent_strong_style
        verse_ref_muted_style = "rgb(130,120,110)"
        highlight_style = "rgb(31,31,31) on rgb(240,220,150) underline"
        marker_style = "bold rgb(31,31,31) on rgb(240,220,150)"
        source_style = "rgb(47,92,160)" # Biblical blue
        translit_style = "italic rgb(100,90,80)"
    else:
        text_style = "grey93"
        strong_text_style = "bold white"
        muted_text_style = "grey50"
        accent_style = "dark_orange"
        accent_strong_style = "bold dark_orange"
        success_style = "bold spring_green3"
        warning_style = "bold sandy_brown"
        border_style = "grey27"
        verse_ref_style = "bold orange3"
        verse_ref_muted_style = "grey42"
        highlight_style = "bold grey93 on grey19"
        marker_style = "bold grey93 on grey19"
        source_style = "bold sky_blue3"
        translit_style = "italic grey62"

    return Theme({
        "text": text_style,
        "text.strong": strong_text_style,
        "text.muted": muted_text_style,
        "info": accent_style,
        "warning": warning_style,
        "success": success_style,
        "ui.action": accent_style,
        "ui.action.key": accent_strong_style,
        "ui.border": border_style,
        "ui.meta": muted_text_style,
        "search.hit": highlight_style,
        "verse.marker": marker_style,
        "verse.ref": verse_ref_style,
        "verse.ref.muted": verse_ref_muted_style,
        "verse.text": text_style,
        "verse.text.focus": strong_text_style,
        "verse.text.muted": muted_text_style,
        "verse.border": border_style,
        "source.text": source_style,
        "source.translit": translit_style,
        "source.border": border_style,
        "lexicon.num": accent_strong_style,
        "lexicon.word": success_style,
        "place.name": warning_style,
        "dict.topic": "rgb(110,60,160)" if theme_mode == "light" else "bold orchid",
        "interlinear.strongs": accent_style,
        "interlinear.translit": translit_style,
    })

ACTIVE_THEME_MODE = resolve_theme_mode(sys.argv[1:])
custom_theme = build_theme(ACTIVE_THEME_MODE)
console_base_style = "rgb(31,31,31) on rgb(249,247,242)" if ACTIVE_THEME_MODE == "light" else "grey93 on grey11"

def detect_console_width():
    for candidate in (os.environ.get("COLUMNS"),):
        try:
            width = int(candidate)
            if width >= 40:
                return width
        except (TypeError, ValueError):
            pass
    try:
        width = os.get_terminal_size(sys.stdout.fileno()).columns
        if width >= 40:
            return width
    except OSError:
        pass
    width = shutil.get_terminal_size((100, 24)).columns
    return max(40, width)

class BackgroundFillWriter:
    def __init__(self, stream, fill_sequence):
        self.stream = stream
        self.fill_sequence = fill_sequence

    def write(self, data):
        if self.fill_sequence:
            data = data.replace("\n", f"{self.fill_sequence}\n")
        return self.stream.write(data)

    def flush(self):
        return self.stream.flush()

    def isatty(self):
        return self.stream.isatty()

    def __getattr__(self, name):
        return getattr(self.stream, name)

def line_fill_sequence():
    if resolve_no_color(sys.argv[1:]) or not sys.stdout.isatty():
        return ""
    return "\033[48;5;231m\033[K" if ACTIVE_THEME_MODE == "light" else "\033[40m\033[K"

console = Console(
    color_system="256",
    theme=custom_theme,
    style=console_base_style,
    no_color=resolve_no_color(sys.argv[1:]),
    width=detect_console_width(),
    file=BackgroundFillWriter(sys.stdout, line_fill_sequence()),
    emoji=True,
)

def fill_terminal_row(text, style="text"):
    rendered = Text(text, style=style)
    remaining = max(0, console.width - rendered.cell_len)
    if remaining:
        rendered.append(" " * remaining, style=style)
    return rendered

def study_note_excerpt(value, limit=520):
    text = re.sub(r"\s+", " ", str(value or "").replace("__", "")).strip()
    if len(text) <= limit:
        return text
    cut = text.rfind(" ", 0, limit)
    if cut < int(limit * 0.65):
        cut = limit
    return text[:cut].rstrip(" ,.;:") + "..."

def study_analysis_text(parsed):
    analysis = Text()
    source = f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "-"
    lemma = f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "-"
    code = parsed["strongs"] or parsed["morph"] or "-"
    gloss = parsed["gloss"] or parsed["english"] or "-"
    analysis.append("Src ", style="dim")
    analysis.append(source, style="source.text")
    analysis.append("\nLemma ", style="dim")
    analysis.append(lemma, style="lexicon.word")
    analysis.append("\nCode ", style="dim")
    analysis.append(code, style="interlinear.strongs")
    analysis.append(" · Gloss ", style="dim")
    analysis.append(gloss, style="text")
    return analysis

def study_lexicon_text(lemma, details):
    note = Text()
    note.append(lemma or "-", style="lexicon.word")
    if details:
        note.append("\n")
        note.append(details, style="text")
    return note

# ---------------------------------------------------------------------------
# Database and application coordinator
# ---------------------------------------------------------------------------
class LexDB:
    def __init__(self, db_path):
        self.db_path = db_path

    def query(self, sql, params=()):
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            cursor.execute(sql, params)
            return cursor.fetchall()

    def word_frequencies(self):
        """Return surface-word frequencies from the active Bible text."""
        counts = Counter()
        with sqlite3.connect(self.db_path) as conn:
            for (text,) in conn.execute("SELECT text FROM bible"):
                counts.update(re.findall(r"[A-Za-z]+(?:'[A-Za-z]+)?", text.lower()))
        return counts

class LexAgent:
    # LexAgent owns all local data access and terminal rendering. The CLI parser
    # at the bottom should stay thin and dispatch into these feature methods.
    def __init__(self, bible_id=None):
        if bible_id is None:
            bible_id = load_bible_preference()

        self.db = LexDB(LEXICON_DB_PATH)
        bible_path = get_bible_path(bible_id)
        self.bible_db = LexDB(bible_path if os.path.exists(bible_path) else LEXICON_DB_PATH)

        # Load Latin database for Vulgate semantic tagging
        self.latin_db = None
        if bible_id == "vulg" or (hasattr(self, 'bible_prefix') and self.bible_prefix == "vulg"):
            latin_path = get_lex_path("latin.db")
            if os.path.exists(latin_path):
                self.latin_db = LexDB(latin_path)

        # Determine reference prefix from bible metadata
        self.bible_prefix = "esv"
        prefix_res = self.bible_db.query("SELECT value FROM metadata WHERE key='reference_prefix'")
        if prefix_res:
            self.bible_prefix = prefix_res[0][0]

        # Build dynamic book mapping from the active database references
        self.canon_map = {}
        self.reverse_canon_map = {}
        # 1. First, index everything actually IN the database
        books_res = self.bible_db.query("SELECT DISTINCT reference FROM bible")
        for (ref,) in books_res:
            ref_body = ref.split(":", 1)[1] if ":" in ref else ref
            ref_parts = ref_body.rsplit(":", 2)
            if len(ref_parts) == 3:
                db_book = ref_parts[0]
                # Map the compact version of the DB string to itself
                self.canon_map[re.sub(r"[^a-z0-9]+", "", db_book.lower())] = db_book
                # Default reverse map to itself
                self.reverse_canon_map[db_book] = db_book

        # 2. Map all standard aliases and full names to the DB's preferred string
        # For each canonical book, find if any of its aliases are in the DB.
        # If so, map all other aliases to that same DB identifier.
        for book in BIBLE_BOOKS:
            # Gather all keys that refer to this book
            aliases = [a for a, f in BOOK_SCOPE_ALIASES.items() if f == book]
            aliases.append(re.sub(r"[^a-z0-9]+", "", book.lower()))

            # Find if any of these are in the DB
            db_target = None
            for a in aliases:
                if a in self.canon_map:
                    db_target = self.canon_map[a]
                    break

            if db_target:
                for a in aliases:
                    self.canon_map[a] = db_target
                # Also map reverse for study mode
                canon_target = "Psalm" if book == "Psalms" else book
                self.reverse_canon_map[db_target] = canon_target

        self.encyclopedia_db = LexDB(ENCYCLOPEDIA_DB_PATH) if os.path.exists(ENCYCLOPEDIA_DB_PATH) else None
        self.cross_refs_db = LexDB(CROSS_REFS_DB_PATH if os.path.exists(CROSS_REFS_DB_PATH) else LEXICON_DB_PATH)
        self.strongs_db = LexDB(STRONGS_DB_PATH if os.path.exists(STRONGS_DB_PATH) else LEXICON_DB_PATH)
        self.dictionary_db = LexDB(DICTIONARY_DB_PATH if os.path.exists(DICTIONARY_DB_PATH) else LEXICON_DB_PATH)
        self.creeds_db = LexDB(CREEDS_DB_PATH if os.path.exists(CREEDS_DB_PATH) else LEXICON_DB_PATH)
        self.places_db = LexDB(PLACES_DB_PATH if os.path.exists(PLACES_DB_PATH) else LEXICON_DB_PATH)
        self.henry_db = LexDB(HENRY_DB_PATH) if os.path.exists(HENRY_DB_PATH) else None
        self.calvin_db = LexDB(CALVIN_DB_PATH) if os.path.exists(CALVIN_DB_PATH) else None
        self.naves_db = LexDB(NAVES_DB_PATH) if os.path.exists(NAVES_DB_PATH) else None
        self.lxx_db = LexDB(LXX_DB_PATH) if os.path.exists(LXX_DB_PATH) else None
        self.last_ref = self.load_history()
        self._interlinear_index = None
        self._ordered_refs = None
        self._interlinear_strongs = None
        self._step_greek = None
        self._step_hebrew = None

    # -----------------------------------------------------------------------
    # Shared utilities and lazy data loading
    # -----------------------------------------------------------------------
    def load_history(self):
        if os.path.exists(HISTORY_FILE):
            try:
                with open(HISTORY_FILE, "r") as f: return f.read().strip()
            except: pass
        return None

    def save_history(self, ref):
        try:
            with open(HISTORY_FILE, "w") as f: f.write(ref)
        except: pass

    def save_query_history(self, command):
        command = command.strip()
        if not command:
            return
        entry = {
            "ts": time.strftime("%Y-%m-%d %H:%M:%S"),
            "version": self.bible_prefix,
            "command": command,
        }
        try:
            existing = self.load_query_history(newest_first=False)
            existing.append(entry)
            with open(QUERY_HISTORY_FILE, "w", encoding="utf-8") as f:
                for row in existing[-QUERY_HISTORY_LIMIT:]:
                    f.write(json.dumps(row, ensure_ascii=False) + "\n")
        except Exception:
            pass

    def load_query_history(self, limit=None, newest_first=True):
        rows = []
        if not os.path.exists(QUERY_HISTORY_FILE):
            return rows
        try:
            with open(QUERY_HISTORY_FILE, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        item = json.loads(line)
                        if isinstance(item, dict) and item.get("command"):
                            rows.append(item)
                    except json.JSONDecodeError:
                        rows.append({"ts": "", "version": "", "command": line})
        except Exception:
            return []
        if newest_first:
            rows = list(reversed(rows))
        if limit:
            rows = rows[:max(1, limit)]
        return rows

    def clear_query_history(self):
        try:
            if os.path.exists(QUERY_HISTORY_FILE):
                os.remove(QUERY_HISTORY_FILE)
            return True
        except Exception:
            return False

    def display_query_history(self, limit=25):
        rows = self.load_query_history(limit=limit)
        if not rows:
            empty = Text("No Lex query history yet.", style="text")
            console.print(
                Panel(
                    Group(empty, "", self.history_commands_table()),
                    title="History",
                    subtitle="~/.lex_query_history",
                    border_style="ui.border",
                    padding=(1, 2),
                    expand=False,
                )
            )
            return True

        table = Table(border_style="ui.border", box=box.SIMPLE_HEAVY)
        table.add_column("#", justify="right", style="ui.action.key", no_wrap=True)
        table.add_column("When", style="ui.meta", no_wrap=True)
        table.add_column("Bible", style="bold cyan", no_wrap=True)
        table.add_column("Command", style="text", overflow="fold")
        for idx, row in enumerate(rows, 1):
            table.add_row(
                str(idx),
                row.get("ts", ""),
                (row.get("version") or "-").upper(),
                row.get("command", ""),
            )
        meta = Text()
        meta.append("Stored in ", style="ui.meta")
        meta.append("~/.lex_query_history", style="text")
        meta.append("  |  ", style="ui.meta")
        meta.append("lex history --clear", style="ui.action.key")
        console.print(
            Panel(
                Group(table, "", meta, "", self.history_commands_table()),
                title=f"Recent Lex Commands ({len(rows)})",
                border_style="ui.border",
                padding=(1, 2),
                expand=False,
            )
        )
        return True

    def history_commands_table(self):
        sheet = Table(title="History Commands", border_style="ui.border", box=box.SIMPLE_HEAVY)
        sheet.add_column("Command", style="bold gold3", no_wrap=True)
        sheet.add_column("Use", style="text", overflow="fold")
        sheet.add_row("lex history", "Show recent local Lex commands")
        sheet.add_row("lex history --limit 10", "Show fewer recent commands")
        sheet.add_row("lex history --clear", "Clear local Lex command history")
        sheet.add_row("lex history clear", "Alternate clear form")
        return sheet

    def clean_text(self, text):
        # Strip inline Strong's markers and similar annotation tokens from read text.
        text = re.sub(r'\\par\b', ' ', text)
        text = re.sub(r'\{\\cf\d+\s+([^{}]*)\}', r'\1', text)
        text = re.sub(r'\\[a-zA-Z]+\d*\s?', '', text)
        text = re.sub(r'[\[<][GH]\d+[>\]]', '', text)
        text = re.sub(r'\*[a-z]+', '', text)
        text = re.sub(r'\byourln\b', 'your', text, flags=re.IGNORECASE)
        text = re.sub(r'\bonld\b', 'on', text, flags=re.IGNORECASE)
        text = re.sub(r'\[/?[a-z]+\]', '', text)
        text = re.sub(r'\s{2,}', ' ', text)
        return text.strip()

    def escape_fts_query(self, query):
        cleaned = re.sub(r"[^\w\s]", " ", query).strip()
        if not cleaned:
            return None
        return f"\"{' '.join(cleaned.split())}\""

    def fts_terms_query(self, query):
        terms = re.findall(r"\w+", query)
        if not terms:
            return None
        return " AND ".join(f'"{term}"' for term in terms)

    def fts_any_terms_query(self, query):
        terms = re.findall(r"\w+", query)
        if not terms:
            return None
        return " OR ".join(f'"{term}"' for term in terms)

    def edit_distance(self, source, target):
        """Damerau-Levenshtein distance, including adjacent transpositions."""
        if source == target:
            return 0
        if not source:
            return len(target)
        if not target:
            return len(source)
        previous_previous = None
        previous = list(range(len(target) + 1))
        for source_idx, source_char in enumerate(source, 1):
            current = [source_idx]
            for target_idx, target_char in enumerate(target, 1):
                value = min(
                    current[target_idx - 1] + 1,
                    previous[target_idx] + 1,
                    previous[target_idx - 1] + (source_char != target_char),
                )
                if (
                    previous_previous is not None
                    and source_idx > 1
                    and target_idx > 1
                    and source_char == target[target_idx - 2]
                    and source[source_idx - 2] == target_char
                ):
                    value = min(value, previous_previous[target_idx - 2] + 1)
                current.append(value)
            previous_previous, previous = previous, current
        return previous[-1]

    def fuzzy_search_query(self, query):
        """Correct missing search tokens against words in the selected Bible."""
        tokens = re.findall(r"[A-Za-z]+(?:'[A-Za-z]+)?|\d+", query)
        if not tokens:
            return query, []
        if not hasattr(self, "_bible_word_frequencies"):
            self._bible_word_frequencies = self.bible_db.word_frequencies()
        frequencies = self._bible_word_frequencies
        vocabulary = list(frequencies)
        corrected = []
        changes = []
        for token in tokens:
            lowered = token.lower()
            if lowered in frequencies or len(lowered) < 4 or not lowered.isalpha():
                corrected.append(token)
                continue
            max_distance = 1 if len(lowered) <= 5 else 2 if len(lowered) <= 9 else 3
            candidates = (
                word for word in vocabulary
                if abs(len(word) - len(lowered)) <= max_distance
                and word[0] == lowered[0]
            )
            best = None
            for candidate in candidates:
                distance = self.edit_distance(lowered, candidate)
                if distance > max_distance:
                    continue
                score = (distance, -frequencies[candidate], candidate)
                if best is None or score < best[0]:
                    best = (score, candidate)
            replacement = best[1] if best else token
            corrected.append(replacement)
            if replacement != token:
                changes.append((token, replacement))
        return " ".join(corrected), changes

    def parse_book_scope(self, token):
        raw = token.strip()
        if not raw.startswith("-") or raw.startswith("--"):
            return None
        scope = raw.lstrip("-").lower().strip()
        if not scope:
            return None
        scope = re.sub(r"[^a-z0-9]+", "-", scope).strip("-")
        if scope in BOOK_SCOPE_GROUPS:
            return {
                "label": scope,
                "books": BOOK_SCOPE_GROUPS[scope],
            }
        if scope in BOOK_SCOPE_ALIASES:
            book = BOOK_SCOPE_ALIASES[scope]
            return {
                "label": book,
                "books": [book],
            }
        aliases = sorted(BOOK_SCOPE_ALIASES, key=len, reverse=True)
        for start_alias in aliases:
            prefix = f"{start_alias}-"
            if not scope.startswith(prefix):
                continue
            end_alias = scope[len(prefix):]
            if end_alias not in BOOK_SCOPE_ALIASES:
                continue
            start_book = BOOK_SCOPE_ALIASES[start_alias]
            end_book = BOOK_SCOPE_ALIASES[end_alias]
            start_idx = BIBLE_BOOK_INDEX[start_book]
            end_idx = BIBLE_BOOK_INDEX[end_book]
            if start_idx > end_idx:
                start_idx, end_idx = end_idx, start_idx
                start_book, end_book = end_book, start_book
            return {
                "label": f"{start_book}-{end_book}",
                "books": BIBLE_BOOKS[start_idx:end_idx + 1],
            }
        return None

    def parse_search_query_and_scope(self, query):
        try:
            tokens = shlex.split(query)
        except ValueError:
            tokens = query.split()
        kept = []
        scope = None
        for token in tokens:
            parsed_scope = self.parse_book_scope(token)
            if parsed_scope:
                scope = parsed_scope
            else:
                kept.append(token)
        return " ".join(kept).strip(), scope

    def highlight_search_terms(self, text, query):
        result = Text()
        terms = sorted(set(re.findall(r"\w+", query)), key=len, reverse=True)
        if not terms:
            result.append(text, style="verse.text")
            return result
        pattern = re.compile("(" + "|".join(re.escape(term) for term in terms) + ")", re.IGNORECASE)
        pos = 0
        for match in pattern.finditer(text):
            if match.start() > pos:
                result.append(text[pos:match.start()], style="verse.text")
            result.append(match.group(0), style="search.hit")
            pos = match.end()
        if pos < len(text):
            result.append(text[pos:], style="verse.text")
        return result

    def normalize_term(self, text):
        return re.sub(r"[^a-z0-9]+", "", text.lower())

    def normalize_strongs_key(self, key):
        match = re.match(r"([gh])0*(\d+)$", key.lower())
        if not match:
            return None, None, None
        prefix, num = match.groups()
        return f"{prefix}{int(num)}", f"{prefix.upper()}{int(num)}", f"{prefix.upper()}{int(num):04d}"

    def load_json_file(self, path):
        if not os.path.exists(path):
            return None
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)

    def load_manuscript_asset(self, relative_path):
        """Load a bundled/cached manuscript shard, fetching it only if absent."""
        if not relative_path or relative_path.startswith(('/', '\\')) or '..' in relative_path.split('/'):
            return None
        local_path = get_lex_path(relative_path)
        cache_path = os.path.join(MANUSCRIPT_CACHE_DIR, *relative_path.split('/'))
        for path in (local_path, cache_path):
            try:
                payload = self.load_json_file(path)
                if payload is not None:
                    return payload
            except (OSError, ValueError, json.JSONDecodeError):
                continue

        request = urllib.request.Request(
            MANUSCRIPT_WEB_BASE + relative_path,
            headers={"User-Agent": f"Lex/{VERSION}"},
        )
        try:
            with urllib.request.urlopen(request, timeout=15) as response:
                payload = json.loads(response.read().decode("utf-8"))
            os.makedirs(os.path.dirname(cache_path), exist_ok=True)
            temp_path = f"{cache_path}.{os.getpid()}.tmp"
            with open(temp_path, "w", encoding="utf-8") as cache_file:
                json.dump(payload, cache_file, ensure_ascii=False, separators=(",", ":"))
            os.replace(temp_path, cache_path)
            return payload
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            self._manuscript_error = str(exc)
            return None

    def get_interlinear_index(self):
        if self._interlinear_index is None:
            data = self.load_json_file(INTERLINEAR_PATH) or []
            self._interlinear_index = {}
            for row in data:
                ref = row.get("r")
                if not ref:
                    continue
                existing = self._interlinear_index.get(ref)
                # Some source rows are heading/context rows with the same ref.
                # Prefer rows with phrase data so study mode gets real tokens.
                if existing is None or (row.get("p") and not existing.get("p")):
                    self._interlinear_index[ref] = row
            self._ordered_refs = [
                row["r"] for row in data
                if row.get("r", "").startswith(f"{self.bible_prefix}:") and row.get("r", "").count(":") == 3 and not row.get("h")
            ]
        return self._interlinear_index

    def get_ordered_refs(self):
        if self._ordered_refs is None:
            # If interlinear data isn't available for this version, load refs from the DB
            if not os.path.exists(INTERLINEAR_PATH) or self.bible_prefix != "esv":
                res = self.bible_db.query("SELECT reference FROM bible WHERE reference NOT LIKE '%:0' ORDER BY id")
                self._ordered_refs = [row[0] for row in res]
            else:
                self.get_interlinear_index()
        return self._ordered_refs or []

    def get_reverse_naves(self, db_ref):
        parts = self.parse_reference_parts(db_ref)
        if not parts: return []

        book = self.reverse_canon_map.get(parts["book"], parts["book"])
        abbr = NAVES_BOOK_ABBR.get(book)
        if not abbr: return []

        # Format: "JHN 3:16" or "JHN 3:16-18"
        # We search for the specific verse in the entry text
        pattern = f"%{abbr} {parts['chapter']}:{parts['verse']}%"

        if not hasattr(self, "naves_db") or not self.naves_db:
            return []

        res = self.naves_db.query(
            "SELECT subject FROM topics WHERE entry LIKE ? ORDER BY subject",
            (pattern,)
        )
        return [row[0] for row in res]

    def get_interlinear_strongs(self):
        if self._interlinear_strongs is None:
            self._interlinear_strongs = self.load_json_file(INTERLINEAR_STRONGS_PATH) or {}
        return self._interlinear_strongs

    def get_step_greek(self):
        if self._step_greek is None:
            self._step_greek = self.load_json_file(STEP_GREEK_PATH) or {}
        return self._step_greek

    def get_step_hebrew(self):
        if self._step_hebrew is None:
            self._step_hebrew = self.load_json_file(STEP_HEBREW_PATH) or {}
        return self._step_hebrew

    def parse_history_ref(self, ref):
        if not ref:
            return None
        verse_match = re.match(r"^(?:[a-z0-9]+:)?(.+?):(\d+):(\d+)$", ref, re.IGNORECASE)
        if verse_match:
            book, chap, verse = verse_match.groups()
            return {"kind": "verse", "book": book, "chapter": int(chap), "verse": int(verse), "reference": ref}
        chapter_match = re.match(r"^(.*?)\s+(\d+)$", ref)
        if chapter_match:
            book, chap = chapter_match.groups()
            return {"kind": "chapter", "book": book, "chapter": int(chap), "reference": ref}
        return None

    def parse_reference_parts(self, db_ref):
        if ":" not in db_ref:
            return None
        version, ref_body = db_ref.split(":", 1)
        parts = ref_body.rsplit(":", 2)
        if len(parts) != 3:
            return None
        book, chapter, verse = parts
        return {
            "version": version,
            "book": book,
            "chapter": int(chapter),
            "verse": int(verse),
            "reference": db_ref,
        }

    def convert_to_tsk_ref(self, book, chapter, verse=None):
        prefix = TSK_BOOK_ABBR.get(book, book)
        if not prefix.endswith("."):
            prefix += "."
        return f"{prefix}{chapter}.{verse}" if verse else f"{prefix}{chapter}."

    def parse_tsk_ref_parts(self, tsk_ref):
        match = re.match(r"^([1-3]?[A-Za-z]+)\.(\d+)\.(\d+)$", tsk_ref)
        if not match:
            return None
        book_abbr, chapter, verse = match.groups()
        book = TSK_TO_BOOK.get(book_abbr)
        if not book:
            return None
        db_book = self.canon_map.get(re.sub(r"[^a-z0-9]+", "", book.lower()), book)
        return db_book, int(chapter), int(verse)

    def parse_tsk_ref(self, tsk_ref):
        first_ref = tsk_ref.split("-", 1)[0]
        parts = self.parse_tsk_ref_parts(first_ref)
        if not parts:
            return None
        book, chapter, verse = parts
        return f"{self.bible_prefix}:{book}:{int(chapter)}:{int(verse)}"

    def tsk_ref_to_range(self, tsk_ref):
        start_raw, end_raw = (tsk_ref.split("-", 1) + [None])[:2]
        start = self.parse_tsk_ref_parts(start_raw)
        if not start:
            return None
        end = self.parse_tsk_ref_parts(end_raw) if end_raw else start
        if not end:
            return start, start
        return start, end

    def format_tsk_display_ref(self, tsk_ref):
        ref_range = self.tsk_ref_to_range(tsk_ref)
        if not ref_range:
            return tsk_ref
        start, end = ref_range
        start_label = f"{self.reverse_canon_map.get(start[0], start[0])} {start[1]}:{start[2]}"
        end_book = self.reverse_canon_map.get(end[0], end[0])
        if start == end:
            return start_label
        if start[0] == end[0] and start[1] == end[1]:
            return f"{start_label}-{end[2]}"
        return f"{start_label}-{end_book} {end[1]}:{end[2]}"

    def get_tsk_crossrefs(self, db_ref):
        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return []
        tsk_ref = self.convert_to_tsk_ref(parts["book"], parts["chapter"], parts["verse"])
        return self.cross_refs_db.query(
            "SELECT to_ref, votes FROM cross_refs WHERE from_ref = ? ORDER BY votes DESC, to_ref",
            (tsk_ref,)
        )

    def get_crossref_text(self, tsk_ref):
        ref_range = self.tsk_ref_to_range(tsk_ref)
        if not ref_range:
            return None
        start, end = ref_range
        start_ref = f"{self.bible_prefix}:{start[0]}:{start[1]}:{start[2]}"
        end_ref = f"{self.bible_prefix}:{end[0]}:{end[1]}:{end[2]}"
        rows = self.bible_db.query(
            """
            SELECT reference, text FROM bible
            WHERE id >= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
              AND id <= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
            ORDER BY id
            """,
            (start_ref, end_ref),
        )
        if not rows:
            return None
        verses = []
        for ref, text in rows:
            parts = self.parse_reference_parts(ref)
            verse_text = self.clean_text(text)
            if parts and len(rows) > 1:
                verses.append(f"{parts['verse']}. {verse_text}")
            else:
                verses.append(verse_text)
        return " ".join(verses)

    def get_crossref_preview(self, tsk_ref):
        return self.get_crossref_text(tsk_ref)

    def get_navigation_reference(self, current_ref, direction):
        refs = self.get_ordered_refs()
        try:
            idx = refs.index(current_ref)
        except ValueError:
            return None
        target_idx = idx + (1 if direction == "next" else -1)
        if target_idx < 0 or target_idx >= len(refs):
            return None
        return refs[target_idx]

    def get_adjacent_chapter_reference(self, book, chapter, direction):
        chapters = []
        seen = set()
        for ref in self.get_ordered_refs():
            parts = self.parse_reference_parts(ref)
            if not parts:
                continue
            key = (parts["book"], parts["chapter"])
            if key not in seen:
                seen.add(key)
                chapters.append(key)
        try:
            idx = chapters.index((book, chapter))
        except ValueError:
            return None
        target_idx = idx + (1 if direction == "next" else -1)
        if target_idx < 0 or target_idx >= len(chapters):
            return None
        target_book, target_chapter = chapters[target_idx]
        for ref in self.get_ordered_refs():
            parts = self.parse_reference_parts(ref)
            if parts and parts["book"] == target_book and parts["chapter"] == target_chapter:
                return ref
        return None

    def resolve_navigation_query(self, direction):
        parsed = self.parse_history_ref(self.last_ref)
        if not parsed:
            return None
        if parsed["kind"] == "verse":
            ref = self.get_navigation_reference(parsed["reference"], direction)
            if not ref:
                return None
            parts = self.parse_reference_parts(ref)
            return f"{parts['book']} {parts['chapter']}:{parts['verse']}" if parts else None
        ref = self.get_adjacent_chapter_reference(parsed["book"], parsed["chapter"], direction)
        if not ref:
            return None
        parts = self.parse_reference_parts(ref)
        return f"{parts['book']} {parts['chapter']}" if parts else None

    def normalize_ref(self, q):
        # User-facing references are intentionally forgiving here. The DB still
        # uses canonical "version:Book:Chapter:Verse" strings internally.
        q_clean = q.lower().strip()

        # 1. Try standard spaced reference: '1 John 3:16', 'John 3:16', '1 Jn 3', 'John 3:16-20'
        # Book is anything before the last set of digits
        pattern_spaced = r'^([1-4]?\s?[a-z\s.]+?)\s+(\d+)(?:[\s:.](\d+)(?:-(\d+))?)?$'
        match = re.match(pattern_spaced, q_clean)

        # 2. Try compact reference: '1cor13', 'jn3:16', 'jn3:16-20'
        if not match:
            pattern_compact = r'^([1-4]?[a-z\s.]+?)(\d+)(?:[\s:.](\d+)(?:-(\d+))?)?$'
            match = re.match(pattern_compact, q_clean)

        if match:
            groups = match.groups()
            b = groups[0].strip()
            c = groups[1]
            v = groups[2] if len(groups) > 2 else None
            v_end = groups[3] if len(groups) > 3 else None

            # For slugs and compact keys, we want to normalize the book part
            b_slug = re.sub(r"[^a-z0-9]+", "-", b.lower()).strip("-")
            b_compact = re.sub(r"[^a-z0-9]+", "", b.lower())

            # 1. Try dynamic canon map from active DB
            b_name = self.canon_map.get(b_compact) or self.canon_map.get(b_slug)

            # 2. Try static aliases
            if not b_name:
                b_name = BOOK_SCOPE_ALIASES.get(b_compact) or BOOK_SCOPE_ALIASES.get(b_slug)

            # 2b. Try LXX/deuterocanonical labels.
            lxx_name = LXX_REFERENCE_BOOK_ALIASES.get(b_compact) or LXX_REFERENCE_BOOK_ALIASES.get(b_slug)
            if lxx_name and self.bible_prefix == "lxx":
                b_name = lxx_name
            elif not b_name:
                b_name = lxx_name

            # 3. Handle cases where the number might be separate in our mapping: '1-john' vs '1john'
            if not b_name and b_slug.startswith(("1-", "2-", "3-")):
                alt_compact = b_slug.replace("-", "", 1)
                b_name = self.canon_map.get(alt_compact) or BOOK_SCOPE_ALIASES.get(alt_compact)

            # 4. Fallback to title case
            if not b_name:
                b_name = b.title()

            # Final resolution to ensure we use the exact DB book name
            b_name_key = re.sub(r"[^a-z0-9]+", "", b_name.lower())
            b_name = self.canon_map.get(b_name_key, b_name)

            ref_norm = f"{b_name}:{c}"
            if v:
                ref_norm += f":{v}"

            return ref_norm, b_name, c, v, v_end

        return None, None, None, None, None

    # -----------------------------------------------------------------------
    # Landing pages, help, and credits
    # -----------------------------------------------------------------------
    def display_intro(self):
        logo = Text(
            r"""
██╗     ███████╗██╗  ██╗
██║     ██╔════╝╚██╗██╔╝
██║     █████╗   ╚███╔╝
██║     ██╔══╝   ██╔██╗
███████╗███████╗██╔╝ ██╗
╚══════╝╚══════╝╚═╝  ╚═╝
""",
            style="bold gold3",
        )
        title = Text("Lex: The Elegant Bible Terminal", style="text.strong")
        tagline = Text("Master Admin Study Tool for the Source Code of the Universe", style="bold cyan")
        positioning = Text(
            "Read the canon. Inspect the languages. Traverse the tradition.",
            style="dim",
        )

        metrics = Table.grid(padding=(0, 2))
        metrics.add_column(justify="center", style="bold gold3")
        metrics.add_column(justify="center", style="bold green")
        metrics.add_column(justify="center", style="bold cyan")
        metrics.add_column(justify="center", style="bold magenta")
        metrics.add_row("66 books", "TSK graph", "Strong's + STEPBible", "Creeds + ISBE")

        primary = Table.grid(padding=(0, 2))
        primary.add_column(style="bold green", no_wrap=True)
        primary.add_column(style="text")
        primary.add_row("Read:", "lex read John 3:16  (Context with navigation)")
        primary.add_row("Study:", "lex study John 3:16  (Interlinear + lexicon)")
        primary.add_row("Search:", 'lex search "mustard seed"  (Ranked search results)')
        primary.add_row("Manuscript Map:", "lex manuscript <verse|name>  (Readings, witnesses, profile)")
        primary.add_row("Strong's Lookup:", "lex strongs love  or  lex G3056")

        also = Table.grid(padding=(0, 2))
        also.add_column(style="bold gold3", no_wrap=True)
        also.add_column(style="text")
        also.add_row("Quick Read:", "lex John 3:16")
        also.add_row("Quick Study:", "lex John 3:16 -i")
        also.add_row("Verse Web:", "lex web John 3:16")
        also.add_row("Commentary:", "lex commentary John 3:16")
        also.add_row("Export:", "lex export John 3:16  (Interactive export menu)")
        also.add_row("Nave's Topics:", "lex naves grace  or  lex naves John 3:16")
        also.add_row("Lexicon:", "lex G3056  or  lex logos")
        also.add_row("Creeds:", "lex creed")
        also.add_row("Define:", "lex define grace")
        also.add_row("Versions:", "lex -v  (Interactive version menu)")
        also.add_row("Switch Bible:", "lex version kjv  or  lex -v nasb")
        also.add_row("Select Once:", "lex -B lxx John 1:1")

        config = Table.grid(padding=(0, 2))
        config.add_column(style="ui.action.key", no_wrap=True)
        config.add_column(style="text")
        config.add_row("lex -light", "Use and remember the light terminal theme")
        config.add_row("lex -dark", "Use and remember the dark terminal theme")
        config.add_row("lex -auto", "Clear the saved theme and auto-detect again")
        config.add_row("lex history", "Show recent local Lex commands with a compact cheat sheet")
        config.add_row("LEX_THEME=light lex", "Use a theme for one shell command")
        config.add_row("LEX_NO_COLOR=1 lex", "Print plain output without Lex colors")

        launch = Table.grid(padding=(0, 1))
        launch.add_column(style="dim")
        launch.add_column(style="bold gold3")
        launch.add_column(style="dim")
        launch.add_row("MODE", "LOCAL-FIRST", "No browser tabs. No drift. Just the sources.")

        nav = Text()
        nav.append("Quick Navigation: ", style="bold cyan")
        nav.append("lex --next", style="gold3")
        nav.append(" | ")
        nav.append("lex --prev", style="gold3")
        nav.append("  (Relative to your last read reference)", style="dim")

        credits = Table.grid(padding=(0, 2))
        credits.add_column(style="bold cyan", no_wrap=True)
        credits.add_column(style="dim")
        credits.add_row("Credits:", "ESV text, TSK/OpenBible, Strong's, STEPBible, UBS, Easton, ISBE, TheologAI historical docs")
        credits.add_row("License:", "Lex code MIT; data remains under source terms. Run lex --credits")

        footer = Text("Start with a verb, or type a reference directly.", style="italic dim")
        console.print(
            Panel(
                Group(
                    Align.center(logo),
                    Align.center(title),
                    Align.center(tagline),
                    Align.center(positioning),
                    "",
                    Align.center(metrics),
                    "",
                    Text("Start with a verb.", style="text.strong"),
                    "",
                    Text("Primary", style="bold cyan"),
                    primary,
                    "",
                    Text("Also Available", style="bold cyan"),
                    also,
                    "",
                    Text("Config", style="bold cyan"),
                    config,
                    "",
                    nav,
                    "",
                    Align.center(launch),
                    "",
                    credits,
                    Align.center(footer),
                ),
                title=f"Lex {VERSION}",
                subtitle="source-aware bible study, shipped as a command",
                border_style="bold cyan",
                padding=(1, 3),
                expand=False,
            )
        )

    def display_credits(self):
        table = Table(title="Lex Credits and Data Licenses", box=None, show_lines=True)
        table.add_column("Component", style="bold cyan", no_wrap=True)
        table.add_column("Source / Repo", style="text", overflow="fold")
        table.add_column("License / Terms", style="gold3", overflow="fold")
        table.add_row(
            "Lex CLI code",
            "Local project code: /home/n8te/lex_v3.py",
            "Recommended: MIT for application code only",
        )
        table.add_row(
            "Bible text",
            "Local bible-data / ESV-derived SQLite: bible_versions/esv.db; source package notes: bible-data",
            "Permission/copyright-controlled translation text; do not relicense as MIT",
        )
        table.add_row(
            "TSK cross refs",
            "Treasury of Scripture Knowledge / OpenBible-style cross-reference data",
            "Verify upstream terms before redistribution",
        )
        table.add_row(
            "Strong's lexicon",
            "OpenScriptures Strong's Hebrew and Greek Dictionaries: github.com/openscriptures/strongs",
            "Local XHTML says GPL-3.0; another local source note says Public Domain. Verify source chain before distribution",
        )
        table.add_row(
            "STEPBible language data",
            "STEPBible Data: github.com/STEPBible/STEPBible-Data; www.STEPBible.org",
            "CC BY 4.0; credit STEP Bible",
        )
        table.add_row(
            "UBS resources",
            "UBS Open License resources: local ubs-open-license dataset",
            "CC BY-SA 4.0; preserve attribution and ShareAlike obligations",
        )
        table.add_row(
            "Bible geography",
            "OpenBible Bible Geocoding Data: openbible.info/geo",
            "CC BY 4.0; some map/image data may carry ODbL or separate CC terms",
        )
        table.add_row(
            "Dictionary",
            "Easton's Bible Dictionary entries in lexicon.db",
            "Public domain",
        )
        table.add_row(
            "Encyclopedia",
            "International Standard Bible Encyclopedia OCR import; currently local Volume II Clement-Heresh",
            "Public domain source; OCR/import quality and volume coverage still in progress",
        )
        table.add_row(
            "Creeds/confessions",
            "TheologAI historical documents dataset: local theolog-ai/data/historical-documents",
            "Public domain per local TheologAI README; preserve source attribution",
        )
        table.add_row(
            "Interlinear data",
            "Local esv-data interlinear + STEPBible/Strong's-backed resources",
            "Mixed source terms; preserve Bible text, STEPBible, and Strong's source obligations",
        )

        note = Markdown(
            """
**Recommended licensing model:** MIT for Lex application code; source-specific terms for all data.

For redistribution, include upstream license files and a `NOTICE`/`DATA_LICENSES.md`.
Do not represent the ESV text, UBS resources, STEPBible data, or generated databases as MIT-licensed.

See: `~/bible-lexicon-data/docs/LICENSING.md`
"""
        )
        console.print(Panel(Group(table, "", note), border_style="ui.border", padding=(1, 2)))

    def display_vulgate_study(self, db_ref, animate=None):
        """Display Vulgate text with Latin semantic tagging from latin.db"""
        if not self.latin_db:
            return False

        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return False

        # Get Vulgate text
        text_row = self.bible_db.query("SELECT text FROM bible WHERE reference = ?", (db_ref,))
        if not text_row:
            return False
        vulgate_text = text_row[0][0]

        self.pause_study_section(animate)
        console.print(
            Panel(
                Text(vulgate_text, style="source.text"),
                title=f"Vulgate (Clementine) • {db_ref}",
                border_style="source.border",
                padding=(1, 2),
            )
        )

        # Simple word lookup from latin.db
        words = re.findall(r'\b\w+\b', vulgate_text.lower())
        table = Table(title="Latin Semantic Tags", box=None, expand=True)
        table.add_column("Latin", style="source.text", no_wrap=True)
        table.add_column("POS", style="dim", width=8)
        table.add_column("Definition", style="text", overflow="fold")

        seen = set()
        for word in words[:12]:  # limit for display
            if word in seen:
                continue
            seen.add(word)
            entry = self.latin_db.query(
                "SELECT headword, part_of_speech, definition FROM latin WHERE normalized = ? LIMIT 1",
                (word,)
            )
            if entry:
                headword, pos, definition = entry[0]
                definition = definition[:80] + "..." if len(definition) > 80 else definition
                table.add_row(headword, pos or "-", definition)

        if table.rows:
            console.print(table)
        else:
            console.print(Panel("No Latin dictionary entries found for these words.", border_style="dim"))

        return True

    def display_study_landing(self):
        md = """
# Lex Study
*Interlinear reading without leaving the terminal*

Study mode aligns the English verse with the source text, transliteration, lemma,
morphology, and Strong's-backed lexicon notes.

**What It Shows**

*   **Verse Context:** a compact read panel around the target verse
*   **Source Alignment:** English phrase, source token, lemma, and code
*   **Lexicon Notes:** Strong's and STEPBible definitions for Greek, Hebrew, and Aramaic
*   **Navigation:** read a verse, then move with `lex --prev` and `lex --next`

**Try These**

*   `lex study John 1:1`
*   `lex study Genesis 1:1`
*   `lex study Daniel 2:4`
*   `lex John 3:16 -i`

---
*Read the text. Inspect the words. Stay in one tool.*
"""
        console.print(Panel(Markdown(md), title="🔤 Study Mode", border_style="ui.border", expand=False))

    def display_read_landing(self):
        md = """
# Lex Read
*Scripture reading with fast terminal navigation*

Read mode centers a passage in context and keeps your place for `lex --prev`
and `lex --next`.

**What It Does**

*   **Verse View:** shows the target verse with nearby context
*   **Chapter View:** prints the full chapter in order
*   **History:** saves your last reading position for navigation
*   **Bridge to Study:** jump from reading into analysis with `lex study ...`

**Try These**

*   `lex read John 3:16`
*   `lex jn 1:1`
*   `lex study rev 1:2`
*   `lex 2 jn 1:2`
*   `lex read Genesis 1`
*   `lex John 1:1`
*   `lex --next`

---
*Open the text fast. Move without friction. Study when needed.*
"""
        console.print(Panel(Markdown(md), title="📖 Read Mode", border_style="ui.border", expand=False))

    def display_search_howto(self):
        md = """
# Search Help

Use explicit search mode:

*   `lex search "mustard seed"`
*   `lex search kingdom heaven`

Search ranks an exact phrase first, then includes verses containing all words.
If nothing matches, Lex corrects likely misspellings and, as a final fallback,
shows verses containing any of the words. Corrections are always shown.

Search only matches verse text, not book names stored in reference metadata.
Use a scope when you mean a book: `lex search light -john`.

## Page Controls

*   `lex search covenant --page 2`
*   `lex search covenant --limit 20`

In an interactive terminal, search opens a compact action bar:

*   `1`, `2`, `3` - study that numbered result
*   `r 1`, `r 2` - read that numbered result
*   `n` / `p` - next or previous page
*   `e` - export menu
*   `q` - quit

Export menu:

*   `d` - DOCX
*   `f` - PDF
*   `p` - PowerPoint
*   `o` - open exports folder
*   `q` - back

Exports are saved under `~/Documents/lex_exports` and Lex tries to open them after saving.

## Book Scopes

Add a single-dash scope after the search term:

*   `lex search covenant -jeremiah`
*   `lex search beast -daniel-revelation`
*   `lex search covenant -major`
*   `lex search resurrection -nt`

Book ranges follow canonical order, so `-jeremiah-revelation` searches from
Jeremiah through Revelation. Book names use lowercase words joined by hyphens:

*   `-song-of-solomon`
*   `-1-john`
*   `-1-corinthians-2-corinthians`

## Group Scopes

*   `-ot` / `-old-testament`
*   `-nt` / `-new-testament`
*   `-law` / `-pentateuch` / `-torah`
*   `-history`
*   `-wisdom` / `-poetry`
*   `-major` / `-major-prophets`
*   `-minor` / `-minor-prophets`
*   `-prophets`
*   `-gospels`
*   `-epistles` / `-letters`
*   `-pauline`
*   `-general-epistles`

Free-text search no longer runs from bare input.

## Terminal Theme Config

Lex tries to detect whether your terminal background is light or dark. You can
override it when the automatic choice is wrong:

*   `lex -light` - switch to light mode and remember it
*   `lex -dark` - switch to dark mode and remember it
*   `lex -auto` - remove the saved manual setting and detect again

The saved setting lives in `~/.lex_config.json`. If you choose `-light` or
`-dark`, Lex will keep using that theme on future launches until you run
`lex -auto` or choose the other theme.

For one command only, use an environment variable:

*   `LEX_THEME=light lex John 3:16`
*   `LEX_THEME=dark lex search covenant`

If you need plain text with no Lex colors:

*   `LEX_NO_COLOR=1 lex John 3:16`
"""
        console.print(Panel(Markdown(md), title="🔎 Search", border_style="ui.border", expand=False))

    def display_topic_howto(self):
        md = """
# Nave's Topical Bible

Browse over 20,000 topics and 100,000 scripture references. This is a
re-implementation of the classic Nave's index using a high-signal
blueprint layout.

## Basic Usage

*   `lex topic church`
*   `lex topic "second coming"` (Quotes required for multiple words)
*   `lex naves grace` (Alias for the same command)

## Deep Search

If a topic title isn't found exactly, Lex will automatically perform a
semantic full-text search across the content of all entries to find
relevant thematic clusters.
"""
        console.print(Markdown(md))

    def display_commentary_howto(self):
        md = """
# Biblical Commentaries

Traverse the Christian tradition with Matthew Henry and John Calvin.
Commentary lookups provide historical context and theological depth
alongside your scripture study.

## Basic Usage

*   `lex commentary John 3:16`
*   `lex commentary Romans 1` (Fetches entire chapter notes)
*   `lex commentary gen 1:1` (Supports standard abbreviations)

## High-Signal View

Notes from multiple traditions are displayed in sequential themed blocks
(Henry in Blue, Calvin in Magenta) for immediate visual comparison.
"""
        console.print(Markdown(md))

    def display_manuscript_howto(self):
        md = """
# Manuscript Map

Look up a Bible verse to compare its available manuscript readings, or look up
a manuscript by semantic name / Gregory-Aland number:

*   `lex manuscript John 1:1`
*   `lex manuscript Isaiah 53:11`
*   `lex manuscript P66`
*   `lex manuscript 1Qisaa`

Use `--limit 25` to show more readings and witnesses. Shards are read from the
local data pack or cache first; missing shards are fetched from Lex Web and
cached for later offline use.
"""
        console.print(Panel(Markdown(md), title="Manuscripts", border_style="ui.border", expand=False))

    def manuscript_verse_path(self, query):
        ref_norm, book, chapter, verse, _ = self.normalize_ref(query)
        if not ref_norm or not verse:
            return None, None
        canonical_book = self.reverse_canon_map.get(book, book)
        if canonical_book == "Psalms":
            canonical_book = "Psalm"
        slug = re.sub(r"[^a-z0-9]+", "-", canonical_book.lower()).strip("-")
        collection = "verses" if canonical_book in MANUSCRIPT_NT_BOOKS else "ot-verses"
        reference = f"{canonical_book} {int(chapter)}:{int(verse)}"
        return f"witnesses/{collection}/{slug}/{slug}-{int(chapter)}-{int(verse)}.json", reference

    def render_manuscript_verse(self, payload, limit=10):
        coverage = payload.get("coverage") or {}
        reference = payload.get("reference") or "Manuscript readings"
        summary = Table.grid(padding=(0, 2))
        summary.add_column(style="ui.meta", no_wrap=True)
        summary.add_column(style="text", overflow="fold")
        summary.add_row("Reference", reference)
        summary.add_row(
            "Coverage",
            f"{coverage.get('unique_witnesses', 0)} unique witnesses · "
            f"{coverage.get('reading_groups', 0)} reading groups · "
            f"{coverage.get('rendered_instances', 0)} rendered instances",
        )
        source_rows = payload.get("sources") or ([payload.get("source")] if payload.get("source") else [])
        source_label = " · ".join(
            f"{source.get('name', 'Unknown source')} ({source.get('license') or source.get('license_note') or 'license noted upstream'})"
            for source in source_rows
        )
        if source_label:
            summary.add_row("Source", source_label)
        console.print(Panel(summary, title=f"Manuscripts · {reference}", border_style="source.border"))

        base_text = payload.get("base_text") or {}
        if base_text.get("text"):
            console.print(Panel(
                Text(base_text["text"], style="source.text"),
                title=base_text.get("label") or "Collation base",
                subtitle=base_text.get("note") or None,
                border_style="source.border",
            ))

        readings = payload.get("readings") or []
        if readings:
            table = Table(title="Reading groups", border_style="ui.border", expand=True)
            table.add_column("#", justify="right", style="ui.meta", width=3)
            table.add_column("Reading", style="source.text", overflow="fold")
            table.add_column("Witnesses", style="text", overflow="fold")
            table.add_column("Flags", style="warning", overflow="fold")
            for reading in readings[:limit]:
                witnesses = []
                for witness in reading.get("witnesses") or []:
                    label = witness.get("label") or witness.get("id") or "?"
                    instance = witness.get("instance")
                    witnesses.append(f"{label}({instance})" if instance and instance != "1" else label)
                table.add_row(
                    str(reading.get("count") or 0),
                    reading.get("reading") or reading.get("display") or "-",
                    ", ".join(witnesses) or "-",
                    ", ".join(reading.get("flags") or []) or "-",
                )
            console.print(table)
            if len(readings) > limit:
                console.print(f"[ui.meta]Showing {limit} of {len(readings)} reading groups; rerun with --limit {min(50, len(readings))}.[/]")

        witnesses = payload.get("witnesses") or []
        if witnesses:
            table = Table(title="Witness map", border_style="ui.border", expand=True)
            table.add_column("Witness", style="verse.ref", no_wrap=True)
            table.add_column("Tradition", style="ui.meta", no_wrap=True)
            table.add_column("Reading", style="source.text", overflow="fold")
            table.add_column("Flags", style="warning", overflow="fold")
            for witness in witnesses[:limit]:
                label = witness.get("label") or witness.get("id") or "?"
                instance = witness.get("instance")
                if instance and instance != "1":
                    label += f" ({instance})"
                table.add_row(
                    label,
                    witness.get("tradition") or "-",
                    witness.get("plain") or witness.get("display") or "-",
                    ", ".join(witness.get("flags") or []) or "-",
                )
            console.print(table)
            if len(witnesses) > limit:
                console.print(f"[ui.meta]Showing {limit} of {len(witnesses)} witnesses.[/]")
        for note in payload.get("notes") or []:
            console.print(f"[ui.meta]{note}[/]")
        return True

    def render_manuscript_profile(self, payload, limit=10):
        profile = payload.get("profile") or {}
        name = profile.get("primaryName") or profile.get("gaNum") or profile.get("id") or "Manuscript"
        details = Table.grid(padding=(0, 2))
        details.add_column(style="ui.meta", no_wrap=True)
        details.add_column(style="text", overflow="fold")
        details.add_row("Semantic name", name)
        details.add_row("Internal ID", str(profile.get("id") or "-"))
        date = " – ".join(str(value) for value in (profile.get("origEarly"), profile.get("origLate")) if value)
        details.add_row("Date", date or "-")
        details.add_row("Language", profile.get("language") or "-")
        if profile.get("tradition"):
            details.add_row("Tradition", profile["tradition"])
        external = profile.get("external") or {}
        if external.get("intfWorkspace"):
            details.add_row("INTF", external["intfWorkspace"])
        if external.get("sourceRepository"):
            details.add_row("Source", external["sourceRepository"])
        console.print(Panel(details, title=f"Manuscript · {name}", border_style="source.border"))

        shelves = profile.get("shelfInstances") or []
        if shelves:
            table = Table(title="Shelf instances", border_style="ui.border", expand=True)
            table.add_column("Institution", style="text")
            table.add_column("Shelf", style="verse.ref")
            table.add_column("Place", style="ui.meta")
            table.add_column("Contents", style="text", overflow="fold")
            for shelf in shelves[:limit]:
                table.add_row(
                    shelf.get("institution") or "-",
                    shelf.get("shelfNumber") or "-",
                    ", ".join(value for value in (shelf.get("place"), shelf.get("country")) if value) or "-",
                    shelf.get("contentOverview") or shelf.get("leaves") or "-",
                )
            console.print(table)

        coverage = payload.get("coverage") or {}
        refs = coverage.get("poc_references") or []
        if refs:
            console.print(Panel(
                "  ·  ".join(refs[:limit]),
                title=f"Mapped references · {coverage.get('poc_reference_count', len(refs))}",
                border_style="ui.border",
            ))
        samples = payload.get("samples") or []
        if samples:
            table = Table(title="Sample readings", border_style="ui.border", expand=True)
            table.add_column("Reference", style="verse.ref", no_wrap=True)
            table.add_column("Reading", style="source.text", overflow="fold")
            table.add_column("Flags", style="warning", overflow="fold")
            for sample in samples[:limit]:
                table.add_row(
                    sample.get("reference") or "-",
                    sample.get("plain") or sample.get("display") or "-",
                    ", ".join(sample.get("flags") or []) or "-",
                )
            console.print(table)
        for link in external.get("links") or []:
            console.print(f"[ui.meta]{link.get('label', 'External resource')}: {link.get('url', '')}[/]")
        for note in payload.get("notes") or []:
            console.print(f"[ui.meta]{note}[/]")
        return True

    def display_manuscript(self, query, limit=10):
        path, reference = self.manuscript_verse_path(query)
        if path:
            payload = self.load_manuscript_asset(path)
            if payload:
                return self.render_manuscript_verse(payload, limit=min(max(limit, 1), 50))
            console.print(f"[warning]No manuscript shard is available for {reference}.[/]")
            if getattr(self, "_manuscript_error", None):
                console.print("[ui.meta]The local cache was empty and Lex Web could not be reached.[/]")
            return False

        manifest = self.load_manuscript_asset("witnesses-profiles-manifest.json")
        key = re.sub(r"^ga\s*", "", query.strip(), flags=re.IGNORECASE).lower()
        profile_path = (manifest or {}).get("manuscripts", {}).get(key)
        if not profile_path:
            console.print(f"[warning]No manuscript profile matches {query}. Try P66, 1Qisaa, WLC, or a Bible reference.[/]")
            return False
        payload = self.load_manuscript_asset(profile_path)
        if not payload:
            console.print(f"[warning]The manuscript profile for {query} could not be loaded.[/]")
            return False
        return self.render_manuscript_profile(payload, limit=min(max(limit, 1), 50))

    def display_strongs_howto(self):
        md = """
# Strong's Lookup

Find Strong's entries by number, transliteration, or English gloss:

*   `lex strongs love`
*   `lex strongs word`
*   `lex strongs God`
*   `lex G3056`
"""
        console.print(Panel(Markdown(md), title="🔤 Strong's Lookup", border_style="ui.border", expand=False))

    # -----------------------------------------------------------------------
    # Bible reading and navigation rendering
    # -----------------------------------------------------------------------
    def format_display_ref(self, db_ref):
        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return db_ref
        return f"{parts['book']} {parts['chapter']}:{parts['verse']}"

    def display_read_nav(self, book, chap, verse=None):
        ref = f"{book} {chap}:{verse}" if verse else f"{book} {chap}"
        study_ref = f"{book} {chap}:{verse or 1}"
        console.print(fill_terminal_row(f"lex --prev  |  lex --next  |  lex study {study_ref}  |  lex export {ref}", "ui.meta"))

    def should_animate(self, animate):
        if animate is not None:
            return animate
        return console.is_terminal and not os.environ.get("NO_COLOR")

    def render_verse_context(self, rows, target_ref, book, chap, verse):
        body = Text()
        for _, ref, text in rows:
            parts = self.parse_reference_parts(ref)
            verse_no = str(parts["verse"]) if parts else self.format_display_ref(ref)
            is_target = ref == target_ref
            marker = ">" if is_target else " "
            label_style = "verse.marker" if is_target else "verse.ref.muted"
            text_style = "verse.text.focus" if is_target else "verse.text.muted"
            body.append(f"{marker} {verse_no.rjust(3)} ", style=label_style)
            body.append(f"{self.clean_text(text)}\n", style=text_style)
        console.print(
            Panel(
                body,
                title=f"📖 {book} {chap}:{verse}",
                subtitle="context",
                border_style="verse.border",
                padding=(1, 2),
            )
        )
        self.display_read_nav(book, chap, verse)

    def render_chapter(self, rows, book, chap):
        body = Text()
        for ref, text in rows:
            parts = self.parse_reference_parts(ref)
            verse_no = str(parts["verse"]) if parts else self.format_display_ref(ref)
            body.append(f"{verse_no.rjust(3)} ", style="verse.ref")
            body.append(f"{self.clean_text(text)}\n\n", style="verse.text")
        console.print(
            Panel(
                body,
                title=f"📖 {book} {chap}",
                subtitle=f"{len(rows)} verses",
                border_style="verse.border",
                padding=(1, 2),
            )
        )
        self.display_read_nav(book, chap)

    def display_verse(self, query, interlinear=False, animate=None):
        ref_norm, book, chap, verse, v_end = self.normalize_ref(query)
        if not ref_norm: return False

        # Ensure we use the DB-specific book name in the LIKE query
        db_book = self.canon_map.get(re.sub(r"[^a-z0-9]+", "", book.lower()), book)

        if v_end:
            # Handle range
            start_ref = f"{self.bible_prefix}:{db_book}:{chap}:{verse}"
            end_ref = f"{self.bible_prefix}:{db_book}:{chap}:{v_end}"
            res = self.bible_db.query(
                """
                SELECT reference, text
                FROM bible
                WHERE id >= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
                  AND id <= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
                ORDER BY id
                """,
                (start_ref, end_ref)
            )
            if res:
                self.render_chapter(res, book, f"{chap}:{verse}-{v_end}")
                return True

        if verse:
            res = self.bible_db.query(
                "SELECT MIN(id), reference, text FROM bible WHERE reference LIKE ? GROUP BY reference LIMIT 1",
                (f"%:{db_book}:{chap}:{verse}",)
            )
        else:
            res = self.bible_db.query(
                """
                SELECT reference, text
                FROM bible
                WHERE id IN (
                    SELECT MIN(id)
                    FROM bible
                    WHERE reference LIKE ? AND reference NOT LIKE '%:0'
                    GROUP BY reference
                )
                ORDER BY id
                """,
                (f"%:{db_book}:{chap}:%",)
            )
        if res:
            if verse:
                target_id, ref, text = res[0]
                target_parts = self.parse_reference_parts(ref)
                context_ids = []
                current = ref
                prev2 = self.get_navigation_reference(current, "prev")
                prev1 = self.get_navigation_reference(prev2, "prev") if prev2 else None
                next1 = self.get_navigation_reference(current, "next")
                next2 = self.get_navigation_reference(next1, "next") if next1 else None
                for candidate in [prev1, prev2, current, next1, next2]:
                    if candidate:
                        candidate_parts = self.parse_reference_parts(candidate)
                        if target_parts and candidate_parts and candidate_parts["book"] != target_parts["book"]:
                            continue
                        row = self.bible_db.query(
                            "SELECT MIN(id), reference, text FROM bible WHERE reference = ? GROUP BY reference",
                            (candidate,)
                        )
                        if row:
                            context_ids.append(row[0])
                self.render_verse_context(context_ids, ref, book, chap, verse)
                if interlinear:
                    self.display_study(ref, animate=animate, context=False)
                self.save_history(ref)
            else:
                self.render_chapter(res, book, chap)
                self.save_history(f"{book} {chap}")
            return True
        return False
    def display_verse_web(self, query, limit=12):
        ref_norm, book, chap, verse, v_end = self.normalize_ref(query)
        if not ref_norm or not verse:
            return False

            console.print("[warning]Verse web needs a single verse, e.g. lex web John 3:16[/]")
            return False
        rows = self.bible_db.query(
            "SELECT MIN(id), reference, text FROM bible WHERE reference LIKE ? GROUP BY reference LIMIT 1",
            (f"%:{book}:{chap}:{verse}",)
        )
        if not rows:
            return False
        _, db_ref, verse_text = rows[0]
        clean_verse = self.clean_text(verse_text)
        refs = self.get_tsk_crossrefs(db_ref)[:max(1, min(limit, 24))]

        center = Text()
        center.append(f"{book} {chap}:{verse}\n", style="verse.ref")
        center.append(clean_verse, style="text.strong")

        console.print(
            Panel(
                Align.center(center),
                title="✦ Scripture Web ✦",
                subtitle="ranked local TSK connections",
                border_style="verse.border",
                padding=(1, 2),
            )
        )

        if not refs:
            console.print("[warning]No local cross-reference links found for this verse.[/]")
            return True

        table = Table(title="Major Connections", box=None, expand=True)
        table.add_column("Rank", style="dim", justify="right", no_wrap=True)
        table.add_column("Link", style="verse.ref", no_wrap=True)
        table.add_column("Weight", style="bold cyan", justify="right", no_wrap=True)
        table.add_column("Text", style="verse.text", overflow="fold")
        for idx, (to_ref, votes) in enumerate(refs, 1):
            text = self.get_crossref_text(to_ref)
            table.add_row(str(idx), self.format_tsk_display_ref(to_ref), str(votes), text or "")
        console.print(table)

        spark = Text()
        for idx, (to_ref, votes) in enumerate(refs[:8], 1):
            if idx > 1:
                spark.append("  ", style="dim")
            spark.append("●", style="gold3" if idx == 1 else "cyan")
            spark.append(f" {self.format_tsk_display_ref(to_ref)}", style="dim")
        console.print(Panel(spark, title="Connection Trail", border_style="ui.border"))
        console.print(f"[dim]Open a link: lex read <ref>  |  Read top 10 refs: lex read all  |  Study center: lex study {book} {chap}:{verse}[/]")
        self.save_history(db_ref)
        return True

    # -----------------------------------------------------------------------
    # Study mode: source text, interlinear rows, lexicons, and TSK links
    # -----------------------------------------------------------------------
    def lookup_lexicon_entry(self, strongs_id):
        short_key, strongs_db_key, step_key = self.normalize_strongs_key(strongs_id)
        interlinear = self.get_interlinear_strongs().get(short_key) if short_key else None
        step = self.get_step_greek().get(step_key) if strongs_id.lower().startswith("g") else self.get_step_hebrew().get(step_key)
        db = self.strongs_db.query("SELECT number, word, pronunciation, definition FROM strongs WHERE number = ?", (strongs_db_key,)) if strongs_db_key else []
        return {
            "interlinear": interlinear,
            "step": step,
            "db": db[0] if db else None,
        }

    def extract_english_glosses(self, entry):
        if not entry:
            return []
        raw = entry.get("r", "")
        if "|English:" not in raw:
            return []
        english = raw.split("|English:", 1)[1]
        glosses = []
        for part in english.split(","):
            gloss = part.strip().lower()
            if gloss and gloss != "misc":
                glosses.append(gloss)
        return glosses

    def parse_interlinear_token(self, token):
        parts = token.split("|")
        while len(parts) < 10:
            parts.append("")
        strongs = parts[3].upper() if parts[3] else ""
        try:
            source_order = int(parts[0])
        except ValueError:
            source_order = None
        surface = parts[6]
        if surface in {"→", "←"}:
            surface = ""
        return {
            "source_order": source_order,
            "strongs": strongs,
            "morph": parts[4],
            "english": parts[5],
            "surface": surface,
            "translit": parts[7],
            "lemma": parts[8],
            "lemma_translit": parts[9],
            "gloss": parts[10] if len(parts) > 10 else "",
        }

    def detect_source_language(self, parsed_tokens):
        codes = [token["strongs"] for token in parsed_tokens if token["strongs"]]
        if any(code.startswith("G") for code in codes):
            return "Greek"
        if any(code.startswith("H") for code in codes):
            return "Hebrew / Aramaic"
        if any(re.search(r"[\u0590-\u05ff]", token["surface"]) for token in parsed_tokens):
            return "Hebrew / Aramaic"
        return "Source Text"

    def display_source_text(self, parsed_tokens):
        source_tokens = sorted(
            [token for token in parsed_tokens if token["surface"]],
            key=lambda token: token["source_order"] if token["source_order"] is not None else 9999,
        )
        source_words = [token["surface"] for token in source_tokens]
        if not source_words:
            return
        translit_words = [token["translit"] for token in source_tokens if token["translit"]]
        body = Text()
        body.append(" ".join(source_words), style="source.text")
        if translit_words:
            body.append("\n\n", style="ui.meta")
            body.append(" ".join(translit_words), style="source.translit")
        console.print(
            Panel(
                body,
                title=f"🔡 {self.detect_source_language(parsed_tokens)}",
                border_style="source.border",
                padding=(1, 2),
            )
        )

    def get_lxx_book_code(self, book):
        canon_book = self.reverse_canon_map.get(book, book)
        if canon_book == "Psalm":
            canon_book = "Psalms"
        return LXX_BOOK_CODES.get(canon_book)

    def get_lxx_tokens(self, db_ref):
        if not self.lxx_db:
            return []
        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return []
        book_code = self.get_lxx_book_code(parts["book"])
        if not book_code:
            return []
        try:
            rows = self.lxx_db.query(
                """
                SELECT
                    t.word_num,
                    t.text,
                    t.lemma,
                    t.strong,
                    t.morph,
                    t.pos,
                    t.gloss,
                    t.head,
                    t.dependency,
                    m.strong AS candidate_strong,
                    m.confidence AS candidate_confidence,
                    m.occurrence_count AS candidate_count,
                    m.lemma_token_count AS candidate_total,
                    m.lexicon_source AS candidate_source
                FROM lxx_text t
                LEFT JOIN lxx_lemma_strong_map m
                  ON m.book = t.book
                 AND m.lemma = t.lemma
                 AND m.confidence = (
                    SELECT MAX(m2.confidence)
                    FROM lxx_lemma_strong_map m2
                    WHERE m2.book = t.book
                      AND m2.lemma = t.lemma
                 )
                WHERE t.book = ? AND t.chapter = ? AND t.verse = ?
                ORDER BY t.word_num
                """,
                (book_code, parts["chapter"], parts["verse"]),
            )
        except sqlite3.Error:
            return []
        tokens = []
        for (
            word_num, text, lemma, strong, morph, pos, gloss, head, dependency,
            candidate_strong, candidate_confidence, candidate_count,
            candidate_total, candidate_source,
        ) in rows:
            tokens.append({
                "source_order": word_num,
                "surface": text or "",
                "lemma": lemma or "",
                "strongs": strong or "",
                "candidate_strongs": candidate_strong or "",
                "candidate_confidence": candidate_confidence or 0,
                "candidate_count": candidate_count or 0,
                "candidate_total": candidate_total or 0,
                "candidate_source": candidate_source or "",
                "morph": morph or "",
                "pos": pos or "",
                "gloss": gloss or "",
                "head": head,
                "dependency": dependency or "",
            })
        return tokens

    def describe_lxx_coverage(self, tokens):
        if not tokens:
            return "no parsed LXX rows"
        total = len(tokens)
        strongs_count = sum(1 for token in tokens if token["strongs"])
        candidate_count = sum(1 for token in tokens if not token["strongs"] and token.get("candidate_strongs"))
        morph_count = sum(1 for token in tokens if token["morph"])
        labels = []
        if strongs_count == 0:
            labels.append("Strong's unavailable")
        elif strongs_count < total:
            labels.append(f"Strong's partial {strongs_count}/{total}")
        if candidate_count:
            labels.append(f"candidate Strong's {candidate_count}/{total}")
        if morph_count == 0:
            labels.append("morphology unavailable")
        elif morph_count < total:
            labels.append(f"morphology partial {morph_count}/{total}")
        return "; ".join(labels) if labels else "parsed tokens complete"

    def display_lxx_study(self, db_ref, animate=None):
        tokens = self.get_lxx_tokens(db_ref)
        if not tokens:
            return False

        self.pause_study_section(animate)
        source_words = [token["surface"] for token in tokens if token["surface"]]
        if source_words:
            console.print(
                Panel(
                    Text(" ".join(source_words), style="source.text"),
                    title="LXX Greek",
                    subtitle=f"parsed Septuagint layer | {self.describe_lxx_coverage(tokens)}",
                    border_style="source.border",
                    padding=(1, 2),
                )
            )

        table = Table(
            title=f"LXX Interlinear: {self.format_display_ref(db_ref)}",
            box=None,
            expand=True,
            pad_edge=False,
        )
        table.add_column("#", style="dim", justify="right", no_wrap=True, width=3)
        table.add_column("Greek", style="source.text", overflow="fold", ratio=3)
        table.add_column("Lemma", style="lexicon.word", overflow="fold", ratio=3)
        table.add_column("Gloss", style="text", overflow="fold", ratio=3)
        table.add_column("Morph", style="interlinear.strongs", overflow="fold", ratio=2)
        table.add_column("Syntax", style="dim", overflow="fold", ratio=2)
        for token in tokens[:30]:
            strongs = token["strongs"]
            is_candidate = False
            if not strongs and token.get("candidate_strongs"):
                strongs = f"{token['candidate_strongs']}?"
                is_candidate = True

            morph_info = token["morph"] or (f"POS:{token['pos']}" if token["pos"] else "")
            display_morph = strongs or morph_info or "-"
            if strongs and morph_info:
                display_morph = f"{strongs} ({morph_info})"

            syntax = token["dependency"] or "-"
            if token["head"] is not None:
                syntax = f"{syntax} -> {token['head']}"

            table.add_row(
                str(token["source_order"]),
                token["surface"] or "-",
                token["lemma"] or "-",
                token["gloss"] or "-",
                display_morph,
                syntax,
            )
        console.print(table)

        strong_tokens = [token for token in tokens if token["strongs"]]
        if strong_tokens:
            narrow_notes = console.width < 88
            lex_table = Table(
                title="LXX Greek Lexicon Notes",
                box=None,
                expand=True,
                pad_edge=False,
            )
            lex_table.add_column("Strongs", style="lexicon.num", no_wrap=True, width=7)
            if narrow_notes:
                lex_table.add_column("Notes", overflow="fold", ratio=1)
            else:
                lex_table.add_column("Lemma", style="lexicon.word", overflow="fold", ratio=2)
                lex_table.add_column("Details", overflow="fold", ratio=8)
            seen = set()
            for token in strong_tokens:
                strongs = token["strongs"]
                if strongs in seen:
                    continue
                seen.add(strongs)
                entry = self.lookup_lexicon_entry(strongs)
                lemma = token["lemma"] or (entry["db"][1] if entry["db"] else "")
                pieces = []
                if token["morph"]:
                    pieces.append(token["morph"])
                if entry["step"]:
                    step_def = entry["step"].get("definition", "")
                    step_def = re.sub(r'<br\s*/?>', ' ', step_def, flags=re.IGNORECASE)
                    step_def = re.sub(r'<[^>]+>', '', step_def)
                    pieces.append(study_note_excerpt(step_def))
                elif entry["interlinear"]:
                    pieces.append(study_note_excerpt(entry["interlinear"].get("d", "")))
                elif entry["db"]:
                    pieces.append(study_note_excerpt(entry["db"][3]))
                details = " | ".join(piece for piece in pieces if piece) or "-"
                if narrow_notes:
                    lex_table.add_row(strongs, study_lexicon_text(lemma or "-", details))
                else:
                    lex_table.add_row(strongs, lemma or "-", details)
                if len(seen) >= 12:
                    break
            console.print(lex_table)
        candidate_tokens = [token for token in tokens if not token["strongs"] and token.get("candidate_strongs")]
        if candidate_tokens:
            lex_table = Table(
                title="LXX Candidate Strong's Notes",
                box=None,
                expand=True,
                pad_edge=False,
            )
            lex_table.add_column("Candidate", style="lexicon.num", no_wrap=True, width=9)
            lex_table.add_column("Lemma", style="lexicon.word", overflow="fold", ratio=2)
            lex_table.add_column("Evidence", overflow="fold", ratio=8)
            seen = set()
            for token in candidate_tokens:
                strongs = token["candidate_strongs"]
                key = (token["lemma"], strongs)
                if key in seen:
                    continue
                seen.add(key)
                entry = self.lookup_lexicon_entry(strongs)
                details = []
                if entry["step"]:
                    step_def = entry["step"].get("definition", "")
                    step_def = re.sub(r'<br\s*/?>', ' ', step_def, flags=re.IGNORECASE)
                    step_def = re.sub(r'<[^>]+>', '', step_def)
                    details.append(step_def[:300])
                elif entry["interlinear"]:
                    details.append(entry["interlinear"].get("d", "")[:300])
                elif entry["db"]:
                    details.append(entry["db"][3][:300])
                if token["candidate_total"]:
                    details.append(
                        f"{token['candidate_count']}/{token['candidate_total']} lemma-token bridge matches"
                    )
                lex_table.add_row(
                    f"{strongs}?",
                    token["lemma"] or "-",
                    " | ".join(detail for detail in details if detail) or "-",
                )
                if len(seen) >= 12:
                    break
            console.print(lex_table)
        return True

    def display_study_tsk(self, db_ref, parsed_tokens):
        refs = self.get_tsk_crossrefs(db_ref)
        if not refs:
            return
        anchor_words = []
        seen_words = set()
        for token in parsed_tokens:
            word = token["english"] or token["gloss"] or token["lemma"] or token["surface"]
            word = word.strip(" ,.;:!?").lower()
            if len(word) < 3 or word in seen_words:
                continue
            seen_words.add(word)
            anchor_words.append(word)
            if len(anchor_words) >= 10:
                break
        table = Table(
            title="🔗 Treasury of Scripture Knowledge",
            box=None,
            expand=True,
            pad_edge=False,
        )
        table.add_column("Ref", style="verse.ref", no_wrap=True, ratio=2, min_width=10)
        table.add_column("Votes", style="dim", justify="right", no_wrap=True, width=5)
        table.add_column("Preview", overflow="fold", ratio=7)
        for to_ref, votes in refs:
            preview = self.get_crossref_preview(to_ref)
            table.add_row(self.format_tsk_display_ref(to_ref), str(votes), preview[:140] if preview else "")
        console.print(table)
        if anchor_words:
            console.print("[dim]Verse-level TSK links; local data has no per-word anchor. Key terms: {}[/]".format(", ".join(anchor_words)))

    def pause_study_section(self, animate):
        if self.should_animate(animate):
            time.sleep(0.16)

    def study_export_dir(self):
        path = os.path.expanduser("~/Documents/lex_exports/studies")
        os.makedirs(path, exist_ok=True)
        return path

    def study_export_filename(self, db_ref, ext):
        parts = self.parse_reference_parts(db_ref)
        label = self.format_display_ref(db_ref) if parts else db_ref
        safe = re.sub(r"[^A-Za-z0-9._-]+", "_", f"lex_study_{label}.{ext}").strip("_")
        return os.path.join(self.study_export_dir(), safe)

    def build_study_export_data(self, db_ref):
        row = self.get_interlinear_index().get(db_ref)
        if not row or not row.get("p"):
            return None
        parsed_tokens = [self.parse_interlinear_token(token) for token in row["p"]]
        parts = self.parse_reference_parts(db_ref)
        display_ref = self.format_display_ref(db_ref) if parts else db_ref
        verse_row = self.bible_db.query(
            "SELECT text FROM bible WHERE reference = ? ORDER BY id LIMIT 1",
            (db_ref,)
        )
        source_tokens = sorted(
            [token for token in parsed_tokens if token["surface"]],
            key=lambda token: token["source_order"] if token["source_order"] is not None else 9999,
        )
        lex_notes = []
        seen = set()
        for parsed in parsed_tokens:
            strongs = parsed["strongs"]
            if not strongs or strongs in seen:
                continue
            seen.add(strongs)
            entry = self.lookup_lexicon_entry(strongs)
            lemma = parsed["lemma"] or (entry["db"][1] if entry["db"] else "")
            pieces = []
            if parsed["morph"]:
                pieces.append(parsed["morph"])
            if entry["step"]:
                step_def = entry["step"].get("definition", "")
                step_def = re.sub(r'<br\s*/?>', ' ', step_def, flags=re.IGNORECASE)
                step_def = re.sub(r'<[^>]+>', '', step_def)
                pieces.append(step_def[:1000])
            elif entry["interlinear"]:
                pieces.append(entry["interlinear"].get("d", "")[:1000])
            elif entry["db"]:
                pieces.append(entry["db"][3][:1000])
            if entry["step"] and entry["step"].get("translit"):
                lemma = f"{lemma} ({entry['step']['translit']})"
            elif entry["db"]:
                lemma = f"{lemma} ({entry['db'][2]})"
            lex_notes.append({"strongs": strongs, "lemma": lemma or "-", "details": " | ".join(piece for piece in pieces if piece) or "-"})
            if len(lex_notes) >= 18:
                break
        tsk_refs = []
        for to_ref, votes in self.get_tsk_crossrefs(db_ref)[:24]:
            preview = self.get_crossref_preview(to_ref)
            tsk_refs.append({"reference": to_ref, "votes": votes, "preview": preview or ""})
        return {
            "db_ref": db_ref,
            "display_ref": display_ref,
            "verse": self.clean_text(verse_row[0][0]) if verse_row else "",
            "language": self.detect_source_language(parsed_tokens),
            "source": " ".join(token["surface"] for token in source_tokens),
            "transliteration": " ".join(token["translit"] for token in source_tokens if token["translit"]),
            "interlinear": parsed_tokens[:30],
            "lex_notes": lex_notes,
            "topical_refs": self.get_reverse_naves(db_ref),
            "tsk_refs": tsk_refs,
        }

    def export_study_docx(self, db_ref):
        try:
            from docx import Document
        except ImportError:
            console.print("[warning]DOCX export needs python-docx installed.[/]")
            return None
        data = self.build_study_export_data(db_ref)
        if not data:
            return None
        path = self.study_export_filename(db_ref, "docx")
        doc = Document()
        doc.add_heading(f"Lex Study: {data['display_ref']}", level=1)
        if data["verse"]:
            doc.add_paragraph(data["verse"])
        doc.add_heading(data["language"], level=2)
        if data["source"]:
            doc.add_paragraph(data["source"])
        if data["transliteration"]:
            doc.add_paragraph(data["transliteration"])
        doc.add_heading("Interlinear", level=2)
        table = doc.add_table(rows=1, cols=5)
        for cell, title in zip(table.rows[0].cells, ["English", "Source", "Lemma", "Code", "Gloss"]):
            cell.text = title
        for parsed in data["interlinear"]:
            row = table.add_row().cells
            row[0].text = parsed["english"] or "-"
            row[1].text = f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "-"
            row[2].text = f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "-"
            row[3].text = parsed["strongs"] or parsed["morph"] or "-"
            row[4].text = parsed["gloss"] or parsed["english"] or "-"
        doc.add_heading("Lexicon Notes", level=2)
        for note in data["lex_notes"]:
            doc.add_paragraph(f"{note['strongs']} - {note['lemma']}: {note['details']}")

        if data["topical_refs"]:
            doc.add_heading("Topical Associations (Nave's)", level=2)
            doc.add_paragraph(" • ".join(data["topical_refs"]))

        doc.add_heading("Treasury of Scripture Knowledge", level=2)
        for ref in data["tsk_refs"]:
            doc.add_paragraph(f"{ref['reference']} ({ref['votes']}): {ref['preview']}")
        doc.save(path)
        return path

    def export_study_pdf(self, db_ref):
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table as PdfTable
        except ImportError:
            console.print("[warning]PDF export needs reportlab installed.[/]")
            return None
        data = self.build_study_export_data(db_ref)
        if not data:
            return None
        path = self.study_export_filename(db_ref, "pdf")
        doc = SimpleDocTemplate(path, pagesize=letter, rightMargin=42, leftMargin=42, topMargin=42, bottomMargin=42)
        styles = getSampleStyleSheet()
        self.setup_pdf_styles(styles)
        story = [self.pdf_paragraph(f"Lex Study: {data['display_ref']}", styles["Title"])]
        if data["verse"]:
            story.extend([self.pdf_paragraph(data["verse"], styles["BodyText"]), Spacer(1, 10)])
        story.append(self.pdf_paragraph(data["language"], styles["Heading2"]))
        if data["source"]:
            source_style = styles["Hebrew"] if re.search(r"[\u0590-\u05ff]", data["source"]) and "Hebrew" in styles else styles["BodyText"]
            story.append(self.pdf_paragraph(data["source"], source_style))
        if data["transliteration"]:
            story.append(self.pdf_paragraph(data["transliteration"], styles["Italic"]))
        story.extend([Spacer(1, 10), self.pdf_paragraph("Interlinear", styles["Heading2"])])
        table_rows = [["English", "Source", "Lemma", "Code", "Gloss"]]
        for parsed in data["interlinear"]:
            table_rows.append([
                self.pdf_paragraph(parsed["english"] or "-", styles["BodyText"]),
                self.pdf_paragraph(f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "-", styles["Hebrew"] if parsed["surface"] and re.search(r"[\u0590-\u05ff]", parsed["surface"]) and "Hebrew" in styles else styles["BodyText"]),
                self.pdf_paragraph(f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "-", styles["Hebrew"] if parsed["lemma"] and re.search(r"[\u0590-\u05ff]", parsed["lemma"]) and "Hebrew" in styles else styles["BodyText"]),
                self.pdf_paragraph(parsed["strongs"] or parsed["morph"] or "-", styles["BodyText"]),
                self.pdf_paragraph(parsed["gloss"] or parsed["english"] or "-", styles["BodyText"]),
            ])
        story.append(PdfTable(table_rows, repeatRows=1))
        story.append(self.pdf_paragraph("Lexicon Notes", styles["Heading2"]))
        for note in data["lex_notes"]:
            story.append(self.pdf_paragraph(f"{note['strongs']} - {note['lemma']}: {note['details']}", styles["BodyText"]))

        if data["topical_refs"]:
            story.append(self.pdf_paragraph("Topical Associations (Nave's)", styles["Heading2"]))
            story.append(self.pdf_paragraph(" • ".join(data["topical_refs"]), styles["BodyText"]))

        story.append(self.pdf_paragraph("Treasury of Scripture Knowledge", styles["Heading2"]))
        for ref in data["tsk_refs"]:
            story.append(self.pdf_paragraph(f"{ref['reference']} ({ref['votes']}): {ref['preview']}", styles["BodyText"]))
        doc.build(story)
        return path

    def prompt_study_export(self, db_ref):
        while True:
            self.render_action_bar(
                "Export",
                [
                    ("d", "DOCX study packet"),
                    ("f", "PDF study packet"),
                    ("p", "PPTX verse slide"),
                    ("o", "open studies folder"),
                    ("q", "back"),
                ],
            )
            action = Prompt.ask("Export action", choices=["d", "f", "p", "o", "q"], default="q").lower()
            if action == "q":
                return
            if action == "o":
                self.open_exports_folder(self.study_export_dir())
                continue

            if action == "d":
                path = self.export_study_docx(db_ref)
            elif action == "f":
                path = self.export_study_pdf(db_ref)
            else:
                path = self.export_study_pptx(db_ref)
            if path:
                self.open_export(path)
                return

    def prompt_bulk_export(self, refs_str, mode):
        # Define supported formats per mode
        mode_formats = {
            "read": {"i": ("png", "PNG"), "p": ("pptx", "PPTX")},
            "study": {"f": ("pdf", "PDF"), "d": ("docx", "DOCX")},
            "web": {"f": ("pdf", "PDF"), "d": ("docx", "DOCX")},
        }
        supported = mode_formats.get(mode, {})

        # Build menu content
        menu_items = [Text("Select export format:"), Text("")]
        for key, (_, label) in supported.items():
            menu_items.append(Text(f"{key} - {label}"))
        menu_items.extend([Text(""), Text("q - Quit")])

        console.print(Panel(
            Group(*menu_items),
            title=f"Export Format ({mode.capitalize()})",
            border_style="cyan",
            expand=False
        ))

        choices = list(supported.keys()) + ["q"]
        action = Prompt.ask("Format", choices=choices, default="q").lower()

        if action == "q":
            return

        fmt = supported[action][0]

        ratio = "16:9"
        if fmt == "png":
            console.print(Panel(
                Group(
                    Text("Select PNG aspect ratio:"),
                    Text(""),
                    Text("1 - 16:9 (Letterboxed)"),
                    Text("2 - 1:1 (Square)"),
                ),
                title="PNG Ratio",
                border_style="cyan",
                expand=False
            ))
            ratio_action = Prompt.ask("Ratio", choices=["1", "2"], default="1")
            ratio = "16:9" if ratio_action == "1" else "1:1"

        theme = "auto"
        if fmt in ["png", "pptx"]:
            console.print(Panel(
                Group(
                    Text("Select export theme:"),
                    Text(""),
                    Text("1 - Light"),
                    Text("2 - Dark"),
                ),
                title="Export Theme",
                border_style="cyan",
                expand=False
            ))
            theme_action = Prompt.ask("Theme", choices=["1", "2"], default="1")
            theme = "light" if theme_action == "1" else "dark"

        self.handle_bulk_export(refs_str, mode=mode, format=fmt, ratio=ratio, theme=theme)

    def prompt_study_actions(self, db_ref):
        current_ref = db_ref
        while True:
            self.render_action_bar(
                "Study Actions",
                [
                    ("n / p", "next or previous verse"),
                    ("r", "read context"),
                    ("w", "verse web"),
                    ("e", "export"),
                    ("q", "done"),
                ],
            )
            action = Prompt.ask("Study action", choices=["n", "p", "r", "w", "e", "q"], default="q").lower()
            if action == "q":
                return
            if action == "e":
                self.prompt_study_export(current_ref)
                continue
            if action == "r":
                self.display_verse(self.format_display_ref(current_ref))
                continue
            if action == "w":
                self.display_verse_web(self.format_display_ref(current_ref))
                continue
            next_ref = self.get_navigation_reference(current_ref, "next" if action == "n" else "prev")
            if next_ref:
                current_ref = next_ref
                self.display_study(current_ref, actions=False)

    def display_study_query(self, query, animate=None):
        ref_norm, book, chap, verse, v_end = self.normalize_ref(query)
        if not ref_norm:
            return False

        # Ensure we use the DB-specific book name in the LIKE query
        db_book = self.canon_map.get(re.sub(r"[^a-z0-9]+", "", book.lower()), book)

        # Default to verse 1 if not specified
        v_num = verse or "1"

        res = self.bible_db.query(
            "SELECT MIN(id), reference, text FROM bible WHERE reference LIKE ? GROUP BY reference LIMIT 1",
            (f"%:{db_book}:{chap}:{v_num}",)
        )
        if res:
            target_id, ref, text = res[0]
            self.display_study(ref, animate=animate)
            self.save_history(ref)
            return True
        return False

    def display_study_context(self, db_ref):
        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return

        prev2 = None
        prev1 = self.get_navigation_reference(db_ref, "prev")
        if prev1:
            prev2 = self.get_navigation_reference(prev1, "prev")
        next1 = self.get_navigation_reference(db_ref, "next")
        next2 = self.get_navigation_reference(next1, "next") if next1 else None

        rows = []
        for candidate in [prev2, prev1, db_ref, next1, next2]:
            if not candidate:
                continue
            candidate_parts = self.parse_reference_parts(candidate)
            if candidate_parts and candidate_parts["book"] != parts["book"]:
                continue
            row = self.bible_db.query(
                "SELECT MIN(id), reference, text FROM bible WHERE reference = ? GROUP BY reference",
                (candidate,)
            )
            if row:
                rows.append(row[0])

        if rows:
            self.render_verse_context(
                rows,
                db_ref,
                parts["book"],
                parts["chapter"],
                parts["verse"],
            )

    def find_source_variant_label(self, book, chapter, verse):
        for start_book, start_ch, start_v, end_book, end_ch, end_v, label in SOURCE_VARIANT_RANGES:
            if book not in {start_book, end_book}:
                continue
            if start_book == end_book:
                if book == start_book and (start_ch, start_v) <= (chapter, verse) <= (end_ch, end_v):
                    return label
                continue
            if book == start_book and chapter == start_ch and verse >= start_v:
                return label
            if book == end_book and chapter == end_ch and verse <= end_v:
                return label
        return None

    def resolve_study_source_route(self, db_ref):
        parts = self.parse_reference_parts(db_ref)
        if not parts:
            return None
        if parts["version"] == "esv":
            return {
                "source_ref": db_ref,
                "target_base": "ESV interlinear",
                "confidence": "direct",
                "status": "direct",
                "label": None,
                "warning": None,
            }
        if parts["version"] != "gen":
            return None

        canon_book = self.reverse_canon_map.get(parts["book"], parts["book"])
        book_for_group = "Psalms" if canon_book == "Psalm" else canon_book
        if book_for_group in PROTESTANT_OT_BOOKS:
            target_base = "Masoretic Hebrew / Aramaic tradition"
        elif book_for_group in PROTESTANT_NT_BOOKS:
            target_base = "Textus Receptus / Reformation Greek tradition"
        else:
            return {
                "source_ref": None,
                "target_base": "unmapped",
                "confidence": "unmapped",
                "status": "unmapped",
                "label": "Geneva source layer: unmapped",
                "warning": "This Geneva book is readable here, but Lex does not yet have a selected source-language study layer for it.",
            }

        source_ref = f"esv:{canon_book}:{parts['chapter']}:{parts['verse']}"
        variant_label = self.find_source_variant_label(canon_book, parts["chapter"], parts["verse"])
        if variant_label:
            return {
                "source_ref": source_ref,
                "target_base": target_base,
                "confidence": "variant/proxy",
                "status": "variant",
                "label": f"Geneva source target: {target_base}",
                "warning": (
                    f"{variant_label}: Geneva/TR tradition includes this unit, while many modern critical editions "
                    "omit or bracket it. Any local interlinear row shown below is a labeled proxy, not Geneva's textual authority."
                ),
            }
        return {
            "source_ref": source_ref,
            "target_base": target_base,
            "confidence": "proxy",
            "status": "proxy",
            "label": f"Geneva source target: {target_base}",
            "warning": "Local interlinear data is used as a proxy for now; it is not being treated as Geneva's textual authority.",
        }

    def display_source_route_notice(self, route):
        if not route or not route.get("label"):
            return
        body = Text()
        body.append(f"{route['label']}\n", style="text.strong")
        body.append(f"Confidence: {route['confidence']}", style="ui.meta")
        warning = route.get("warning")
        if warning:
            body.append("\n\n", style="ui.meta")
            body.append(warning, style="warning")
        console.print(Panel(body, title="Source Path", border_style="warning", padding=(1, 2)))

    def display_study(self, db_ref, animate=None, actions=None, context=True):
        parts = self.parse_reference_parts(db_ref)
        if parts and parts["version"] == "lxx":
            if self.display_lxx_study(db_ref, animate=animate):
                return True
            console.print(Panel("No local LXX study data found for this verse.", border_style="warning"))
            return False

        if parts and parts["version"] == "vulg":
            if self.latin_db and self.display_vulgate_study(db_ref, animate=animate):
                return True
            console.print(Panel("Vulgate text loaded. Latin semantic tagging available from latin.db.", border_style="info"))
            # Fall through to basic read view with Latin support

        route = self.resolve_study_source_route(db_ref)
        if parts and parts["version"] not in ("esv", "vulg") and not route:
            console.print(Panel(
                f"Interlinear study data is only available for ESV and Vulgate right now.\n\n"
                f"Read view is using {parts['version'].upper()}; run `lex study {parts['book']} {parts['chapter']}:{parts['verse']}` without `-v {parts['version']}` for the ESV interlinear packet.",
                border_style="warning"
            ))
            return False
        study_ref = route["source_ref"] if route and route.get("source_ref") else db_ref
        if route and not route.get("source_ref"):
            if context:
                self.display_study_context(db_ref)
            self.display_source_route_notice(route)
            return False

        index = self.get_interlinear_index()
        row = index.get(study_ref)

        # If no exact match (likely due to version prefix mismatch), try to find by book:chap:verse
        if not row or not row.get("p"):
            study_parts = self.parse_reference_parts(study_ref)
            if study_parts:
                # Use reverse canon map to get the canonical book name (e.g. "Ps" -> "Psalm")
                canon_book = self.reverse_canon_map.get(study_parts["book"], study_parts["book"])
                generic_suffix = f":{canon_book}:{study_parts['chapter']}:{study_parts['verse']}"

                # The index keys in interlinear JSON look like "esv:Genesis:1:1" or "esv:Psalm:1:1"
                for k, v in index.items():
                    if k.endswith(generic_suffix):
                        row = v
                        break

        if not row or not row.get("p"):
            console.print(Panel("No local interlinear data found for this verse.", border_style="warning"))
            return False
        parsed_tokens = [self.parse_interlinear_token(token) for token in row["p"]]
        if context:
            self.display_study_context(db_ref)
        self.display_source_route_notice(route)
        self.pause_study_section(animate)
        self.display_source_text(parsed_tokens)
        self.pause_study_section(animate)
        table_title = f"🔤 Study: {db_ref}"
        if study_ref != db_ref:
            table_title = f"🔤 Study: {db_ref} via {study_ref}"
        if console.width < 88:
            verse_table = Table(
                title=table_title,
                box=None,
                expand=True,
                pad_edge=False,
            )
            verse_table.add_column("#", style="dim", justify="right", no_wrap=True, width=3)
            verse_table.add_column("Token", style="text.strong", overflow="fold", ratio=2, min_width=10)
            verse_table.add_column("Analysis", overflow="fold", ratio=6, min_width=28)
            for idx, parsed in enumerate(parsed_tokens[:30], start=1):
                verse_table.add_row(
                    str(idx),
                    parsed["english"] or "•",
                    study_analysis_text(parsed),
                )
        else:
            verse_table = Table(
                title=table_title,
                box=None,
                expand=True,
                pad_edge=False,
            )
            verse_table.add_column("Eng", style="text.strong", overflow="fold", ratio=2)
            verse_table.add_column("Src", style="source.text", overflow="fold", ratio=3)
            verse_table.add_column("Lemma", style="lexicon.word", overflow="fold", ratio=3)
            verse_table.add_column("Code", style="interlinear.strongs", no_wrap=True, width=6)
            verse_table.add_column("Gloss", style="text", overflow="fold", ratio=2)
            for parsed in parsed_tokens[:30]:
                code = parsed["strongs"] or parsed["morph"] or "-"
                gloss = parsed["gloss"] or parsed["english"] or "-"
                verse_table.add_row(
                    parsed["english"] or "•",
                    f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "•",
                    f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "•",
                    code,
                    gloss,
                )
        console.print(verse_table)

        self.pause_study_section(animate)
        narrow_notes = console.width < 88
        lex_table = Table(
            title="📚 Lexicon Notes",
            box=None,
            expand=True,
            pad_edge=False,
        )
        lex_table.add_column("Strongs", style="lexicon.num", no_wrap=True, width=7)
        if narrow_notes:
            lex_table.add_column("Notes", overflow="fold", ratio=1)
        else:
            lex_table.add_column("Lemma", style="lexicon.word", overflow="fold", ratio=2)
            lex_table.add_column("Details", overflow="fold", ratio=8)
        seen = set()
        for parsed in parsed_tokens:
            strongs = parsed["strongs"]
            if not strongs or strongs in seen:
                continue
            seen.add(strongs)
            entry = self.lookup_lexicon_entry(strongs)
            lemma = parsed["lemma"] or (entry["db"][1] if entry["db"] else "")
            pieces = []
            if parsed["morph"]:
                pieces.append(parsed["morph"])
            if entry["step"]:
                step_def = entry["step"].get("definition", "")
                step_def = re.sub(r'<br\s*/?>', ' ', step_def, flags=re.IGNORECASE)
                step_def = re.sub(r'<[^>]+>', '', step_def)
                pieces.append(study_note_excerpt(step_def))
            elif entry["interlinear"]:
                pieces.append(study_note_excerpt(entry["interlinear"].get("d", "")))
            elif entry["db"]:
                pieces.append(study_note_excerpt(entry["db"][3]))
            if entry["step"] and entry["step"].get("translit"):
                lemma = f"{lemma} ({entry['step']['translit']})"
            elif entry["db"]:
                lemma = f"{lemma} ({entry['db'][2]})"
            details = " | ".join(piece for piece in pieces if piece) or "-"
            if narrow_notes:
                lex_table.add_row(strongs, study_lexicon_text(lemma or "-", details))
            else:
                lex_table.add_row(strongs, lemma or "-", details)
            if len(seen) >= 12:
                break
        console.print(lex_table)

        topical_refs = self.get_reverse_naves(db_ref)
        if topical_refs:
            self.pause_study_section(animate)
            topics_text = Text()
            for idx, topic in enumerate(topical_refs):
                if idx > 0:
                    topics_text.append("  •  ", style="dim")
                topics_text.append(topic, style="dict.topic")
            console.print(Panel(topics_text, title="🏷️ Topical Associations (Nave's)", border_style="ui.border"))

        self.pause_study_section(animate)
        self.display_study_tsk(study_ref, parsed_tokens)
        use_actions = console.is_terminal if actions is None else actions
        if use_actions:
            self.prompt_study_actions(db_ref)
        return True

    def display_dictionary_howto(self):
        md = "# 📖 Bible Dictionary Help\n\n- `lex define \"Grace\"`"
        console.print(Panel(Markdown(md), border_style="violet"))

    # -----------------------------------------------------------------------
    # Creeds and historical documents
    # -----------------------------------------------------------------------
    def format_creed_source(self, topic, source):
        source_map = {
            "Athanasian Creed": "5th c. | trinitarian creed",
            "Augsburg Confession": "1530 | Lutheran confession",
            "Baltimore Catechism": "1885 | Roman Catholic catechism",
            "Belgic Confession": "1561 | Reformed confession",
            "Canons of Dort": "1619 | Reformed canons",
            "Chalcedonian Definition": "451 | Christological definition",
            "Confession of Dositheus": "1672 | Eastern Orthodox confession",
            "Council of Trent": "1545-1563 | Catholic council decrees",
            "Heidelberg Catechism": "1563 | Reformed catechism",
            "London Baptist Confession of Faith": "1689 | Baptist confession",
            "The Apostles' Creed": "early | baptismal creed",
            "The Longer Catechism of the Orthodox Church": "1830s | Orthodox catechism",
            "The Nicene Creed": "325/381 | ecumenical creed",
            "Thirty-Nine Articles": "1571 | Anglican articles",
            "Westminster Confession of Faith": "1646 | Presbyterian confession",
            "Westminster Larger Catechism": "1648 | Presbyterian catechism",
            "Westminster Shorter Catechism": "1647 | Presbyterian catechism",
        }
        return source_map.get(topic, source or "undated | creed text")

    def creed_sort_key(self, topic):
        order = {
            "The Apostles' Creed": 200,
            "The Nicene Creed": 325,
            "Chalcedonian Definition": 451,
            "Athanasian Creed": 500,
            "Council of Trent": 1545,
            "Augsburg Confession": 1530,
            "Belgic Confession": 1561,
            "Heidelberg Catechism": 1563,
            "Thirty-Nine Articles": 1571,
            "Canons of Dort": 1619,
            "Westminster Confession of Faith": 1646,
            "Westminster Shorter Catechism": 1647,
            "Westminster Larger Catechism": 1648,
            "London Baptist Confession of Faith": 1689,
            "Confession of Dositheus": 1672,
            "Baltimore Catechism": 1885,
            "The Longer Catechism of the Orthodox Church": 1830,
        }
        return order.get(topic, 9999), topic

    def creed_tradition(self, topic):
        groups = {
            "The Apostles' Creed": "Ecumenical Creeds",
            "The Nicene Creed": "Ecumenical Creeds",
            "Chalcedonian Definition": "Ecumenical Creeds",
            "Athanasian Creed": "Ecumenical Creeds",
            "Augsburg Confession": "Lutheran",
            "Belgic Confession": "Reformed",
            "Canons of Dort": "Reformed",
            "Heidelberg Catechism": "Reformed",
            "Westminster Confession of Faith": "Reformed",
            "Westminster Shorter Catechism": "Reformed",
            "Westminster Larger Catechism": "Reformed",
            "London Baptist Confession of Faith": "Baptist",
            "Thirty-Nine Articles": "Anglican",
            "Council of Trent": "Roman Catholic",
            "Baltimore Catechism": "Roman Catholic",
            "Confession of Dositheus": "Eastern Orthodox",
            "The Longer Catechism of the Orthodox Church": "Eastern Orthodox",
        }
        return groups.get(topic, "Other")

    def creed_tradition_sort_key(self, topic):
        tradition_order = {
            "Ecumenical Creeds": 0,
            "Lutheran": 1,
            "Reformed": 2,
            "Anglican": 3,
            "Baptist": 4,
            "Roman Catholic": 5,
            "Eastern Orthodox": 6,
            "Other": 99,
        }
        tradition = self.creed_tradition(topic)
        year, title = self.creed_sort_key(topic)
        return tradition_order.get(tradition, 99), year, title

    def creed_year_label(self, topic, source):
        return self.format_creed_source(topic, source).split("|", 1)[0].strip()

    def extract_creed_title(self, content):
        match = re.match(r'^\[(.*?)\]\s*', content, re.DOTALL)
        return match.group(1).strip() if match else None

    def strip_creed_title(self, content):
        return re.sub(r'^\[.*?\]\s*', '', content, count=1, flags=re.DOTALL).strip()

    def is_empty_creed_content(self, content):
        return not self.strip_creed_title(content or "")

    def load_historical_document(self, topic):
        filename = HISTORICAL_DOC_FILES.get(topic)
        if not filename:
            return None
        return self.load_json_file(os.path.join(HISTORICAL_DOCS_DIR, filename))

    def build_creed_sections_from_file(self, topic):
        data = self.load_historical_document(topic)
        if not data:
            return []
        sections = []
        source = data.get("title", topic)
        for item in data.get("sections", []):
            title = item.get("title") or item.get("chapter") or item.get("question") or topic
            if item.get("q") or item.get("a"):
                body = "\n\n".join(
                    part for part in [
                        f"**Q.** {item.get('q')}" if item.get("q") else "",
                        f"**A.** {item.get('a')}" if item.get("a") else "",
                    ]
                    if part
                )
                if item.get("question"):
                    title = f"Q{item['question']}: {item.get('q', '').strip()}"
            else:
                body = item.get("content", "")
            proofs = self.extract_scripture_refs(body)
            sections.append({"title": str(title), "source": source, "body_parts": [body] if body else [], "proofs": proofs})
        return sections

    def extract_scripture_refs(self, text):
        patterns = [
            r'\b(?:[1-3]\s+)?[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\s+\d+[:;]\d+(?:-\d+)?(?:,\s*\d+(?:-\d+)?)*',
            r'\b(?:[1-3]\s+)?[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\s+\d+\b',
        ]
        refs = []
        for pattern in patterns:
            refs.extend(re.findall(pattern, text))
        cleaned = []
        seen = set()
        for ref in refs:
            normalized = ref.replace(";", ":").strip(" .")
            if normalized not in seen:
                seen.add(normalized)
                cleaned.append(normalized)
        filtered = []
        for ref in cleaned:
            if ":" not in ref and any(full.startswith(f"{ref}:") for full in cleaned):
                continue
            filtered.append(ref)
        return filtered

    def is_proof_only_row(self, content):
        body = self.strip_creed_title(content)
        refs = self.extract_scripture_refs(body)
        if not refs:
            return False
        stripped = body
        for ref in refs:
            stripped = stripped.replace(ref.replace(":", ";"), " ")
            stripped = stripped.replace(ref, " ")
        stripped = re.sub(r'[\d\W_]+', ' ', stripped)
        return len(stripped.strip()) <= 18

    def build_creed_sections(self, topic):
        rows = self.creeds_db.query("SELECT rowid, content, source FROM creeds WHERE topic = ? ORDER BY rowid", (topic,))
        if not rows:
            return self.build_creed_sections_from_file(topic)
        sections = []
        current = None
        for _, content, source in rows:
            title = self.extract_creed_title(content) or topic
            body = self.strip_creed_title(content)
            proof_refs = self.extract_scripture_refs(body)
            proof_only = self.is_proof_only_row(content)
            if current is None or current["title"] != title:
                current = {"title": title, "source": source, "body_parts": [], "proofs": []}
                sections.append(current)
            if not proof_only and body:
                current["body_parts"].append(body)
            for ref in proof_refs:
                if ref not in current["proofs"]:
                    current["proofs"].append(ref)
        if not any(section["body_parts"] or section["proofs"] for section in sections):
            file_sections = self.build_creed_sections_from_file(topic)
            if file_sections:
                return file_sections
        return sections

    def should_render_creed_as_document(self, topic, sections):
        short_topics = {
            "The Apostles' Creed",
            "The Nicene Creed",
            "Athanasian Creed",
            "Chalcedonian Definition",
        }
        total_body = sum(len("\n\n".join(section["body_parts"])) for section in sections)
        return topic in short_topics or (len(sections) <= 4 and total_body <= 5000)

    def get_creed_original(self, topic, section_title):
        doc = CREED_ORIGINALS.get(topic)
        if not doc:
            return None, None
        return doc["sections"].get(section_title), doc["language"]

    def display_creed_note(self, topic):
        note = CREED_NOTES.get(topic)
        if note:
            console.print(Panel(Markdown(note), title="Textual / Tradition Note", border_style="yellow"))

    def display_creed_original_document(self, topic, sections):
        doc = CREED_ORIGINALS.get(topic)
        if not doc:
            return False
        console.print(Panel(f"{topic}\nSource: {self.format_creed_source(topic, sections[0]['source'])}\nOriginal: {doc['language']}", border_style="bold green"))
        table = Table(title=f"{topic}: English / {doc['language']}", box=box.ROUNDED, expand=True, show_lines=True)
        table.add_column(doc["language"], style="cyan", overflow="fold", ratio=1)
        table.add_column("English", style="text", overflow="fold", ratio=1)
        for section in sections:
            body = "\n\n".join(section["body_parts"]).strip()
            orig_body = doc["sections"].get(section["title"])
            if not body and not orig_body:
                continue

            left = f"{orig_body or '[not yet loaded]'}"
            right = f"{body}"
            table.add_row(left, right)
        console.print(table)
        self.display_creed_note(topic)
        return True

    def display_creed_document(self, topic, sections):
        if not sections:
            return False
        if topic in CREED_ORIGINALS:
            return self.display_creed_original_document(topic, sections)
        parts = []
        proof_set = []
        seen = set()
        for section in sections:
            body = "\n\n".join(section["body_parts"]).strip()
            if body:
                parts.append(f"## {section['title']}\n\n{body}")
            for ref in section["proofs"]:
                if ref not in seen:
                    seen.add(ref)
                    proof_set.append(ref)
        proofs = ""
        if proof_set:
            proofs = "\n\n---\n\n**Scripture Proofs**\n\n" + "\n".join(f"- {ref}" for ref in proof_set[:40])
        console.print(
            Panel(
                Markdown(
                    f"# {topic}\n\n"
                    f"**Source:** {self.format_creed_source(topic, sections[0]['source'])}\n\n---\n\n"
                    + "\n\n".join(parts)
                    + proofs
                ),
                border_style="bold green"
            )
        )
        self.display_creed_note(topic)
        return True

    def display_creed_sections(self, topic):
        sections = self.build_creed_sections(topic)
        if not sections:
            return False
        if self.should_render_creed_as_document(topic, sections):
            return self.display_creed_document(topic, sections)
        while True:
            table = Table(title=f"📜 {topic}", box=None)
            table.add_column("ID", style="verse.ref")
            table.add_column("Section", style="bold green")
            table.add_column("Proofs", style="dim cyan")
            for i, section in enumerate(sections, 1):
                proof_count = str(len(section["proofs"])) if section["proofs"] else "-"
                table.add_row(str(i), section["title"], proof_count)
            console.print(table)
            if not sys.stdin.isatty():
                console.print(f"[dim]Use an interactive terminal to browse sections for: lex creed {topic}[/]")
                return True
            try:
                choice = Prompt.ask("Select section, or q to quit", default="1").strip().lower()
            except EOFError:
                console.print(f"[dim]Use an interactive terminal to browse sections for: lex creed {topic}[/]")
                return True
            if choice == "q":
                return True
            if not choice.isdigit():
                console.print("[warning]Enter a section number or q.[/]")
                continue
            section_idx = int(choice) - 1
            if 0 <= section_idx < len(sections):
                self.display_creed_reader(topic, sections, start_idx=section_idx)
            else:
                console.print("[warning]Section number out of range.[/]")
        return True

    def display_creed_navigator(self, query=None):
        if query:
            matches = self.find_creed_topics(query)
            if len(matches) == 1:
                return self.display_creed_sections(matches[0][0])
            if len(matches) > 1:
                table = Table(title=f"📜 Matching Creeds: {query}", box=None)
                table.add_column("Document", style="bold green")
                table.add_column("Sections", style="dim cyan")
                table.add_column("Source", style="dim")
                for topic, source in matches:
                    table.add_row(topic, str(len(self.build_creed_sections(topic))), self.format_creed_source(topic, source))
                console.print(table)
                console.print("[dim]Use: lex creed <document name>[/]")
                return True
            res = self.creeds_db.query(
                """
                SELECT topic, content, source
                FROM creeds
                WHERE content LIKE ? AND content NOT LIKE '[]%'
                LIMIT 8
                """,
                (f'%{query}%',)
            )
            for t, c, s in res:
                title = self.extract_creed_title(c) or t
                refs = self.extract_scripture_refs(c)
                snippet = self.strip_creed_title(c)[:700]
                if refs:
                    snippet += "\n\n**Scripture Proofs:** " + "; ".join(refs[:8])
                display_source = self.format_creed_source(t, s)
                console.print(Panel(Markdown(f"# {t}: {title}\n\n**Source:** {display_source}\n\n{snippet}"), border_style="green"))
            return bool(res)

        creeds_list = self.creeds_db.query(
            """
            SELECT topic, source
            FROM creeds
            GROUP BY topic, source
            """
        )
        creeds_list = sorted(creeds_list, key=lambda row: self.creed_tradition_sort_key(row[0]))
        table = Table(title="📜 Creeds Navigator", box=None)
        table.add_column("ID", style="verse.ref")
        table.add_column("Tradition", style="bold cyan")
        table.add_column("Year", style="dim")
        table.add_column("Document", style="bold green")
        table.add_column("Sections", style="dim cyan")
        section_counts = {topic: len(self.build_creed_sections(topic)) for topic, _ in creeds_list}
        last_tradition = None
        for i, (t, s) in enumerate(creeds_list):
            tradition = self.creed_tradition(t)
            tradition_label = tradition if tradition != last_tradition else ""
            table.add_row(
                str(i+1),
                tradition_label,
                self.creed_year_label(t, s),
                t,
                str(section_counts.get(t, 0)),
            )
            last_tradition = tradition
        console.print(table)
        if not sys.stdin.isatty():
            console.print("[dim]Use: lex creed <document name>[/]")
            return True
        try:
            choice = Prompt.ask("Select ID, or q to quit", default="1").strip().lower()
        except EOFError:
            console.print("[dim]Use: lex creed <document name>[/]")
            return True
        if choice == "q":
            return True
        if not choice.isdigit():
            console.print("[warning]Enter a document ID or q.[/]")
            return True
        doc_idx = int(choice) - 1
        if 0 <= doc_idx < len(creeds_list):
            self.display_creed_sections(creeds_list[doc_idx][0])
        else:
            console.print("[warning]Document ID out of range.[/]")
        return True

    def find_creed_topics(self, query):
        normalized_query = self.normalize_term(query)
        if not normalized_query:
            return []
        rows = self.creeds_db.query(
            """
            SELECT topic, source
            FROM creeds
            GROUP BY topic, source
            """
        )
        exact = []
        partial = []
        for topic, source in rows:
            normalized_topic = self.normalize_term(topic)
            if normalized_topic == normalized_query:
                exact.append((topic, source))
            elif normalized_query in normalized_topic:
                partial.append((topic, source))
        return sorted(exact or partial, key=lambda row: self.creed_sort_key(row[0]))

    def display_creed_reader(self, topic, sections, start_idx=0):
        if not sections:
            return
        art_idx = start_idx
        while True:
            section = sections[art_idx]
            body = "\n\n".join(section["body_parts"]).strip() or "_No article body stored for this section._"
            proofs = ""
            if section["proofs"]:
                proofs = "\n\n---\n\n**Scripture Proofs**\n\n" + "\n".join(f"- {ref}" for ref in section["proofs"][:24])
            console.clear()
            original, original_language = self.get_creed_original(topic, section["title"])
            if original:
                table = Table(title=f"{topic}: {section['title']}", box=None)
                table.add_column(original_language, style="cyan", overflow="fold")
                table.add_column("English", style="text", overflow="fold")
                table.add_row(original, body)
                console.print(table)
                if proofs:
                    console.print(Panel(Markdown(proofs), border_style="green"))
            else:
                console.print(
                    Panel(
                        Markdown(
                            f"# {topic}: {section['title']}\n\n"
                            f"**Source:** {self.format_creed_source(topic, section['source'])}\n\n---\n\n{body}{proofs}"
                        ),
                        border_style="bold green"
                    )
                )
            console.print(f"[dim]Section {art_idx+1}/{len(sections)} of '{topic}'[/]")
            console.print("[dim][n] Next | [p] Prev | [m] Sections | [q] Quit[/]")

            nav = Prompt.ask("Navigate", choices=["n", "p", "m", "q"], default="q").lower()
            if nav == "n" and art_idx < len(sections)-1: art_idx += 1
            elif nav == "p" and art_idx > 0: art_idx -= 1
            elif nav == "m": break
            elif nav == "q": sys.exit(0)

    # -----------------------------------------------------------------------
    # Scripture search, Strong's lookup, dictionary, and encyclopedia
    # -----------------------------------------------------------------------
    def search_scope_clause(self, scope):
        if not scope:
            return "", ()
        clauses = " OR ".join(["reference GLOB ?"] * len(scope["books"]))
        params = tuple(f"*:{book}:*" for book in scope["books"])
        return f" AND ({clauses})", params

    def text_fts_query(self, fts_query):
        # The FTS table also indexes `reference`; column qualification prevents
        # `lex search john` from returning every verse whose reference is John.
        return f"text : ({fts_query})"

    def query_search_results(self, fts_query, limit, offset, scope=None):
        scope_clause, scope_params = self.search_scope_clause(scope)
        return self.bible_db.query(
            f"""
            SELECT reference, text
            FROM bible_fts
            WHERE bible_fts MATCH ?
            AND reference NOT GLOB '*:0'
            {scope_clause}
            ORDER BY rank
            LIMIT ? OFFSET ?
            """,
            (self.text_fts_query(fts_query), *scope_params, limit, offset)
        )

    def query_ranked_search_results(self, phrase_query, terms_query, limit, offset, scope=None):
        """Return phrase hits first, followed by remaining all-term hits."""
        scope_clause, scope_params = self.search_scope_clause(scope)
        return self.bible_db.query(
            f"""
            WITH candidates AS (
                SELECT rowid, reference, text, 0 AS tier, bm25(bible_fts) AS score
                FROM bible_fts
                WHERE bible_fts MATCH ? AND reference NOT GLOB '*:0'{scope_clause}
                UNION ALL
                SELECT rowid, reference, text, 1 AS tier, bm25(bible_fts) AS score
                FROM bible_fts
                WHERE bible_fts MATCH ? AND reference NOT GLOB '*:0'{scope_clause}
            ), ranked AS (
                SELECT rowid, reference, text, MIN(tier) AS tier, MIN(score) AS score
                FROM candidates
                GROUP BY rowid, reference, text
            )
            SELECT reference, text
            FROM ranked
            ORDER BY tier, score
            LIMIT ? OFFSET ?
            """,
            (
                self.text_fts_query(phrase_query), *scope_params,
                self.text_fts_query(terms_query), *scope_params,
                limit, offset,
            )
        )

    def count_search_results(self, fts_query, scope=None):
        scope_clause, scope_params = self.search_scope_clause(scope)
        rows = self.bible_db.query(
            f"SELECT COUNT(*) FROM bible_fts WHERE bible_fts MATCH ? AND reference NOT GLOB '*:0'{scope_clause}",
            (self.text_fts_query(fts_query), *scope_params)
        )
        return rows[0][0] if rows else 0

    def resolve_search(self, query, page=1, limit=10):
        search_query, scope = self.parse_search_query_and_scope(query)
        effective_query = search_query
        phrase_query = self.escape_fts_query(effective_query)
        terms_query = self.fts_terms_query(effective_query)
        if not phrase_query or not terms_query:
            return None
        corrections = []
        page = max(1, page)
        limit = min(max(1, limit), 50)

        total = self.count_search_results(terms_query, scope=scope)
        if not total:
            corrected_query, corrections = self.fuzzy_search_query(search_query)
            if corrections:
                corrected_phrase = self.escape_fts_query(corrected_query)
                corrected_terms = self.fts_terms_query(corrected_query)
                corrected_total = self.count_search_results(corrected_terms, scope=scope)
                if corrected_total:
                    effective_query = corrected_query
                    phrase_query = corrected_phrase
                    terms_query = corrected_terms
                    total = corrected_total

        mode = "term"
        active_query = terms_query
        phrase_total = 0
        if total:
            phrase_total = self.count_search_results(phrase_query, scope=scope)
            if phrase_query != terms_query and phrase_total:
                mode = "phrase first" if phrase_total < total else "phrase"
                active_query = phrase_query if mode == "phrase" else terms_query
            elif phrase_query != terms_query:
                mode = "all terms"
        else:
            active_query = self.fts_any_terms_query(effective_query)
            total = self.count_search_results(active_query, scope=scope) if active_query else 0
            mode = "any terms"

        if total:
            page = min(page, ((total - 1) // limit) + 1)
        offset = (page - 1) * limit
        if mode == "phrase first":
            res = self.query_ranked_search_results(
                phrase_query, terms_query, limit, offset, scope=scope
            )
        else:
            res = self.query_search_results(active_query, limit, offset, scope=scope) if total else []
        if not res:
            return None
        return {
            "query": search_query,
            "effective_query": effective_query,
            "corrections": corrections,
            "active_query": active_query,
            "mode": mode,
            "scope": scope,
            "page": page,
            "limit": limit,
            "total": total,
            "offset": offset,
            "results": res,
            "page_count": ((total - 1) // limit) + 1 if total else 1,
        }

    def render_search_page(self, state, interactive=False):
        body = Text()
        query = state["query"]
        effective_query = state.get("effective_query", query)
        page = state["page"]
        limit = state["limit"]
        offset = state["offset"]
        total = state["total"]
        res = state["results"]
        scope = state.get("scope")
        for idx, (ref, text) in enumerate(res, 1):
            parts = self.parse_reference_parts(ref)
            display_ref = self.format_display_ref(ref) if parts else ref
            body.append(f"{offset + idx:>3}. {display_ref}\n", style="verse.ref")
            body.append_text(self.highlight_search_terms(self.clean_text(text), effective_query))
            body.append("\n\n", style="dim")
        shown_end = offset + len(res)
        scope_label = f"  |  Scope: {scope['label']}" if scope else ""
        footer = ""
        if state.get("corrections"):
            changes = ", ".join(f"{old} → {new}" for old, new in state["corrections"])
            footer += f"Corrected: {changes}\n"
        footer += f"Mode: {state['mode']}{scope_label}  |  Showing {offset + 1}-{shown_end} of {total}"
        query_arg = shlex.quote(query)
        if scope:
            query_arg = f"{query_arg} -{scope['label'].lower().replace(' ', '-').replace('--', '-')}"
        if interactive:
            footer += "  |  Choose an action below"
        elif shown_end < total:
            footer += f"\nNext page: lex search {query_arg} --page {page + 1}"
            if limit != 10:
                footer += f" --limit {limit}"
        if not interactive and page > 1:
            footer += f"\nPrevious page: lex search {query_arg} --page {page - 1}"
            if limit != 10:
                footer += f" --limit {limit}"
        body.append(footer, style="ui.meta")
        console.print(
            Panel(
                body,
                title=f"🔍 Search: {query}",
                subtitle=f"page {page}/{state['page_count']}",
                border_style="ui.border",
                padding=(1, 2),
            )
        )
        if interactive:
            self.render_action_bar(
                "Actions",
                [
                    ("1-10", "study result"),
                    ("r #", "read result"),
                    ("n / p", "page"),
                    ("e", "export"),
                    ("q", "quit"),
                ],
            )

    def search_export_dir(self):
        path = os.path.expanduser("~/Documents/lex_exports")
        os.makedirs(path, exist_ok=True)
        return path

    def render_action_bar(self, title, actions):
        grid = Table.grid(padding=(0, 2))
        grid.add_column(no_wrap=True)
        grid.add_column(style="ui.meta")
        for key, label in actions:
            grid.add_row(f"[ui.action.key]{key}[/]", label)
        console.print(Panel(grid, title=title, border_style="ui.action", padding=(0, 1), expand=True))

    def open_export(self, path):
        if not path:
            return
        if sys.platform.startswith("win"):
            try:
                os.startfile(path)
                console.print(f"[success]Saved and opened:[/] {path}")
            except Exception:
                console.print(f"[success]Saved:[/] {path}")
            return
        opener = None
        candidates = ("open", "xdg-open", "gio", "kde-open") if sys.platform == "darwin" else ("xdg-open", "gio", "kde-open", "open")
        for candidate in candidates:
            candidate_path = shutil.which(candidate)
            if candidate_path:
                opener = [candidate_path]
                if candidate == "gio":
                    opener.append("open")
                break
        if not opener:
            console.print(f"[dim]Saved file:[/] {path}")
            return
        try:
            subprocess.Popen([*opener, path], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            console.print(f"[success]Saved and opened:[/] {path}")
        except Exception:
            console.print(f"[success]Saved:[/] {path}")

    def open_exports_folder(self, path=None):
        folder = path or self.search_export_dir()
        self.open_export(folder)

    def pdf_safe_text(self, value):
        text = "" if value is None else str(value)
        replacements = {
            "•": "-",
            "→": "->",
            "←": "<-",
            "\u00a0": " ",
        }
        for old, new in replacements.items():
            text = text.replace(old, new)
        return html.escape(text)

    def setup_pdf_styles(self, styles):
        paragraph_style_cls = None
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.lib.styles import ParagraphStyle
            paragraph_style_cls = ParagraphStyle
            win_fonts = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
            bundled_fonts = os.path.join(RUNTIME_DATA_DIR, "fonts")
            common_unicode_fonts = [
                os.environ.get("LEX_PDF_FONT"),
                os.path.join(bundled_fonts, "NotoSans-Regular.ttf"),
                os.path.join(bundled_fonts, "DejaVuSans.ttf"),
                "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
                "/Library/Fonts/Arial Unicode.ttf",
                "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
                "/usr/share/fonts/opentype/noto/NotoSans-Regular.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
                os.path.join(win_fonts, "arialuni.ttf"),
                os.path.join(win_fonts, "arial.ttf"),
                os.path.join(win_fonts, "segoeui.ttf"),
            ]
            font_paths = {
                "LexSans": common_unicode_fonts,
                "LexSansHebrew": [
                    os.environ.get("LEX_PDF_HEBREW_FONT"),
                    os.path.join(bundled_fonts, "NotoSansHebrew-Regular.ttf"),
                    "/usr/share/fonts/truetype/noto/NotoSansHebrew-Regular.ttf",
                    "/usr/share/fonts/opentype/noto/NotoSansHebrew-Regular.ttf",
                    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
                    "/Library/Fonts/Arial Unicode.ttf",
                    os.path.join(win_fonts, "arialuni.ttf"),
                    os.path.join(win_fonts, "arial.ttf"),
                    *common_unicode_fonts,
                ],
            }
            for font_name, candidates in font_paths.items():
                for font_path in candidates:
                    if font_path and os.path.exists(font_path) and font_name not in pdfmetrics.getRegisteredFontNames():
                        pdfmetrics.registerFont(TTFont(font_name, font_path))
                        break
            base_font = "LexSans" if "LexSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica"
            hebrew_font = "LexSansHebrew" if "LexSansHebrew" in pdfmetrics.getRegisteredFontNames() else base_font
        except Exception:
            base_font = "Helvetica"
            hebrew_font = base_font
        for style_name in ["Title", "Heading1", "Heading2", "Heading3", "Normal", "BodyText", "Italic"]:
            if style_name in styles:
                styles[style_name].fontName = base_font
        if "Hebrew" not in styles and paragraph_style_cls:
            styles.add(paragraph_style_cls(name="Hebrew", parent=styles["BodyText"]))
        if "Hebrew" in styles:
            styles["Hebrew"].fontName = hebrew_font
        return base_font, hebrew_font

    def pdf_paragraph(self, text, style):
        from reportlab.platypus import Paragraph
        return Paragraph(self.pdf_safe_text(text), style)

    def search_export_filename(self, state, ext):
        scope = state.get("scope")
        scope_part = f"_{scope['label']}" if scope else ""
        raw = f"lex_search_{state['query']}{scope_part}_p{state['page']}.{ext}"
        safe = re.sub(r"[^A-Za-z0-9._-]+", "_", raw).strip("_")
        return os.path.join(self.search_export_dir(), safe)

    def search_export_rows(self, state):
        rows = []
        for idx, (ref, text) in enumerate(state["results"], 1):
            parts = self.parse_reference_parts(ref)
            display_ref = self.format_display_ref(ref) if parts else ref
            rows.append({
                "number": state["offset"] + idx,
                "reference": display_ref,
                "text": self.clean_text(text),
            })
        return rows

    def export_search_docx(self, state):
        try:
            from docx import Document
        except ImportError:
            console.print("[warning]DOCX export needs python-docx installed.[/]")
            return None
        path = self.search_export_filename(state, "docx")
        doc = Document()
        doc.add_heading(f"Lex Search: {state['query']}", level=1)
        scope = state.get("scope")
        meta = f"Mode: {state['mode']} | Page: {state['page']}/{state['page_count']} | Showing {state['offset'] + 1}-{state['offset'] + len(state['results'])} of {state['total']}"
        if scope:
            meta += f" | Scope: {scope['label']}"
        doc.add_paragraph(meta)
        for row in self.search_export_rows(state):
            doc.add_heading(f"{row['number']}. {row['reference']}", level=2)
            doc.add_paragraph(row["text"])
        doc.save(path)
        return path

    def export_search_pdf(self, state):
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        except ImportError:
            console.print("[warning]PDF export needs reportlab installed.[/]")
            return None
        path = self.search_export_filename(state, "pdf")
        doc = SimpleDocTemplate(path, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
        styles = getSampleStyleSheet()
        self.setup_pdf_styles(styles)
        story = [self.pdf_paragraph(f"Lex Search: {state['query']}", styles["Title"])]
        scope = state.get("scope")
        meta = f"Mode: {state['mode']} | Page: {state['page']}/{state['page_count']} | Showing {state['offset'] + 1}-{state['offset'] + len(state['results'])} of {state['total']}"
        if scope:
            meta += f" | Scope: {scope['label']}"
        story.extend([self.pdf_paragraph(meta, styles["Normal"]), Spacer(1, 12)])
        for row in self.search_export_rows(state):
            story.append(self.pdf_paragraph(f"{row['number']}. {row['reference']}", styles["Heading2"]))
            story.append(self.pdf_paragraph(row["text"], styles["BodyText"]))
            story.append(Spacer(1, 10))
        doc.build(story)
        return path

    def export_search_pptx(self, state):
        path = self.search_export_filename(state, "pptx")
        verses = []
        for row in self.search_export_rows(state):
            verses.append((row["reference"], row["text"]))
        return self.export_verses_pptx(verses, path, f"Lex Search: {state['query']}")

    def export_study_pptx(self, db_ref):
        path = self.study_export_filename(db_ref, "pptx")
        res = self.bible_db.query("SELECT reference, text FROM bible WHERE reference = ? LIMIT 1", (db_ref,))
        if not res:
            return None
        ref, text = res[0]
        display_ref = self.format_display_ref(ref)
        return self.export_verses_pptx([(display_ref, self.clean_text(text))], path, f"Lex Study: {display_ref}")

    def export_verses_pptx(self, verses, path, title):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
        except ImportError:
            console.print("[warning]PowerPoint export needs python-pptx installed.[/]")
            return None

        prs = Presentation()
        # Set to 16:9 Widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Dynamic Theme Colors
        if ACTIVE_THEME_MODE == "light":
            bg_color = RGBColor(249, 247, 242)
            box_color = RGBColor(240, 235, 225)
            text_color = RGBColor(31, 31, 31)
            footer_bg = RGBColor(220, 210, 190)
            footer_text = RGBColor(100, 100, 100)
        else:
            bg_color = RGBColor(28, 28, 28) # grey11 roughly
            box_color = bg_color # Bucket effect: match background
            text_color = RGBColor(235, 235, 235)
            footer_bg = bg_color # Bucket effect: match background
            footer_text = RGBColor(120, 120, 120) # Muted for floating look

        # layout 6 is usually blank in standard templates
        try:
            blank_slide_layout = prs.slide_layouts[6]
        except IndexError:
            blank_slide_layout = prs.slide_layouts[0]

        for ref, text in verses:
            parts = self.parse_reference_parts(ref)
            verse_no = str(parts["verse"]) if parts else ""
            display_ref = self.format_display_ref(ref) if parts else ref

            slide = prs.slides.add_slide(blank_slide_layout)

            # Set background
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color

            # Add Box for text
            box_width = Inches(11)
            box_height = Inches(4.5)
            box_left = (prs.slide_width - box_width) / 2
            box_top = Inches(1.2)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = box_color
            shape.line.fill.background() # No border

            # Add text to box
            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            # Marker and Verse Number
            marker_text = f">  {verse_no.rjust(3)} " if verse_no else ">  "
            p.text = f"{marker_text}{text}"
            p.font.name = "Courier New"
            p.font.size = Pt(32)
            p.font.color.rgb = text_color

            # Footer / Reference Box
            footer_width = Inches(5)
            footer_height = Inches(0.8)
            footer_left = (prs.slide_width - footer_width) / 2
            footer_top = box_top + box_height + Inches(0.4)

            footer_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, footer_left, footer_top, footer_width, footer_height
            )
            footer_shape.fill.solid()
            footer_shape.fill.fore_color.rgb = footer_bg
            footer_shape.line.fill.background()

            ftf = footer_shape.text_frame
            ftf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            fp = ftf.paragraphs[0]
            fp.alignment = PP_ALIGN.CENTER
            fp.text = f"- {display_ref}"
            fp.font.name = "Courier New"
            fp.font.size = Pt(28)
            fp.font.color.rgb = footer_text

        prs.save(path)
        return path

    def display_export_howto(self):
        md = """
# Lex Export
*Bulk Scripture processing for presentations and study*

The `export` command allows you to process multiple verses into high-quality
images, presentations, or study documents.

## Usage

```bash
lex export "John 3:16"            # Interactive menu
lex export "John 3:16" --mode read --format png --ratio 1:1 --theme dark
lex export "Romans 8:28-39" --mode read --format pptx
lex export "Ephesians 1" --mode study --format pdf
```

## Modes & Formats

| Mode | Target Content | Formats |
| :--- | :--- | :--- |
| **read** | Verse text only (CLI aesthetic) | `png`, `pptx` |
| **study**| Interlinear, notes, TSK | `pdf`, `docx` |
| **web**  | Verse + cross-references | `pdf`, `docx` |

## Reference Strings

You can provide a single reference, a range, or a comma-separated list of both.
Chapter-only references (e.g., `John 1`) will export every verse in that chapter.

*   `"John 3:16, John 1:1-5"`
*   `"Genesis 1; Psalm 23"`
*   `"Romans 8:28-39"`

---
*Created for the Church and the Terminal.*
"""
        console.print(Panel(Markdown(md), title="📤 Bulk Export", border_style="ui.border", expand=False))

    def handle_bulk_export(self, refs_str, mode="read", format="pptx", ratio="16:9", theme="auto"):
        # Temporarily override theme for export if not 'auto'
        global ACTIVE_THEME_MODE
        original_theme = ACTIVE_THEME_MODE
        if theme != "auto":
            ACTIVE_THEME_MODE = theme

        try:
            # Split by comma or semicolon
            import re
            ref_list = re.split(r'[,;]', refs_str)
            ref_list = [r.strip() for r in ref_list if r.strip()]

            if not ref_list:
                console.print("[error]No valid references found to export.[/]")
                return

            all_verses = []
            for raw_ref in ref_list:
                ref_norm, book, chap, verse, v_end = self.normalize_ref(raw_ref)
                if not ref_norm:
                    console.print(f"[warning]Skipping invalid reference: {raw_ref}[/]")
                    continue

                db_book = self.canon_map.get(re.sub(r"[^a-z0-9]+", "", book.lower()), book)
                if v_end:
                    start_ref = f"{self.bible_prefix}:{db_book}:{chap}:{verse}"
                    end_ref = f"{self.bible_prefix}:{db_book}:{chap}:{v_end}"
                    res = self.bible_db.query(
                        """
                        SELECT reference, text FROM bible
                        WHERE id >= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
                          AND id <= (SELECT id FROM bible WHERE reference = ? LIMIT 1)
                        ORDER BY id
                        """,
                        (start_ref, end_ref)
                    )
                    for r, t in res:
                        all_verses.append((r, t))
                elif verse:
                    res = self.bible_db.query(
                        "SELECT reference, text FROM bible WHERE reference LIKE ? GROUP BY reference LIMIT 1",
                        (f"%:{db_book}:{chap}:{verse}",)
                    )
                    if res: all_verses.append(res[0])
                else:
                    # Chapter
                    res = self.bible_db.query(
                        "SELECT reference, text FROM bible WHERE reference LIKE ? AND reference NOT LIKE '%:0' ORDER BY id",
                        (f"%:{db_book}:{chap}:%",)
                    )
                    for r, t in res:
                        all_verses.append((r, t))

            if not all_verses:
                console.print("[error]No verses found for the given references.[/]")
                return

            if mode == "read":
                if format == "pptx":
                    path = self.search_export_filename({"query": "bulk_read", "page": 1}, "pptx")
                    display_verses = [(self.format_display_ref(r), self.clean_text(t)) for r, t in all_verses]
                    final_path = self.export_verses_pptx(display_verses, path, "Lex Bulk Read")
                    if final_path: self.open_export(final_path)
                elif format == "png":
                    exported_files = []
                    for r, t in all_verses:
                        display_ref = self.format_display_ref(r)
                        path = self.study_export_filename(r, "png")
                        final_path = self.export_read_png(display_ref, self.clean_text(t), path, ratio=ratio)
                        if final_path:
                            console.print(f"[success]Exported: {final_path}[/]")
                            exported_files.append(final_path)
                    if exported_files:
                        self.open_exports_folder(self.study_export_dir())
                else:
                    console.print(f"[error]Format {format} not supported for read mode. Use png or pptx.[/]")

            elif mode == "study":
                if format == "pdf":
                    path = self.search_export_filename({"query": "bulk_study", "page": 1}, "pdf")
                    final_path = self.export_bulk_study_pdf(all_verses, path)
                    if final_path: self.open_export(final_path)
                elif format == "docx":
                    path = self.search_export_filename({"query": "bulk_study", "page": 1}, "docx")
                    final_path = self.export_bulk_study_docx(all_verses, path)
                    if final_path: self.open_export(final_path)
                else:
                    console.print(f"[error]Format {format} not supported for study mode. Use pdf or docx.[/]")

            elif mode == "web":
                if format == "pdf":
                    path = self.search_export_filename({"query": "bulk_web", "page": 1}, "pdf")
                    final_path = self.export_bulk_web_pdf(all_verses, path)
                    if final_path: self.open_export(final_path)
                elif format == "docx":
                    path = self.search_export_filename({"query": "bulk_web", "page": 1}, "docx")
                    final_path = self.export_bulk_web_docx(all_verses, path)
                    if final_path: self.open_export(final_path)
                else:
                    console.print(f"[error]Format {format} not supported for web mode. Use pdf or docx.[/]")
        finally:
            ACTIVE_THEME_MODE = original_theme
    def export_read_png(self, ref, text, path, ratio="16:9"):
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            console.print("[warning]PNG export needs Pillow installed.[/]")
            return None

        # Dimensions based on ratio
        if ratio == "1:1":
            width = 1080
            height = 1080
            box_h = 700
            box_y = 100
        else: # 16:9
            width = 1920
            height = 1080
            box_h = 600
            box_y = 150

        # Dynamic Theme Colors
        if ACTIVE_THEME_MODE == "light":
            bg_color = (249, 247, 242)
            box_color = (240, 235, 225)
            text_color = (31, 31, 31)
            footer_bg = (220, 210, 190)
            footer_text = (100, 100, 100)
        else:
            bg_color = (28, 28, 28)
            box_color = bg_color
            text_color = (235, 235, 235)
            footer_bg = bg_color
            footer_text = (120, 120, 120)

        img = Image.new('RGB', (width, height), color=bg_color)
        draw = ImageDraw.Draw(img)

        # Fonts
        try:
            # Common Darwin paths for Courier New
            font_path = "/System/Library/Fonts/Supplemental/Courier New.ttf"
            if not os.path.exists(font_path):
                font_path = "/Library/Fonts/Courier New.ttf"

            main_font = ImageFont.truetype(font_path, 48)
            footer_font = ImageFont.truetype(font_path, 40)

            # Attempt to use emoji-capable font if available for text drawing
            emoji_font_path = "/System/Library/Fonts/Apple Color Emoji.ttc"
            if os.path.exists(emoji_font_path):
                # We need a font that supports both text and emojis.
                # As a fallback, we use the main font, but this is a complex limitation.
                pass
        except:
            main_font = ImageFont.load_default()
            footer_font = ImageFont.load_default()

        # Draw Verse Box
        box_padding = 100
        box_w = width - (box_padding * 2)
        box_x = box_padding

        draw.rectangle([box_x, box_y, box_x + box_w, box_y + box_h], fill=box_color)

        # Handle word wrap
        parts = self.parse_reference_parts(ref)
        verse_no = str(parts["verse"]) if parts else ""
        marker_text = f">  {verse_no.rjust(3)} " if verse_no else ">  "
        full_text = f"{marker_text}{text}"

        words = full_text.split(' ')
        lines = []
        current_line = []
        for word in words:
            test_line = ' '.join(current_line + [word])
            # getbbox returns (left, top, right, bottom)
            w = draw.textbbox((0, 0), test_line, font=main_font)[2]
            if w < box_w - 100:
                current_line.append(word)
            else:
                lines.append(' '.join(current_line))
                current_line = [word]
        lines.append(' '.join(current_line))

        # Draw lines centered vertically in box
        line_height = draw.textbbox((0, 0), "Ay", font=main_font)[3] + 20
        total_text_h = len(lines) * line_height
        start_y = box_y + (box_h - total_text_h) / 2

        for i, line in enumerate(lines):
            draw.text((box_x + 50, start_y + (i * line_height)), line, font=main_font, fill=text_color)

        # Draw Footer Box
        footer_w = 600 if ratio == "1:1" else 700
        footer_h = 100
        footer_x = (width - footer_w) / 2
        footer_y = box_y + box_h + (20 if ratio == "1:1" else 50)

        draw.rectangle([footer_x, footer_y, footer_x + footer_w, footer_y + footer_h], fill=footer_bg)

        display_ref = f"- {ref}"
        rw = draw.textbbox((0, 0), display_ref, font=footer_font)[2]
        rh = draw.textbbox((0, 0), display_ref, font=footer_font)[3]
        draw.text((footer_x + (footer_w - rw) / 2, footer_y + (footer_h - rh) / 2 - 5), display_ref, font=footer_font, fill=footer_text)

        img.save(path)
        return path

    def export_bulk_study_docx(self, verses, path):
        try:
            from docx import Document
        except ImportError:
            console.print("[warning]DOCX export needs python-docx installed.[/]")
            return None
        doc = Document()
        for ref, text in verses:
            data = self.build_study_export_data(ref)
            if data:
                self.add_study_to_docx(doc, data)
                doc.add_page_break()
        doc.save(path)
        return path

    def export_bulk_study_pdf(self, verses, path):
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, PageBreak
        except ImportError:
            console.print("[warning]PDF export needs reportlab installed.[/]")
            return None
        doc = SimpleDocTemplate(path, pagesize=letter)
        story = []
        for ref, text in verses:
            data = self.build_study_export_data(ref)
            if data:
                story.extend(self.build_study_pdf_story(data))
                story.append(PageBreak())
        doc.build(story)
        return path

    def export_bulk_web_docx(self, verses, path):
        try:
            from docx import Document
        except ImportError:
            console.print("[warning]DOCX export needs python-docx installed.[/]")
            return None
        doc = Document()
        for ref, text in verses:
            data = self.build_web_export_data(ref)
            if data:
                self.add_web_to_docx(doc, data)
                doc.add_page_break()
        doc.save(path)
        return path

    def export_bulk_web_pdf(self, verses, path):
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, PageBreak
        except ImportError:
            console.print("[warning]PDF export needs reportlab installed.[/]")
            return None
        doc = SimpleDocTemplate(path, pagesize=letter)
        story = []
        for ref, text in verses:
            data = self.build_web_export_data(ref)
            if data:
                story.extend(self.build_web_pdf_story(data))
                story.append(PageBreak())
        doc.build(story)
        return path

    def build_web_export_data(self, db_ref):
        # Resolve DB ref if needed
        if ":" not in db_ref:
            ref_norm, book, chap, verse, v_end = self.normalize_ref(db_ref)
            if not ref_norm: return None
            db_ref = f"{self.bible_prefix}:{ref_norm}"

        parts = self.parse_reference_parts(db_ref)
        if not parts or not parts["verse"]: return None

        verse_row = self.bible_db.query("SELECT text FROM bible WHERE reference = ? LIMIT 1", (db_ref,))
        if not verse_row: return None

        display_ref = self.format_display_ref(db_ref)

        connections = []
        for target, votes in self.get_tsk_crossrefs(db_ref)[:15]:
            text = self.get_crossref_text(target)
            if text:
                connections.append({
                    "reference": self.format_tsk_display_ref(target),
                    "votes": votes,
                    "text": text,
                })

        return {
            "display_ref": display_ref,
            "verse": self.clean_text(verse_row[0][0]),
            "connections": connections
        }

    def add_study_to_docx(self, doc, data):
        doc.add_heading(f"Lex Study: {data['display_ref']}", level=1)
        if data["verse"]:
            doc.add_paragraph(data["verse"])
        doc.add_heading(data["language"], level=2)
        if data["source"]:
            doc.add_paragraph(data["source"])
        if data["transliteration"]:
            doc.add_paragraph(data["transliteration"])
        doc.add_heading("Interlinear", level=2)
        table = doc.add_table(rows=1, cols=5)
        for cell, title in zip(table.rows[0].cells, ["English", "Source", "Lemma", "Code", "Gloss"]):
            cell.text = title
        for parsed in data["interlinear"]:
            row = table.add_row().cells
            row[0].text = parsed["english"] or "-"
            row[1].text = f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "-"
            row[2].text = f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "-"
            row[3].text = parsed["strongs"] or parsed["morph"] or "-"
            row[4].text = parsed["gloss"] or parsed["english"] or "-"
        doc.add_heading("Lexicon Notes", level=2)
        for note in data["lex_notes"]:
            doc.add_paragraph(f"{note['strongs']} - {note['lemma']}: {note['details']}")
        if data["topical_refs"]:
            doc.add_heading("Topical Associations (Nave's)", level=2)
            doc.add_paragraph(" • ".join(data["topical_refs"]))
        doc.add_heading("Treasury of Scripture Knowledge", level=2)
        for ref in data["tsk_refs"]:
            doc.add_paragraph(f"{ref['reference']} ({ref['votes']}): {ref['preview']}")

    def build_study_pdf_story(self, data):
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, Spacer, Table as PdfTable
        styles = getSampleStyleSheet()
        self.setup_pdf_styles(styles)
        story = [self.pdf_paragraph(f"Lex Study: {data['display_ref']}", styles["Title"])]
        if data["verse"]:
            story.extend([self.pdf_paragraph(data["verse"], styles["BodyText"]), Spacer(1, 10)])
        story.append(self.pdf_paragraph(data["language"], styles["Heading2"]))
        if data["source"]:
            story.append(self.pdf_paragraph(data["source"], styles["BodyText"]))
        if data["transliteration"]:
            story.append(self.pdf_paragraph(data["transliteration"], styles["Italic"]))
        story.extend([Spacer(1, 10), self.pdf_paragraph("Interlinear", styles["Heading2"])])
        table_rows = [["English", "Source", "Lemma", "Code", "Gloss"]]
        for parsed in data["interlinear"]:
            table_rows.append([
                self.pdf_paragraph(parsed["english"] or "-", styles["BodyText"]),
                self.pdf_paragraph(f"{parsed['surface']} ({parsed['translit']})" if parsed["surface"] else "-", styles["BodyText"]),
                self.pdf_paragraph(f"{parsed['lemma']} ({parsed['lemma_translit']})" if parsed["lemma"] else "-", styles["BodyText"]),
                self.pdf_paragraph(parsed["strongs"] or parsed["morph"] or "-", styles["BodyText"]),
                self.pdf_paragraph(parsed["gloss"] or parsed["english"] or "-", styles["BodyText"]),
            ])
        story.append(PdfTable(table_rows, repeatRows=1))
        story.append(self.pdf_paragraph("Lexicon Notes", styles["Heading2"]))
        for note in data["lex_notes"]:
            story.append(self.pdf_paragraph(f"{note['strongs']} - {note['lemma']}: {note['details']}", styles["BodyText"]))
        if data["topical_refs"]:
            story.append(self.pdf_paragraph("Topical Associations (Nave's)", styles["Heading2"]))
            story.append(self.pdf_paragraph(" • ".join(data["topical_refs"]), styles["BodyText"]))
        story.append(self.pdf_paragraph("Treasury of Scripture Knowledge", styles["Heading2"]))
        for ref in data["tsk_refs"]:
            story.append(self.pdf_paragraph(f"{ref['reference']} ({ref['votes']}): {ref['preview']}", styles["BodyText"]))
        return story

    def add_web_to_docx(self, doc, data):
        doc.add_heading(f"Lex Verse Web: {data['display_ref']}", level=1)
        doc.add_paragraph(data["verse"])
        doc.add_heading("Cross-Reference Connections", level=2)
        for conn in data["connections"]:
            doc.add_heading(f"{conn['reference']} (Relevance: {conn['votes']})", level=3)
            doc.add_paragraph(conn["text"])

    def build_web_pdf_story(self, data):
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, Spacer
        styles = getSampleStyleSheet()
        self.setup_pdf_styles(styles)
        story = [self.pdf_paragraph(f"Lex Verse Web: {data['display_ref']}", styles["Title"])]
        story.extend([self.pdf_paragraph(data["verse"], styles["BodyText"]), Spacer(1, 12)])
        story.append(self.pdf_paragraph("Cross-Reference Connections", styles["Heading2"]))
        for conn in data["connections"]:
            story.append(self.pdf_paragraph(f"{conn['reference']} (Relevance: {conn['votes']})", styles["Heading3"]))
            story.append(self.pdf_paragraph(conn["text"], styles["BodyText"]))
            story.append(Spacer(1, 8))
        return story

    def prompt_search_export(self, state):
        while True:
            self.render_action_bar(
                "Export",
                [
                    ("d", "DOCX current page"),
                    ("f", "PDF current page"),
                    ("p", "PPTX current page"),
                    ("o", "open exports folder"),
                    ("q", "back"),
                ],
            )
            action = Prompt.ask("Export action", choices=["d", "f", "p", "o", "q"], default="q").lower()
            if action == "q":
                return
            if action == "o":
                self.open_exports_folder()
                continue

            if action == "d":
                path = self.export_search_docx(state)
            elif action == "f":
                path = self.export_search_pdf(state)
            else:
                path = self.export_search_pptx(state)
            if path:
                self.open_export(path)
                return

    def search_result_ref(self, state, user_number):
        if not user_number.isdigit():
            return None
        idx = int(user_number) - state["offset"] - 1
        if idx < 0 or idx >= len(state["results"]):
            return None
        return state["results"][idx][0]

    def display_search(self, query, page=1, limit=10, interactive=None):
        state = self.resolve_search(query, page=page, limit=limit)
        if not state:
            return False
        use_interactive = console.is_terminal if interactive is None else interactive
        use_interactive = use_interactive and page == 1 and state["page_count"] > 1
        if not use_interactive:
            self.render_search_page(state)
            return True
        while True:
            console.clear()
            self.render_search_page(state, interactive=True)
            nav = Prompt.ask("Search action", default="q").strip().lower()
            if nav == "q":
                return True
            if nav == "e":
                self.prompt_search_export(state)
                Prompt.ask("Press Enter to continue", default="")
                continue
            if nav.isdigit():
                ref = self.search_result_ref(state, nav)
                if ref:
                    self.display_study(ref, actions=True)
                continue
            read_match = re.match(r"^r\s+(\d+)$", nav)
            if read_match:
                ref = self.search_result_ref(state, read_match.group(1))
                if ref:
                    self.display_verse(self.format_display_ref(ref))
                    Prompt.ask("Press Enter to return to search", default="")
                continue
            if nav not in {"n", "p"}:
                continue
            next_page = state["page"] + 1 if nav == "n" else state["page"] - 1
            if next_page < 1 or next_page > state["page_count"]:
                continue
            next_state = self.resolve_search(query, page=next_page, limit=limit)
            if next_state:
                state = next_state
        return True

    def display_strongs(self, query):
        if re.match(r'^[GH]\d+$', query.upper()):
            res = self.strongs_db.query("SELECT number, word, pronunciation, definition FROM strongs WHERE number = ?", (query.upper(),))
        else:
            normalized = self.normalize_term(query)
            res = self.strongs_db.query(
                """
                SELECT number, word, pronunciation, definition
                FROM strongs
                WHERE lower(replace(replace(replace(pronunciation, '''', ''), '-', ''), ' ', '')) = ?
                LIMIT 5
                """,
                (normalized,)
            )
            if not res:
                safe_query = self.escape_fts_query(query)
                if not safe_query:
                    return False
                res = self.strongs_db.query(
                    """
                    SELECT s.number, s.word, s.pronunciation, s.definition
                    FROM strongs_fts f
                    JOIN strongs s ON s.number = f.number
                    WHERE strongs_fts MATCH ?
                    LIMIT 5
                    """,
                    (safe_query,)
                )
        for n, w, p, d in res:
            lang = "Greek" if n.startswith('G') else "Hebrew"
            extended_entry = self.lookup_lexicon_entry(n)

            definition_text = d
            if extended_entry["step"] and extended_entry["step"].get("definition"):
                step_def = extended_entry["step"]["definition"]
                step_def = re.sub(r'<br\s*/?>', '\n', step_def, flags=re.IGNORECASE)
                step_def = re.sub(r'<b>(.*?)</b>', r'[bold]\1[/]', step_def, flags=re.IGNORECASE|re.DOTALL)
                step_def = re.sub(r'<i>(.*?)</i>', r'[italic]\1[/]', step_def, flags=re.IGNORECASE|re.DOTALL)
                step_def = re.sub(r'<ref[^>]*>(.*?)</ref>', r'[cyan]\1[/]', step_def, flags=re.IGNORECASE|re.DOTALL)
                step_def = re.sub(r'<[^>]+>', '', step_def)
                source = extended_entry["step"].get("source", "Extended Lexicon")
                definition_text = f"{step_def}\n\n[dim]---\nSource: {source}\nBrief: {d}[/]"


            # Fetch major uses (concordance) from KJV if available
            major_uses = ""
            kjv_path = get_bible_path("kjv")
            if os.path.exists(kjv_path):
                with sqlite3.connect(kjv_path) as conn:
                    c = conn.cursor()
                    pattern = f"%[{n}]%"
                    c.execute("SELECT reference, text FROM bible WHERE text LIKE ? LIMIT 5", (pattern,))
                    rows = c.fetchall()
                    if rows:
                        major_uses = "\n\n[bold underline]Major Uses (KJV):[/]"
                        for ref, txt in rows:
                            disp_ref = self.format_display_ref(ref)
                            target_tag = f"[{n}]"
                            # Mark target
                            highlighted = txt.replace(target_tag, "__TARGET_TAG__")
                            # Strip all other Strong's tags (handles [G123], <G123>, etc.)
                            highlighted = re.sub(r'[<\[][GH]\d+[>\]]', '', highlighted)
                            # Strip remaining HTML
                            highlighted = re.sub(r'<[^>]+>', '', highlighted)
                            # Restore target with highlight
                            highlighted = highlighted.replace("__TARGET_TAG__", f"[bold yellow]{target_tag}[/]")
                            major_uses += f"\n[cyan]{disp_ref}[/] {highlighted.strip()}"

                        c.execute("SELECT COUNT(*) FROM bible WHERE text LIKE ?", (pattern,))
                        total = c.fetchone()[0]
                        if total > 5:
                            major_uses += f"\n[dim]... and {total - 5} more occurrences.[/]"

            definition_text += major_uses
            console.print(Panel(f"[lexicon.word]{w}[/] ({p})\n\n{definition_text}", title=f"📚 {lang} Lexicon: {n}", border_style="blue"))
        return bool(res)

    def display_english_strongs(self, query):
        normalized = self.normalize_term(query)
        if not normalized:
            return False
        exact_results = []
        fuzzy_results = []
        seen = set()
        for strongs_id, entry in self.get_interlinear_strongs().items():
            glosses = self.extract_english_glosses(entry)
            exact_matches = [gloss for gloss in glosses if self.normalize_term(gloss) == normalized]
            fuzzy_matches = [gloss for gloss in glosses if normalized in self.normalize_term(gloss)]
            if not exact_matches and not fuzzy_matches:
                continue
            _, db_key, _ = self.normalize_strongs_key(strongs_id)
            db_rows = self.strongs_db.query(
                "SELECT number, word, pronunciation, definition FROM strongs WHERE number = ?",
                (db_key,)
            ) if db_key else []
            if not db_rows or db_key in seen:
                continue
            seen.add(db_key)
            number, word, pronunciation, definition = db_rows[0]
            item = (number, word, pronunciation, definition, ", ".join((exact_matches or fuzzy_matches)[:3]))
            if exact_matches:
                exact_results.append(item)
            else:
                fuzzy_results.append(item)
        results = exact_results[:8] if exact_results else fuzzy_results[:8]
        safe_query = self.escape_fts_query(query)
        if safe_query:
            db_rows = self.strongs_db.query(
                """
                SELECT strongs.number, strongs.word, strongs.pronunciation, strongs.definition
                FROM strongs_fts
                JOIN strongs USING(number)
                WHERE strongs_fts MATCH ?
                LIMIT 12
                """,
                (safe_query,)
            )
            for n, w, p, d in db_rows:
                if n in seen:
                    continue
                results.append((n, w, p, d, query))
                seen.add(n)
                if len(results) >= 12:
                    break
        if not results:
            return False
        table = Table(title=f"🔤 Strong's Lookup: '{query}'", box=None)
        table.add_column("No.", style="lexicon.num")
        table.add_column("Lemma", style="lexicon.word")
        table.add_column("Pronunciation")
        table.add_column("English", style="text")
        table.add_column("Definition", overflow="fold")
        for number, word, pronunciation, definition, gloss in results:
            table.add_row(number, word, pronunciation, gloss, definition)
        console.print(table)
        return True

    def display_dictionary(self, query):
        normalized = query.strip()
        if not normalized:
            return False
        res = self.dictionary_db.query(
            """
            SELECT topic, content, source
            FROM dictionary
            WHERE lower(topic) = lower(?)
            LIMIT 3
            """,
            (normalized,)
        )
        if not res:
            res = self.dictionary_db.query(
                """
                SELECT topic, content, source
                FROM dictionary
                WHERE lower(topic) LIKE lower(?)
                ORDER BY CASE WHEN lower(topic) LIKE lower(?) THEN 0 ELSE 1 END, topic
                LIMIT 3
                """,
                (f"{normalized}%", f"%{normalized}%")
            )
        if not res:
            safe_query = self.escape_fts_query(query)
            if not safe_query:
                return False
            res = self.dictionary_db.query(
                "SELECT topic, content, source FROM dictionary_fts WHERE dictionary_fts MATCH ? LIMIT 3",
                (safe_query,)
            )
        for t, c, s in res:
            console.print(Panel(Markdown(c), title=f"📖 {t} ({s})", border_style="violet"))
        return bool(res)

    def display_encyclopedia(self, query):
        if not self.encyclopedia_db:
            return False
        normalized = query.strip()
        if not normalized:
            return False
        tables = {row[0] for row in self.encyclopedia_db.query("SELECT name FROM sqlite_master WHERE type='table'")}
        if "encyclopedia" not in tables:
            return False
        res = self.encyclopedia_db.query(
            """
            SELECT topic, content, source
            FROM encyclopedia
            WHERE lower(topic) = lower(?)
            LIMIT 3
            """,
            (normalized,)
        )
        if not res:
            res = self.encyclopedia_db.query(
                """
                SELECT topic, content, source
                FROM encyclopedia
                WHERE lower(topic) LIKE lower(?)
                ORDER BY CASE WHEN lower(topic) LIKE lower(?) THEN 0 ELSE 1 END, topic
                LIMIT 3
                """,
                (f"{normalized}%", f"%{normalized}%")
            )
        if not res and "encyclopedia_fts" in tables:
            safe_query = self.escape_fts_query(query)
            if not safe_query:
                return False
            res = self.encyclopedia_db.query(
                """
                SELECT topic, content, source
                FROM encyclopedia_fts
                WHERE encyclopedia_fts MATCH ?
                LIMIT 3
                """,
                (safe_query,)
            )
        for t, c, s in res:
            console.print(Panel(Markdown(c), title=f"📚 {t} ({s})", border_style="cyan"))
        return bool(res)

    def format_naves_entry(self, entry_text):
        # Precise Regex for Verse References: handles "1 Cor 1:1", "Jhn 3:16", etc.
        pattern = r'(\b(?:[1-3]\s?)?[A-Z][a-z]{0,3}\s\d+:\d+[0-9:,\-]*|\b(?:[1-3]\s?)?[A-Z]{2,4}\s\d+:\d+[0-9:,\-]*)'

        formatted_content = Text()
        lines = entry_text.splitlines()

        for line in lines:
            if not line.strip():
                formatted_content.append("\n")
                continue

            stripped = line.lstrip()
            indent_size = len(line) - len(stripped)

            # Determine style based on indentation depth
            if indent_size == 0:
                style = "category" if ACTIVE_THEME_MODE == "dark" else "text"
            elif indent_size <= 5:
                style = "ui.cyan"
            else:
                style = "ui.meta"

            formatted_content.append(" " * indent_size)

            # Handle Bullets
            current_content = stripped
            if stripped.startswith('-'):
                formatted_content.append("-", style="ui.meta")
                current_content = stripped[1:]

            # Highlight Verses without duplication
            last_pos = 0
            for match in re.finditer(pattern, current_content):
                formatted_content.append(current_content[last_pos:match.start()], style=style)
                formatted_content.append(match.group(0), style="verse.ref")
                last_pos = match.end()

            formatted_content.append(current_content[last_pos:], style=style)
            formatted_content.append("\n")

        return formatted_content

    def display_reverse_naves(self, reference):
        ref_norm, book, chap, verse, v_end = self.normalize_ref(reference)
        if not ref_norm or not verse:
            return False

        db_ref = f"{self.bible_prefix}:{ref_norm}"
        topics = self.get_reverse_naves(db_ref)
        if not topics:
            return False

        console.print(Rule(style="ui.border"))
        console.print(f" [dict.topic]Nave's Topics for {reference}[/dict.topic]")
        console.print(Rule(style="ui.border"))
        console.print("")

        topics_text = Text()
        for idx, topic in enumerate(topics):
            if idx > 0:
                topics_text.append("  •  ", style="dim")
            topics_text.append(topic, style="dict.topic")
        console.print(topics_text)
        console.print("")
        console.print(Rule(style="ui.meta"))
        return True

    def display_naves(self, query):
        if not self.naves_db:
            return False
        normalized = query.strip()
        if not normalized:
            return False

        # 1. Try Exact Match first (Fastest)
        res = self.naves_db.query(
            "SELECT subject, entry FROM topics WHERE subject_upper = ?",
            (normalized.upper(),)
        )

        # 2. If no exact match, try Substring Match or FTS
        if not res:
            res = self.naves_db.query(
                "SELECT subject, entry FROM topics WHERE subject_upper LIKE ? ORDER BY subject LIMIT 50",
                (f"%{normalized.upper()}%",)
            )

        if not res:
            try:
                res = self.naves_db.query(
                    "SELECT subject, entry FROM topics_fts WHERE entry MATCH ? LIMIT 20",
                    (normalized,)
                )
            except: pass

        if not res:
            # 3. Try Reverse lookup (maybe the query is a verse)
            if self.display_reverse_naves(query):
                return True
            return False

        if len(res) == 1:
            subject, entry = res[0]
            console.print(Panel(
                self.format_naves_entry(entry),
                title=subject,
                title_align="left",
                border_style="ui.border",
                padding=(1, 2),
                expand=True,
            ))
            return True
        else:
            # Multiple results - show picker
            table = Table(title=f"Nave's Results for '{query}'", border_style="ui.border", box=box.SIMPLE_HEAVY, expand=True)
            table.add_column("ID", justify="right", style="ui.action.key")
            table.add_column("Subject", style="dict.topic")
            for i, row in enumerate(res, 1):
                table.add_row(str(i), row[0])
            console.print(table)

            try:
                choice = console.input(f"\n [ui.action]Select ID [1-{len(res)}]: [/]").strip()
                idx = int(choice) - 1
                if 0 <= idx < len(res):
                    self.display_naves(res[idx][0]) # Recursive call to show the single entry
                    return True
            except:
                pass
            return True

    def display_commentary(self, query):
        ref_norm, book, chap, verse, v_end = self.normalize_ref(query)
        if not ref_norm:
            return False

        ref_label = f"{book} {chap}"
        if verse:
            ref_label += f":{verse}"

        console.print(Rule(style="ui.border"))
        console.print(f" [dict.topic]Commentary: {ref_label}[/dict.topic]")
        console.print(Rule(style="ui.border"))
        console.print("")

        # Ensure we use the DB-specific book name
        db_book = self.canon_map.get(re.sub(r"[^a-z0-9]+", "", book.lower()), book)

        # Helper to query a commentary DB
        def get_comm(db, b, c, v):
            if not db: return []
            if v:
                # Find sections that cover this verse
                return db.query(
                    "SELECT section_title, markdown, source FROM commentary WHERE book = ? AND chapter = ? AND verse_start <= ? AND verse_end >= ? ORDER BY section_order",
                    (b, c, v, v)
                )
            else:
                # Get entire chapter
                return db.query(
                    "SELECT section_title, markdown, source FROM commentary WHERE book = ? AND chapter = ? ORDER BY section_order",
                    (b, c)
                )

        henry_notes = get_comm(self.henry_db, db_book, chap, verse)
        calvin_notes = get_comm(self.calvin_db, db_book, chap, verse)

        if not henry_notes and not calvin_notes:
            console.print(f"[warning]No commentary notes found for {ref_label}.[/]")
            return False

        # Display in side-by-side or sequential blocks
        for title, md, src in henry_notes:
            console.print(Panel(Markdown(md), title=f"📜 {src}: {title}", border_style="blue"))

        for title, md, src in calvin_notes:
            console.print(Panel(Markdown(md), title=f"📜 {src}: {title}", border_style="magenta"))

        return True

# ---------------------------------------------------------------------------
# CLI dispatch
# ---------------------------------------------------------------------------
# Keep parsing and routing here thin. Feature behavior should live on LexAgent
# so commands can eventually be tested without shelling out.
def main():
    raw_argv = sys.argv[1:]
    if "search" in raw_argv or "serch" in raw_argv:
        command_idx = raw_argv.index("search") if "search" in raw_argv else raw_argv.index("serch")
        protected_argv = []
        for idx, token in enumerate(raw_argv):
            if (
                idx > command_idx
                and token.startswith("-")
                and not token.startswith("--")
                and token not in {"-i", "-d", "-c", "-s", "-v", "-light", "-dark", "-auto", "-B", "--bible"}
            ):
                protected_argv.append(f"__lexscope__{token[1:]}")
            else:
                protected_argv.append(token)
        raw_argv = protected_argv
    parser = argparse.ArgumentParser(
        prog="lex",
        description="Lex: local-first Bible reading, study, search, and manuscript tools.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""commands:
  lex John 3:16                    read a verse in context
  lex study John 1:1               open interlinear and lexicon study
  lex search "holy spirit" -nt     ranked, scoped scripture search
  lex manuscript John 1:1          map readings and manuscript witnesses
  lex manuscript P66               open a manuscript profile
  lex web John 3:16                 map connected verses
  lex strongs love                  search Strong's entries
  lex naves grace                   browse Nave's topical index
  lex commentary John 3:16         read available commentary
  lex -v esv                        set the default Bible version

Run `lex <command>` without an argument for command-specific help.
""",
    )
    parser.add_argument("query", nargs="*")
    parser.add_argument("-i", "--interlinear", action="store_true")
    parser.add_argument("-d", "--define", action="store_true")
    parser.add_argument("-c", "--creed", action="store_true")
    parser.add_argument("-s", "--strongs", action="store_true")
    parser.add_argument("-v", dest="bible_version", type=str, default=None, help="Set Bible version (e.g. -v vulgate, -v kjv, -v esv)")
    parser.add_argument("--version", action="version", version=VERSION)
    parser.add_argument("-B", "--bible", type=str, default=None, choices=BIBLE_VERSIONS.keys(), help="Select Bible version for current command (legacy)")
    parser.add_argument("--update", action="store_true", help="Check for and install data updates")
    theme_group = parser.add_mutually_exclusive_group()
    theme_group.add_argument("-light", dest="theme_mode", action="store_const", const="light")
    theme_group.add_argument("-dark", dest="theme_mode", action="store_const", const="dark")
    theme_group.add_argument("-auto", dest="theme_mode", action="store_const", const="auto")
    parser.add_argument("--credits", action="store_true")
    parser.add_argument("--next", action="store_true")
    parser.add_argument("--prev", action="store_true")
    parser.add_argument("--page", type=int, default=1)
    parser.add_argument("--limit", type=int, default=10)
    parser.add_argument("--mode", type=str, choices=["read", "study", "web"], default="read")
    parser.add_argument("--format", type=str, choices=["png", "pptx", "pdf", "docx"], default=None)
    parser.add_argument("--ratio", type=str, choices=["16:9", "1:1"], default="16:9", help="Aspect ratio for PNG exports")
    parser.add_argument("--animate", dest="animate", action="store_true", default=None)
    parser.add_argument("--no-animate", dest="animate", action="store_false")
    parser.add_argument("--clear", action="store_true", help="Clear command history when used with `lex history`")
    args, unknown = parser.parse_known_args(raw_argv)
    if args.theme_mode == "auto":
        clear_theme_preference()
    elif args.theme_mode:
        save_theme_preference(args.theme_mode)
    args.query = [f"-{q[len('__lexscope__'):]}" if q.startswith("__lexscope__") else q for q in args.query]

    # Initialize Update Manager and ensure data exists
    manager = LexUpdateManager(console, data_dir=DATA_DIR)

    # If the user is running an update, we handle it and exit
    if args.update or (len(args.query) > 0 and args.query[0] == "update"):
        if not manager.perform_update():
            sys.exit(1)
        sys.exit(0)

    # Ensure critical data exists (downloads if missing)
    if not manager.ensure_data():
        sys.exit(1)

    if unknown:
        if args.query and args.query[0] in {"search", "serch", "read"} and all(u.startswith("-") and not u.startswith("--") for u in unknown):
            args.query.extend(unknown)
        else:
            parser.error(f"unrecognized arguments: {' '.join(unknown)}")
    query = " ".join(args.query)

    # Alias "bible" or "version" to persistent switch
    if query.startswith("bible ") or query.startswith("version "):
        target = query.split(" ", 1)[1].strip()
        if target in BIBLE_VERSIONS:
            if not bible_version_available(target):
                print_missing_bible_version(target)
                sys.exit(1)
            save_bible_preference(target)
            console.print(f"[success]Default Bible version set to [bold cyan]{target}[/] ({BIBLE_VERSIONS[target]['name']})[/]")
            sys.exit(0)
        else:
            console.print(f"[error]Unknown Bible version: {target}[/]")
            sys.exit(1)

    # Handle -v Bible-version selection cleanly
    if args.bible_version:
        target = args.bible_version.lower()
        # Resolve alias (vulgate -> vulg)
        if target in BIBLE_VERSIONS and BIBLE_VERSIONS[target].get("alias"):
            target = BIBLE_VERSIONS[target]["alias"]
        if target in BIBLE_VERSIONS:
            if not bible_version_available(target):
                print_missing_bible_version(target)
                sys.exit(1)
            save_bible_preference(target)
            console.print(f"[success]Default Bible version set to [bold cyan]{target}[/] ({BIBLE_VERSIONS[target]['name']})[/]")
            if len(args.query) == 0:
                sys.exit(0)
            # Continue with the query using the new default
        else:
            console.print(f"[error]Unknown Bible version: {target}[/]")
            console.print(f"Available: {', '.join(BIBLE_VERSIONS.keys())}")
            sys.exit(1)
        # Interactive menu removed. Use `lex -v <id>` to switch version cleanly.

    # Handle persistent bible selection: "lex -B kjv" with no query
    if args.bible and not query:
        if not bible_version_available(args.bible):
            print_missing_bible_version(args.bible)
            sys.exit(1)
        save_bible_preference(args.bible)
        console.print(f"[success]Default Bible version set to [bold cyan]{args.bible}[/] ({BIBLE_VERSIONS[args.bible]['name']})[/]")
        sys.exit(0)

    selected_bible = args.bible or load_bible_preference()
    if not bible_version_available(selected_bible):
        print_missing_bible_version(selected_bible)
        sys.exit(1)

    agent = LexAgent(bible_id=selected_bible)

    if args.credits:
        agent.display_credits()
        sys.exit(0)

    if query == "history":
        if args.clear:
            if agent.clear_query_history():
                console.print("[success]Lex query history cleared.[/]")
                sys.exit(0)
            console.print("[error]Could not clear Lex query history.[/]")
            sys.exit(1)
        agent.display_query_history(limit=args.limit)
        sys.exit(0)
    elif query.startswith("history "):
        subquery = query.split(" ", 1)[1].strip()
        if subquery == "clear":
            if agent.clear_query_history():
                console.print("[success]Lex query history cleared.[/]")
                sys.exit(0)
            console.print("[error]Could not clear Lex query history.[/]")
            sys.exit(1)
        agent.display_query_history(limit=args.limit)
        sys.exit(0)

    if raw_argv:
        agent.save_query_history("lex " + shlex.join(raw_argv))

    if args.next or args.prev:
        last = agent.last_ref
        if not last: sys.exit(1)
        query = agent.resolve_navigation_query("next" if args.next else "prev")
        if not query:
            sys.exit(1)

    if not query and not (args.next or args.prev):
        agent.display_intro()
        sys.exit(0)

    if query == "read":
        agent.display_read_landing()
        sys.exit(0)
    elif query == "read all" or query == "all":
        last_ref = agent.last_ref
        if not last_ref:
            console.print("[warning]No previous verse reference found in history.[/]")
            sys.exit(1)
        refs = agent.get_tsk_crossrefs(last_ref)[:10]
        if not refs:
            console.print("[warning]No cross references found in history for the last verse.[/]")
            sys.exit(1)
        any_success = False
        for to_ref, votes in refs:
            db_ref = agent.parse_tsk_ref(to_ref)
            if db_ref:
                display_ref = agent.format_display_ref(db_ref)
                if agent.display_verse(display_ref, interlinear=args.interlinear, animate=args.animate):
                    any_success = True
        if not any_success:
            console.print("[warning]Failed to display any cross references in read mode.[/]")
            sys.exit(1)
        sys.exit(0)
    elif query.startswith("read "):
        query = query[5:].strip()
    elif query == "study":
        agent.display_study_landing()
        sys.exit(0)
    elif query.startswith("study "):
        query = query[6:].strip()
        if not query:
            agent.display_study_landing()
            sys.exit(0)
        # Split the query by separators to support sequential study if multiple are provided
        study_parts = [p.strip() for p in re.split(r'[,;]+', query) if p.strip()]
        if len(study_parts) > 1:
            any_success = False
            for part in study_parts:
                if agent.display_study_query(part, animate=args.animate):
                    any_success = True
            if not any_success:
                console.print("[warning]No local study data found for those references.[/]")
                sys.exit(1)
            sys.exit(0)
        else:
            if not agent.display_study_query(query, animate=args.animate):
                console.print("[warning]No local study data found for that reference.[/]")
                sys.exit(1)
            sys.exit(0)
    elif query in {"search", "serch"}:
        agent.display_search_howto()
        sys.exit(0)
    elif query.startswith("search ") or query.startswith("serch "):
        query = query.split(" ", 1)[1].strip()
        if not query:
            sys.exit(1)
        if not agent.display_search(query, page=args.page, limit=args.limit):
            console.print("[warning]No scripture search results found.[/]")
            sys.exit(1)
        sys.exit(0)
    elif query == "web":
        console.print("[warning]Usage: lex web John 3:16[/]")
        sys.exit(1)
    elif query.startswith("web "):
        q = query[4:].strip()
        if not agent.display_verse_web(q, limit=args.limit):
            console.print("[warning]No verse web found.[/]")
            sys.exit(1)
        sys.exit(0)
    elif query == "manuscript":
        agent.display_manuscript_howto()
        sys.exit(0)
    elif query.startswith("manuscript "):
        q = query[len("manuscript "):].strip()
        if not q:
            agent.display_manuscript_howto()
            sys.exit(1)
        if not agent.display_manuscript(q, limit=args.limit):
            sys.exit(1)
        sys.exit(0)
    elif query == "export":
        agent.display_export_howto()
        sys.exit(0)
    elif query.startswith("export "):
        refs_str = query[7:].strip()
        if not refs_str:
            agent.display_export_howto()
            sys.exit(0)

        fmt = args.format
        if not fmt:
            agent.prompt_bulk_export(refs_str, mode=args.mode)
        else:
            agent.handle_bulk_export(refs_str, mode=args.mode, format=fmt, ratio=args.ratio)
        sys.exit(0)
    elif query == "strongs":
        agent.display_strongs_howto()
        sys.exit(0)
    elif query.startswith("strongs "):
        q = query[8:].strip()
        if not q:
            agent.display_strongs_howto()
            sys.exit(1)
        if re.match(r'^[GH]\d+$', q, re.IGNORECASE):
            if not agent.display_strongs(q):
                console.print("[warning]No Strong's entry found for that number.[/]")
                sys.exit(1)
        else:
            if not agent.display_english_strongs(q):
                if not agent.display_strongs(q):
                    console.print("[warning]No Strong's entries found for that term.[/]")
                    sys.exit(1)
        sys.exit(0)
    elif query == "topic" or query == "naves":
        agent.display_topic_howto()
        sys.exit(0)
    elif query.startswith("topic ") or query.startswith("naves "):
        q = query.replace("topic ", "").replace("naves ", "").strip()
        if not agent.display_naves(q):
            console.print("[warning]No Nave's Topical Bible entry found.[/]")
            sys.exit(1)
        sys.exit(0)
    elif query == "commentary":
        agent.display_commentary_howto()
        sys.exit(0)
    elif query.startswith("commentary "):
        q = query.replace("commentary ", "").strip()
        if not agent.display_commentary(q):
            sys.exit(1)
        sys.exit(0)

    if args.define or query.startswith("define"):
        q = query.replace("define ", "").strip()
        if not q or q == "define": agent.display_dictionary_howto()
        else:
            dictionary_found = agent.display_dictionary(q)
            encyclopedia_found = agent.display_encyclopedia(q)
            if not dictionary_found and not encyclopedia_found:
                console.print("[warning]No dictionary or encyclopedia entry found.[/]")
    elif args.creed or query in {"creed", "creeds"} or query.startswith(("creed ", "creeds ")):
        q = re.sub(r"^creeds?\s*", "", query).strip()
        if not q: agent.display_creed_navigator()
        elif not agent.display_creed_navigator(q):
            console.print("[warning]No creed or confession entry found.[/]")
            sys.exit(1)
    elif args.strongs:
        if not query:
            agent.display_strongs_howto()
            sys.exit(1)
        if re.match(r'^[GH]\d+$', query, re.IGNORECASE):
            if not agent.display_strongs(query):
                console.print("[warning]No Strong's entry found for that number.[/]")
                sys.exit(1)
        else:
            if not agent.display_english_strongs(query):
                if not agent.display_strongs(query):
                    console.print("[warning]No Strong's entries found for that term.[/]")
                    sys.exit(1)
    elif re.match(r'^[GH]\d+', query, re.IGNORECASE):
        if not agent.display_strongs(query):
            console.print("[warning]No Strong's entry found for that number.[/]")
            sys.exit(1)
    elif query:
        # Check if the query has multiple references separated by comma or semicolon
        ref_parts = [p.strip() for p in re.split(r'[,;]+', query) if p.strip()]
        if len(ref_parts) > 1:
            any_success = False
            for part in ref_parts:
                if agent.display_verse(part, interlinear=args.interlinear, animate=args.animate):
                    any_success = True
                else:
                    if agent.display_strongs(part):
                        any_success = True
            if not any_success:
                agent.display_search_howto()
                sys.exit(1)
        else:
            if not agent.display_verse(query, interlinear=args.interlinear, animate=args.animate):
                if not agent.display_strongs(query):
                    agent.display_search_howto()
                    sys.exit(1)
    else:
        agent.display_intro()

if __name__ == "__main__":
    main()
